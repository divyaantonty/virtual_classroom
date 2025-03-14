import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
        
    # Get teacher courses
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    
    # Get unread message count
    unread_count = ParentTeacherMessage.objects.filter(
        teacher_id=teacher_id,
        message_type='parent_to_teacher',
        is_read=False
    ).count()
    
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses,
        'unread_messages': unread_count
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
        
    # Get teacher courses
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    
    # Get unread message count
    unread_count = ParentTeacherMessage.objects.filter(
        teacher_id=teacher_id,
        message_type='parent_to_teacher',
        is_read=False
    ).count()
    
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses,
        'unread_messages': unread_count
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
        
    # Get teacher courses
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    
    # Get unread message count
    unread_count = ParentTeacherMessage.objects.filter(
        teacher_id=teacher_id,
        message_type='parent_to_teacher',
        is_read=False
    ).count()
    
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses,
        'unread_messages': unread_count
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
        
    # Get teacher courses
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    
    # Get unread message count
    unread_count = ParentTeacherMessage.objects.filter(
        teacher_id=teacher_id,
        message_type='parent_to_teacher',
        is_read=False
    ).count()
    
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses,
        'unread_messages': unread_count
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent

def parent_update_profile(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        messages.error(request, 'No parent ID found in session.')
        return redirect('login')

    parent = get_object_or_404(Parent, id=parent_id)

    if request.method == 'POST':
        # Fetch the form data
        parent.auto_generated_username = request.POST.get('auto_generated_username')
        

        # Save the updated teacher details
        parent.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('parent_dashboard')  

    context = {'parent': parent}
    return render(request, 'parent_update_profile.html', context)
from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.contrib import messages  # For displaying messages
from .models import Material, Course, Teacher, TeacherCourse  # Import the TeacherCourse model

def upload_material(request):
    if request.method == 'POST':
        teacher_id = request.session.get('teacher_id')
        
        if not teacher_id:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'success': False, 'error': 'Not logged in as teacher'})
            else:
                messages.error(request, 'You are not logged in as a teacher.')
                return redirect('login')

        try:
            # Check if this is a note upload (AJAX request)
            note_id = request.POST.get('note_id')
            if note_id:
                # Get the note
                note = TeacherNote.objects.get(id=note_id, teacher_id=teacher_id)
                course = note.course

                # Create PDF from note content
                buffer = BytesIO()
                p = canvas.Canvas(buffer)
                
                # Add title
                p.setFont("Helvetica-Bold", 16)
                p.drawString(100, 800, note.title)
                
                # Add content
                p.setFont("Helvetica", 12)
                y = 750
                for line in note.content.split('\n'):
                    if y > 50:  # Ensure we don't write below page margin
                        p.drawString(100, y, line.strip())
                        y -= 20
                
                p.showPage()
                p.save()
                
                # Create PDF file
                pdf_file = ContentFile(buffer.getvalue())
                
                # Create material
                material = Material.objects.create(
                    teacher_id=teacher_id,
                    course=course,
                    description=f"Note: {note.title}",
                )
                
                # Save the PDF file
                material.file.save(f'note_{note.id}.pdf', pdf_file)

                return JsonResponse({'success': True})
            
            # Handle regular file upload
            else:
                course_id = request.POST.get('course')
                description = request.POST.get('description')
                file = request.FILES.get('file')

                # Verify course assignment
                assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
                if int(course_id) not in assigned_courses:
                    messages.error(request, 'You do not have permission to upload material for this course.')
                    return redirect('upload_material')

                if file:
                    course = get_object_or_404(Course, id=course_id)
                    teacher = get_object_or_404(Teacher, id=teacher_id)

                    material = Material.objects.create(
                        teacher=teacher,
                        course=course,
                        file=file,
                        description=description
                    )
                    messages.success(request, 'Material uploaded successfully!')
                    return redirect('teacher_dashboard')
                else:
                    messages.error(request, 'Please upload a file.')
                    return redirect('upload_material')

        except Exception as e:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'success': False, 'error': str(e)})
            else:
                messages.error(request, f'Error saving material: {str(e)}')
                return redirect('upload_material')

    # GET request - show the upload form
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        messages.error(request, 'You must be logged in as a teacher.')
        return redirect('login')

    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)
    return render(request, 'upload_material.html', {'courses': courses})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.contrib import messages  # type: ignore
from .models import Material, CustomUser, Enrollment
from datetime import datetime
from django.utils import timezone

def view_materials(request):
    # Check if a CustomUser (student) is logged in by checking the session
    custom_user_id = request.session.get('custom_user_id')
    
    if not custom_user_id:
        messages.error(request, 'You must be logged in as a student to view materials.')
        return redirect('login')  # Redirect to login page if no session

    # Fetch the CustomUser (student) object using the session ID
    custom_user = get_object_or_404(CustomUser, id=custom_user_id)

    # Ensure the student is registered for any courses
    enrollments = Enrollment.objects.filter(student=custom_user)
    enrolled_courses = [enrollment.course for enrollment in enrollments]
    
    if not enrollments:
        messages.error(request, 'You are not registered for any course.')
        return redirect('student_dashboard')  # Redirect if no course is associated

    # Prepare a list to hold materials that meet the enrollment criteria
    materials = []
    
    # Loop through each enrollment to filter materials
    for enrollment in enrollments:
        # Combine enrollment_date and enrollment_time into a single datetime object
        if enrollment.enrollment_time:
            enrollment_datetime = datetime.combine(enrollment.enrollment_date, enrollment.enrollment_time)
        else:
            # Fallback if time is not present (you may adjust this as needed)
            enrollment_datetime = timezone.make_aware(datetime.combine(enrollment.enrollment_date, datetime.min.time()))
        
        # Make the combined datetime timezone-aware
        enrollment_datetime = timezone.make_aware(enrollment_datetime)
        
        # Fetch materials related to the course and filter by the enrollment datetime
        course_materials = Material.objects.filter(course=enrollment.course, uploaded_at__gte=enrollment_datetime)
        materials.extend(course_materials)

        context = {
        'materials': materials,
        'enrolled_courses': enrolled_courses,  # Add this line
    }

    # Render the materials in the template
    return render(request, 'view_materials.html', context)

# views.py
from .models import Material, Parent, CustomUser

def view_study_materials(request):

    parent = get_object_or_404(Parent, auto_generated_username=request.user.username)
    
    student = get_object_or_404(CustomUser, username=parent.student_username)

    materials = Material.objects.filter(course=student.course)
    
    return render(request, 'view_study_materials.html', {'materials': materials, 'student': student})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.utils import timezone  # type: ignore
from .models import Quizs, Course, TeacherCourse  # Import TeacherCourse for filtering

def create_quiz(request):
    teacher_id = request.session.get('teacher_id')  # Retrieve the teacher's ID from the session

    if not teacher_id:
        # Handle the case where the teacher is not logged in
        return redirect('login')

    # Fetch courses taught by the teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)  # Filter to only assigned courses

    if request.method == 'POST':
        course_id = request.POST.get('course')
        title = request.POST.get('title')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        duration = request.POST.get('duration')

        course = get_object_or_404(Course, id=course_id)
        
        # Create and save the quiz
        quiz = Quizs(course=course, teacher_id=teacher_id, title=title, 
                     start_date=start_date, end_date=end_date, 
                     start_time=start_time, end_time=end_time,duration=duration)
        quiz.save()

        return redirect('add_question', quiz_id=quiz.id)
    
    return render(request, 'create_quiz.html', {'courses': courses})


def add_question(request, quiz_id):
    quiz = get_object_or_404(Quizs, id=quiz_id)

    if request.method == 'POST':
        questions_data = request.POST.getlist('question')
        options_a = request.POST.getlist('option_a')
        options_b = request.POST.getlist('option_b')
        options_c = request.POST.getlist('option_c')
        options_d = request.POST.getlist('option_d')
        correct_options = request.POST.getlist('correct_option')

        # Iterate through all submitted questions
        for i in range(len(questions_data)):
            question_text = questions_data[i]
            option_a = options_a[i]
            option_b = options_b[i]
            option_c = options_c[i]
            option_d = options_d[i]
            correct_option = correct_options[i]

            # Create a new Question object for each set of data
            question = Question(
                quiz=quiz, 
                text=question_text,
                option_a=option_a,
                option_b=option_b,
                option_c=option_c,
                option_d=option_d,
                correct_option=correct_option
            )
            question.save()

        return redirect('add_question', quiz_id=quiz.id)

    # Add this return statement for GET requests
    return render(request, 'add_question.html', {
        'quiz': quiz,
        'quiz_id': quiz_id,
        'first_name': request.session.get('first_name', ''),
        'last_name': request.session.get('last_name', '')
    })

from django.shortcuts import render, get_object_or_404
from .models import Quizs, Question

def take_quiz(request, quiz_id):
    quiz = get_object_or_404(Quizs, id=quiz_id)
    questions = Question.objects.filter(quiz=quiz)

    if request.method == 'POST':
        user_answers = {}  # Store user's answers
        correct_count = 0  # Track the number of correct answers

        # Loop through each question in the quiz
        for question in questions:
            user_answer = request.POST.get(f'question_{question.id}')  # Get user's answer for the question
            user_answers[question.id] = user_answer  # Save the user's answer

            # Check if the user's answer matches the correct answer
            if user_answer == question.correct_option:
                correct_count += 1  # Increment correct count if answer is correct

        # After submission, show the result
        return render(request, 'quiz_result.html', {
            'quiz': quiz,
            'questions': questions,
            'user_answers': user_answers,
            'correct_count': correct_count,
            'total_questions': questions.count()
        })

    # Render the quiz page with questions
    return render(request, 'take_quiz.html', {
        'quiz': quiz,
        'questions': questions
    })

from django.shortcuts import render, redirect
from django.contrib import messages
from django.utils import timezone
from django.db.models import Q
from .models import CustomUser, Quizs, UserAnswers, Enrollment  # Ensure Enrollment is imported

def available_quizzes(request):
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    # Fetch enrollments for the student
    enrollments = Enrollment.objects.filter(student=student)

    if not enrollments.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    current_datetime = timezone.localtime()

    # Create a list of enrolled courses with their enrollment dates and times
    enrolled_courses = [
        (enrollment.course, enrollment.enrollment_date, enrollment.enrollment_time)
        for enrollment in enrollments
    ]

    available_quizzes = []

    # Fetch quizzes that are ongoing or upcoming and have not ended
    for course, enrollment_date, enrollment_time in enrolled_courses:
        # Combine enrollment date and time to create a full datetime object
        enrollment_datetime = timezone.make_aware(datetime.combine(enrollment_date, enrollment_time or datetime.min.time()))

        # Fetch quizzes based on the course and check the date and time constraints
        quizzes_for_course = Quizs.objects.filter(
            course=course,
            start_date__lte=current_datetime.date(),
            end_date__gte=current_datetime.date(),
        ).filter(
            Q(start_time__lte=current_datetime.time()) | Q(start_date__gt=current_datetime.date())
        ).exclude(
            Q(end_date__lt=current_datetime.date()) | 
            (Q(end_date=current_datetime.date()) & Q(end_time__lt=current_datetime.time()))
        ).order_by('start_date', 'start_time')

        # Filter out quizzes that the student has already attempted
        quizzes_for_course = [quiz for quiz in quizzes_for_course if not UserAnswers.objects.filter(user=student, question__quiz=quiz).exists()]

        # Add the filtered quizzes to the available_quizzes list
        available_quizzes.extend(quizzes_for_course)

    return render(request, 'available_quizzes.html', {
        'quizzes': available_quizzes,
        'today': current_datetime.date(),
        'current_time': current_datetime.time()
    })



def submit_quiz(request, quiz_id):
    # Get the quiz and its associated questions
    quiz = get_object_or_404(Quizs, id=quiz_id)
    questions = quiz.questions.all()

    # Retrieve the `custom_user_id` from the session
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You must be logged in to submit the quiz.")
        return redirect('login')

    # Fetch the CustomUser based on the custom user ID from the session
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "User not found.")
        return redirect('login')

    # Initialize counters
    correct_count = 0
    total_questions = questions.count()

    # Prepare to store user answers with questions
    question_data = []

    # Loop through the quiz questions
    for question in questions:
        # Get the user's selected answer from POST data
        user_answer = request.POST.get(f'question_{question.id}')

        # Only process if the user provided an answer
        if user_answer in ['A', 'B', 'C', 'D']:
            # Store the user's answer in the UserAnswers table
            UserAnswers.objects.update_or_create(
                user=custom_user,
                question=question,
                defaults={'selected_option': user_answer}
            )

            # Increment correct count if the answer is correct
            if user_answer == question.correct_option:
                correct_count += 1
        else:
            user_answer = None

        # Retrieve the stored answer from the database
        try:
            stored_answer = UserAnswers.objects.get(user=custom_user, question=question).selected_option
        except UserAnswers.DoesNotExist:
            stored_answer = "Not answered"

        # Append the question and user answer to the list
        question_data.append({
            'question': question,
            'user_answer': stored_answer,
            'correct_answer': question.correct_option,
        })

    # Render the result page, passing the structured data to the template
    context = {
        'quiz': quiz,
        'questions_data': question_data,
        'correct_count': correct_count,
        'total_questions': total_questions,
    }
    return render(request, 'quiz_result.html', context)


def quiz_result(request, quiz_id):
    # Ensure the user is logged in
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You must be logged in to view quiz results.")
        return redirect('login')

    # Get the quiz
    quiz = get_object_or_404(Quizs, id=quiz_id)

    # Retrieve user's answers for the specific quiz
    user_answers = UserAnswers.objects.filter(user_id=custom_user_id, question__quiz=quiz)

    # Prepare data for display
    results = []
    correct_count = 0
    total_questions = quiz.questions.count()

    # Create a dictionary to map question ids to user answers
    user_answers_dict = {user_answer.question.id: user_answer.selected_option for user_answer in user_answers}

    for question in quiz.questions.all():  # Iterate through all questions in the quiz
        selected_option = user_answers_dict.get(question.id, "Not answered")  # Get the selected option or default to "Not answered"
        correct_answer = question.correct_option
        is_correct = (selected_option == correct_answer)

        results.append({
            'question': question,
            'selected_option': selected_option,
            'correct_option': correct_answer,
            'is_correct': is_correct,  # Keep track of whether the answer is correct
        })

        if is_correct:
            correct_count += 1

    context = {
        'quiz': quiz,
        'results': results,
        'correct_count': correct_count,
        'total_questions': total_questions,
    }
    return render(request, 'quiz_result.html', context)


def quiz_questions(request, quiz_id):
    quiz = Quizs.objects.get(id=quiz_id)
    questions = Question.objects.filter(quizs=quiz)  # Adjust based on your models
    return render(request, 'quiz_questions.html', {'quiz': quiz, 'questions': questions})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.http import HttpResponse  # type: ignore
from .models import Assignment, Course, Teacher, TeacherCourse  # Import the TeacherCourse model
from django.core.exceptions import ValidationError  # type: ignore
from datetime import datetime

def create_assignment(request):
    if request.method == 'POST':
        # Fetch form data
        title = request.POST.get('title')
        description = request.POST.get('description')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        course_name_id = request.POST.get('course')
        file = request.FILES.get('file')

        # Validate and convert date and time
        try:
            start_date = datetime.strptime(start_date, '%Y-%m-%d').date()
            end_date = datetime.strptime(end_date, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time, '%H:%M').time()
            end_time = datetime.strptime(end_time, '%H:%M').time()

            teacher_id = request.session.get('teacher_id')
            teacher = get_object_or_404(Teacher, id=teacher_id)

            assignment = Assignment(
                title=title,
                description=description,
                start_date=start_date,
                end_date=end_date,
                start_time=start_time,
                end_time=end_time,
                file=file,
                course_name_id=course_name_id,
                teacher=teacher
            )

            if assignment.start_date > assignment.end_date:
                raise ValidationError("Start date cannot be after end date.")

            assignment.save()

            return redirect('teacher_dashboard')

        except ValidationError as e:
            return render(request, 'create_assignment.html', {
                'error': str(e),
                'courses': courses,  # Pass the courses to the template
                'assignment': request.POST
            })
        except Exception as e:
            return render(request, 'create_assignment.html', {
                'error': str(e),
                'courses': courses,
                'assignment': request.POST
            })

    # If not a POST request, get the assigned courses for the teacher
    teacher_id = request.session.get('teacher_id')
    teacher = get_object_or_404(Teacher, id=teacher_id)

    assigned_courses = TeacherCourse.objects.filter(teacher=teacher).select_related('course')
    courses = [tc.course for tc in assigned_courses]

    return render(request, 'create_assignment.html', {
        'courses': courses,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
    })

from django.contrib import messages
from django.shortcuts import redirect, render
from django.utils import timezone
from datetime import datetime
from .models import Assignment, AssignmentSubmission, CustomUser, Enrollment  # Ensure to import Enrollment

def assignment_submission_view(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')  # Redirect to login page if user ID is not in session

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    # Fetch the enrolled courses for the student along with the enrollment date and time
    enrollments = Enrollment.objects.filter(student=student)

    if not enrollments:
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Create a list to store the courses and their enrollment date and time
    enrolled_courses_with_dates = [(enrollment.course, enrollment.enrollment_date, enrollment.enrollment_time) for enrollment in enrollments]

    # Create a list to hold the course names and the minimum enrollment datetime
    enrolled_courses = [course for course, _, _ in enrolled_courses_with_dates]
    
    # Fetch assignments based on the registered courses and filter by enrollment datetime
    current_date = timezone.now()
    assignment_details = []

    for course, enrollment_date, enrollment_time in enrolled_courses_with_dates:
        # Combine enrollment date and time to create a full datetime object
        if enrollment_time:
            enrollment_datetime = datetime.combine(enrollment_date, enrollment_time)
        else:
            enrollment_datetime = datetime.combine(enrollment_date, datetime.min.time())

        # Make the combined enrollment datetime timezone-aware
        enrollment_datetime = timezone.make_aware(enrollment_datetime)

        # Get assignments for the current course, considering enrollment datetime
        assignments = Assignment.objects.filter(course_name=course, start_date__gte=enrollment_datetime)

        for assignment in assignments:
            current_datetime = timezone.now()
            assignment_end_datetime = datetime.combine(assignment.end_date, assignment.end_time)
            assignment_end_datetime = timezone.make_aware(assignment_end_datetime)

            has_reached_deadline = assignment_end_datetime <= current_datetime

            # Check if the student has submitted the assignment
            submission = AssignmentSubmission.objects.filter(assignment=assignment, student=student).first()

            # Build the details for each assignment
            assignment_detail = {
                'assignment': assignment,
                'submission': submission,
                'has_reached_deadline': has_reached_deadline,
                'can_submit': not has_reached_deadline and not submission,
                'submission_allowed': not has_reached_deadline and submission,  # Option to re-submit before the deadline
            }
            assignment_details.append(assignment_detail)

    # Handle the file upload if the request method is POST
    if request.method == 'POST':
        assignment_id = request.POST.get('assignment_id')  # Get the assignment ID from the form
        file = request.FILES.get('file')
        if file and assignment_id:
            # Fetch the specific assignment based on the assignment ID
            assignment = Assignment.objects.filter(id=assignment_id, course_name__in=enrolled_courses).first()

            if assignment and not has_reached_deadline and not submission:
                # Create a new submission
                submission = AssignmentSubmission(
                    assignment=assignment,
                    student=student,
                    file=file
                )
                submission.save()
                messages.success(request, 'Your submission has been uploaded successfully!')
                return redirect('student_dashboard')  # Redirect to your desired page
            else:
                if has_reached_deadline:
                    messages.error(request, "Assignment deadline has passed, submission not allowed.")
                else:
                    messages.error(request, "Assignment already submitted.")
        else:
            messages.error(request, 'Please upload a file and select an assignment.')

    # Render the assignment details template with the list of assignment details
    return render(request, 'assignment_detail.html', {'assignment_details': assignment_details})


from django.shortcuts import render, redirect  # type: ignore
from django.utils import timezone
from .models import Assignment, Course

def view_assignment(request):
    if 'teacher_id' in request.session:
        teacher_id = request.session['teacher_id']

        assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
        courses = Course.objects.filter(id__in=assigned_courses)

        # Get the selected course ID from the GET request
        selected_course_id = request.GET.get('course_id')

        # Fetch assignments for the selected course or all assignments if no course is selected
        if selected_course_id:
            assignments = Assignment.objects.filter(course_name__id=selected_course_id, teacher_id=teacher_id)
        else:
            assignments = Assignment.objects.filter(teacher_id=teacher_id)

        # Get the current date and time
        current_datetime = timezone.localtime()

        # Add a 'status' attribute to each assignment based on end date and time
        for assignment in assignments:
            if (assignment.end_date < current_datetime.date() or
                (assignment.end_date == current_datetime.date() and assignment.end_time <= current_datetime.time())):
                assignment.status = 'Completed'
            else:
                assignment.status = 'Ongoing'

        return render(request, 'view_assignment.html', {
            'assignments': assignments,
            'courses': courses,
            
            'selected_course_id': selected_course_id,
            
        })
    else:
        return redirect('login')



from django.db.models import Subquery, OuterRef  # type: ignore
from django.shortcuts import render  # type: ignore
from .models import AssignmentSubmission, Course, TeacherCourse  # Import the TeacherCourse model

def evaluate_assignments(request):
    # Check if the teacher is logged in
    if 'teacher_id' in request.session:
        teacher_id = request.session['teacher_id']
        selected_course_id = request.GET.get('course_id')

        # Fetch all courses assigned to the teacher from the TeacherCourse table
        assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
        courses = Course.objects.filter(id__in=assigned_courses)

        # Get the latest submission for each student
        latest_submissions = AssignmentSubmission.objects.filter(
            student=OuterRef('student'),
            assignment=OuterRef('assignment')
        ).order_by('-submitted_at')

        if selected_course_id:
            # Filter by course and ensure we only get the latest submission for each student
            submissions = AssignmentSubmission.objects.filter(
                assignment__course_name__id=selected_course_id,
                assignment__course_name__id__in=assigned_courses,  # Ensure the course is assigned to the teacher
                submitted_at=Subquery(latest_submissions.values('submitted_at')[:1])
            )
        else:
            # Get the latest submission for each student across all assigned courses
            submissions = AssignmentSubmission.objects.filter(
                assignment__course_name__id__in=assigned_courses,  # Ensure the course is assigned to the teacher
                submitted_at=Subquery(latest_submissions.values('submitted_at')[:1])
            )

        context = {
            'submissions': submissions,
            'courses': courses,
            'selected_course_id': selected_course_id,
            
            
        }
        return render(request, 'evaluate_assignment.html', context)
    else:
        return redirect('login')  # Redirect to login if the teacher is not authenticated

from django.shortcuts import get_object_or_404, redirect # type: ignore
from .models import AssignmentSubmission
from django.views.decorators.http import require_POST # type: ignore


@require_POST # type: ignore
def submit_grade(request, submission_id):
    submission = get_object_or_404(AssignmentSubmission, id=submission_id)
    grade = request.POST.get('grade')

    if grade is not None:
        submission.grade = grade  # Store the grade in the grade field
        submission.save()  # Save the submission to update the database

    return redirect('evaluate_assignment')

from django.shortcuts import render, redirect
from .models import FeedbackQuestion, Course

def add_feedback_question(request):
    courses = Course.objects.all()  # Fetch all courses from the Course table

    if request.method == 'POST':
        questions = []
        course_id = request.POST.get('course')  # Retrieve selected course ID
        release_date = request.POST.get('release_date')
        end_date = request.POST.get('end_date')

        course = Course.objects.get(id=course_id) if course_id else None

        for key, value in request.POST.items():
            if key.startswith('question_text_'):
                questions.append(value)

        if questions and release_date and end_date:
            for question_text in questions:
                # Save each question with an optional associated course, release date, and end date
                FeedbackQuestion.objects.create(
                    question_text=question_text,
                    course=course,
                    release_date=release_date,
                    end_date=end_date
                )
            return redirect('admin_dashboard')
        else:
            return render(request, 'add_feedback_question.html', {
                'courses': courses,
                'error': 'Please enter at least one question, and specify release and end dates.'
            })

    return render(request, 'add_feedback_question.html', {'courses': courses})




import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
        
    # Get teacher courses
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    
    # Get unread message count
    unread_count = ParentTeacherMessage.objects.filter(
        teacher_id=teacher_id,
        message_type='parent_to_teacher',
        is_read=False
    ).count()
    
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses,
        'unread_messages': unread_count
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(parent=parent).order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    # Get all messages for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent,
        message_type='parent_to_teacher'  # Only get messages sent by parent
    ).select_related('teacher').order_by('-date')

    context.update({
        'sent_messages': messages_list,
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent
import json
import uuid
import re
from django.shortcuts import render, redirect
from django.contrib.auth.hashers import make_password
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages
from django.utils.crypto import get_random_string
from django.core.validators import validate_email
from django.core.exceptions import ValidationError
import requests
from .models import CustomUser, Course, FinalExam, FinalExamQuestion, Parent, TeacherMessage, TeacherStudent, Material, MaterialSummary
import cv2
import numpy as np
from datetime import datetime, timedelta
import os
from django.shortcuts import render
from django.http import JsonResponse
from .models import Attendance
import base64
from transformers import MarianMTModel, MarianTokenizer
import torch
from .models import StudentFaceData
from django.views.decorators.csrf import csrf_exempt
from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import ClassSchedule
from django.utils import timezone
import mediapipe as mp
from django.core.files.base import ContentFile
from PIL import Image
import io
from skimage.metrics import structural_similarity as ssim  # Add this import
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone
from .models import Course, WhiteboardShare
from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent
from django.views.decorators.http import require_POST
from .models import MindMap
from django.views.decorators.http import require_http_methods
from django.views.decorators.csrf import ensure_csrf_cookie
import os
from dotenv import load_dotenv
from gtts import gTTS
from django.http import FileResponse, HttpResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile
from transformers import SpeechT5Processor, SpeechT5ForTextToSpeech, SpeechT5HifiGan
import soundfile as sf
from datasets import load_dataset
import pyttsx3
import threading
from .models import TeacherNote, Course
from django.contrib import messages
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re
from reportlab.pdfgen import canvas
from io import BytesIO
from django.core.files.base import ContentFile
# Load environment variables
load_dotenv()
def register(request):
    if request.method == 'POST':
        # Extract form data from request.POST
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        date_of_birth = request.POST.get('date_of_birth')
        username = request.POST.get('username')
        password = request.POST.get('password')
       

        # Perform basic validation
        if not all([first_name, last_name, email, contact, date_of_birth, username, password]):
            messages.error(request, 'Please fill out all fields.')
            return render(request, 'register.html')

        if not username.isalnum():
            messages.error(request, "Username should contain only alphabets or numbers.")
            return redirect('register')

        # Check if the username or email is unique for the student
        if CustomUser.objects.filter(username=username).exists():
            messages.error(request, 'Username is already taken.')
            return render(request, 'register.html')

        if CustomUser.objects.filter(email=email).exists():
            messages.error(request, 'Email is already registered.')
            return render(request, 'register.html')

        # Validate email format
        
        # Validate password complexity
       


        # Create and save student user instance
        student_user = CustomUser(
            username=username,
            first_name=first_name,
            last_name=last_name,
            email=email,
            contact=contact,
            date_of_birth=date_of_birth,
            password=make_password(password),
            is_active=True,  # Ensure the user is active upon registration
        )
        student_user.save()

        # Auto-generate a unique parent username using UUID
        parent_username = f'parent_{uuid.uuid4().hex[:8]}'  # Generate an 8-character unique username

        # Generate a random secure password for the parent
        parent_password = get_random_string(12)  # Generate a 12-character password

        # Create the parent instance linked to the student user
        parent = Parent.objects.create(
            auto_generated_username=parent_username,
            auto_generated_password=parent_password,
            student_username=student_user.username  # Link to the student's username
        )

        # Send email with parent's auto-generated credentials to the student's email
        send_mail(
            subject='Registration Successful: Parent Login Credentials',
            message=f'You have successfully registered.\n\n'
                    f'Your parent\'s login details:\n'
                    f'Username: {parent_username}\n'
                    f'Password: {parent_password}\n',
            from_email=settings.DEFAULT_FROM_EMAIL,
            recipient_list=[email],  # Send to the student's email
            fail_silently=False,
        )

        messages.success(request, 'Account created successfully! Parent login credentials have been sent to your email.')
        return redirect('login')
    else:
        return render(request, 'register.html')



from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher, Parent, CustomUser

def login_view(request):
    if request.method == 'POST':
        username = request.POST.get('username')
        password = request.POST.get('password')

        # Check Teacher model
        try:
            teacher = Teacher.objects.get(auto_generated_username=username)
            if teacher.auto_generated_password == password:
                # Manually log in the teacher (using sessions)
                request.session['teacher_id'] = teacher.id  # Store teacher ID in session
                return redirect('teacher_dashboard')
        except Teacher.DoesNotExist:
            pass  # Not a teacher

        # Check Parent model
        try:
            parent = Parent.objects.get(auto_generated_username=username)
            if parent.auto_generated_password == password:
                # Manually log in the parent (using sessions)
                request.session['parent_id'] = parent.id  # Store parent ID in session
                return redirect('parent_dashboard')
        except Parent.DoesNotExist:
            pass  # Not a parent

        # Check CustomUser model
        try:
            custom_user = CustomUser.objects.get(username=username)
            if not custom_user.is_active:
                messages.error(request, 'Your account is deactivated. Please contact support.')
            if custom_user.check_password(password):  # Assuming password is stored as plaintext
                # Manually log in the CustomUser (using sessions)
                request.session['custom_user_id'] = custom_user.id  # Store CustomUser ID in session
                
                # Check if face data exists
                # face_data = StudentFaceData.objects.filter(user=custom_user).first()
                # if not face_data or not face_data.is_face_captured:
                #     return redirect('face_capture')
                
                return redirect('available_courses')
            else:
                messages.error(request, 'Invalid password for CustomUser.')
        except CustomUser.DoesNotExist:
            pass  # Not a CustomUser

        # If no match, show an error message
        messages.error(request, 'Invalid username or password.')

    return render(request, 'login.html')

def assigned_courses(request):
    # Get the teacher ID from the session
    teacher_id = request.session.get('teacher_id')

    if teacher_id:
        # Fetch the assigned courses from the session
        assigned_courses = request.session.get('assigned_courses', [])

        return render(request, 'assigned_courses.html', {'courses': assigned_courses})

    return redirect('teacher_login')  # Redirect to login if no session found


from django.shortcuts import render, redirect
from django.utils import timezone
from .models import Course, Enrollment

def available_courses(request):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        return redirect('login')  # Redirect to the login page if not authenticated
   
    # Get the logged-in user's ID from the session
    custom_user_id = request.session['custom_user_id']
    custom_user = CustomUser.objects.get(id=custom_user_id)
    # Get the current date
    current_date = timezone.now().date()

    # Filter courses to only include those that have not started and haven't reached their end date
    courses = Course.objects.filter(starting_date__gt=current_date, ending_date__gt=current_date)

    # Apply price range filter if set
    price_range = request.GET.get('price_range')
    if price_range:
        if price_range == 'all':
            pass
        elif price_range == '0-500':
            courses = courses.filter(price__lt=500)
        elif price_range == '500-999':
            courses = courses.filter(price__gte=500, price__lt=1000)

        elif price_range == '1000-2000':
            courses = courses.filter(price__gte=1000, price__lt=2000)
        elif price_range == '2000-3000':
            courses = courses.filter(price__gte=2000, price__lt=3000)
        elif price_range == '3000-4000':
            courses = courses.filter(price__gte=3000, price__lt=4000)
        elif price_range == '4000-5000':
            courses = courses.filter(price__gte=4000, price__lt=5000)
        elif price_range == '5000+':
            courses = courses.filter(price__gte=5000)

    # Prepare a list to hold course data along with enrollment status
    course_data = []
    for course in courses:
        # Check if the user is enrolled in the current course
        is_enrolled = Enrollment.objects.filter(student_id=custom_user_id, course=course).exists()
        
        # Append course data with enrollment status and start date check
        course_data.append({
            'course': course,
            'is_enrolled': is_enrolled,
            'can_enroll': current_date < course.starting_date  # Check if course hasn't started yet
        })

    # Pass the course data and current date to the template
    return render(request, 'available_courses.html', {
        'custom_user': custom_user,
        'course_data': course_data,
        'current_date': current_date,
        'selected_price_range': price_range,
        
    })

from django.shortcuts import render, redirect, get_object_or_404
from .models import Course, Enrollment
from django.http import HttpResponseForbidden

def enroll_course(request, course_id):
    # Check if the user is authenticated via session
    if 'custom_user_id' not in request.session:
        # If the user is not authenticated, redirect to the login page
        return redirect('login')  # Redirect to your custom login page

    # Get the authenticated user's ID from the session
    custom_user_id = request.session['custom_user_id']
    
    # Get the course the user is trying to enroll in
    course = get_object_or_404(Course, id=course_id)

    # Check if the user is already enrolled in this course
    if Enrollment.objects.filter(student_id=custom_user_id, course=course).exists():
        # You can add a message here indicating the user is already enrolled
        return redirect('available_courses')  # Redirect back to the available courses page

    # If the user is not enrolled, create a new Enrollment
    enrollment = Enrollment.objects.create(student_id=custom_user_id, course=course)

    # Optionally, you can show a success message or redirect to another page
    return redirect('available_courses')  # Redirect to the course list page

from datetime import datetime
def student_dashboard(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    # Fetch the CustomUser object based on the session ID
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
        if not custom_user.is_active:
            return redirect('login')
    except CustomUser.DoesNotExist:
        return redirect('login')

    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    today = datetime.today()

    # Check for unanswered feedback questions for the logged-in user
    feedback_questions = FeedbackQuestion.objects.all()
    answered_questions = Feedback.objects.filter(user=custom_user_id).values_list('question_id', flat=True)
    
    # Identify new feedback questions (not answered by the user)
    new_feedback_questions = feedback_questions.exclude(id__in=answered_questions)

    return render(request, 'student_dashboard.html', {
        'enrolled_courses': enrolled_courses,
        'custom_user': custom_user,
        'new_feedback_questions': new_feedback_questions,  # This will contain unanswered questions
        'today': today,  # Pass today's date to the template
    })


def teacher_dashboard(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher not authenticated

    # Fetch the Teacher object using teacher_id
    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect if the teacher doesn't exist
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).select_related('course')
    # Pass the teacher object to the template
    return render(request, 'teacher_dashboard.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'teacher_courses': teacher_courses
    })

def parent_dashboard(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    try:
        parent = Parent.objects.get(id=parent_id)
        student = CustomUser.objects.get(username=parent.student_username)
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Get student's quiz answers/grades
        grades = UserAnswers.objects.filter(
            user=student
        ).select_related(
            'question__quiz',
            'question__quiz__course'
        ).order_by('-attempt_date')
        
        # Get assignment submissions for the student
        assignment_submissions = AssignmentSubmission.objects.filter(
            student=student
        ).select_related(
            'assignment',
            'assignment__teacher',
            'assignment__course_name'
        ).order_by('-submitted_at')

        # Calculate submission status for each submission
        for submission in assignment_submissions:
            # Convert submitted_at to date for comparison
            submission_date = submission.submitted_at.date()
            due_date = submission.assignment.end_date
            
            if submission_date > due_date:
                submission.is_overdue = True
                submission.days_overdue = (submission_date - due_date).days
                submission.status_class = 'text-danger'
            elif submission_date == due_date:
                submission.is_overdue = False
                submission.status_class = 'text-warning'
                submission.submitted_on = 'last day'
            else:
                submission.is_overdue = False
                submission.days_early = (due_date - submission_date).days
                submission.status_class = 'text-success'

        # Get attendance records
        attendance_records = Attendance.objects.filter(student=student).order_by('-check_in_time')
        
        # Calculate attendance statistics
        total_records = attendance_records.count()
        present_count = attendance_records.filter(status='present').count()
        absent_count = total_records - present_count
        attendance_percentage = (present_count / total_records * 100) if total_records > 0 else 0

        # Get the student's enrolled courses
        courses = Enrollment.objects.filter(student=student).select_related('course')
        
        # Calculate total course fees from course prices
        total_course_fees = sum(enrollment.course.price for enrollment in courses)
        
        # Use select_related and prefetch_related to optimize queries
        teachers = Teacher.objects.filter(
            Q(parent_messages__parent=parent) | 
            Q(teacher_courses__course__students__username=parent.student_username)
        ).prefetch_related('teacher_courses').distinct()

        messages_list = ParentTeacherMessage.objects.filter(
            Q(parent=parent, message_type='parent_to_teacher') |  # Messages sent by parent
            Q(parent=parent, message_type='teacher_to_parent')    # Messages received by parent
        ).select_related('teacher').order_by('-date')
        
        context = {
            'parent': parent,
            'student': student,
            'courses': courses,
            'grades': grades,
            'assignment_submissions': assignment_submissions,
            'attendance_records': attendance_records,
            'present_count': present_count,
            'absent_count': absent_count,
            'attendance_percentage': attendance_percentage,
            'total_course_fees': total_course_fees,
            'messages_list': messages_list,
            'teachers': teachers,
            'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
        }

        try:
            # Get enrolled courses
            enrolled_courses = Enrollment.objects.filter(
                student__username=parent.student_username
            ).values_list('course', flat=True)

            # Get upcoming events for enrolled courses
            upcoming_events = CalendarEvent.objects.filter(
                course__in=enrolled_courses,
                end_time__gte=timezone.now()
            ).order_by('start_time')

            # Get registered events for the student
            registered_events = EventRegistration.objects.filter(
                user__username=parent.student_username,
                status='registered'
            ).select_related('event')

            try:
                # Get event suggestions for this parent
                event_suggestions = EventSuggestion.objects.filter(
                    parent=parent
                ).select_related('event').order_by('-created_at')

                context.update({
                    'upcoming_events': upcoming_events,
                    'registered_events': registered_events,
                    'event_suggestions': event_suggestions,  # Add suggestions to context
                })

                return render(request, 'parent_dashboard.html', context)

            except Exception as e:
                print(f"Error getting event suggestions: {str(e)}")
                return render(request, 'parent_dashboard.html', context)

        except Exception as e:
            print(f"Error in parent_dashboard: {str(e)}")
            return render(request, 'parent_dashboard.html', context)

    except Parent.DoesNotExist:
        return redirect('login')
    except CustomUser.DoesNotExist:
        student = None
        courses = []
        grades = []
        assignment_submissions = []

    messages_list = ParentTeacherMessage.objects.filter(
            Q(parent=parent, message_type='parent_to_teacher') |  # Messages sent by parent
            Q(parent=parent, message_type='teacher_to_parent')    # Messages received by parent
        ).select_related('teacher').order_by('-date')

    # Get unread count for notifications
    unread_count = messages_list.filter(
        is_read=False, 
        message_type='teacher_to_parent'
    ).count()

    context.update({
        'messages_page': messages_list,  # Changed from sent_messages to messages_page
        'unread_count': unread_count
    })

    return render(request, 'parent_dashboard.html', {
        'parent': parent,
        'student': student,
        'courses': courses,
        'grades': grades,
        'assignment_submissions': assignment_submissions,
        'attendance_records': attendance_records,
        'present_count': present_count,
        'absent_count': absent_count,
        'attendance_percentage': attendance_percentage,
        'messages_page': messages_list,  # Add messages to the final context
        'unread_count': unread_count
    })


def index_view(request):
    return render(request, 'index.html')

def about_view(request):
    return render(request, 'about.html')

def admissions_view(request):
    return render(request, 'admissions.html')


from django.shortcuts import render, redirect
from django.contrib import messages
from .forms import ContactMessageForm  # Adjust the import based on your structure

def contact_view(request):
    if request.method == 'POST':
        form = ContactMessageForm(request.POST)
        if form.is_valid():
            form.save()  # Save the form data to the database
            messages.success(request, 'Your message has been sent successfully!')
            return redirect('index')  # Redirect to the index page after submitting
        else:
            # This line ensures that the form with error messages is re-rendered
            messages.error(request, 'Please correct the errors below.')
    else:
        form = ContactMessageForm()

    return render(request, 'contact.html', {'form': form})  # Replace 'your_template.html' with your actual template name


from django.shortcuts import render, redirect
from .models import ContactMessage
from django.core.mail import send_mail
from django.conf import settings
from django.contrib import messages

# View to display all contact messages
def view_messages(request):
    messages_list = ContactMessage.objects.all()
    return render(request, 'view_messages.html', {'messages_list': messages_list})

def reply_message(request, message_id):
    message_obj = ContactMessage.objects.get(id=message_id)
    if request.method == 'POST':
        subject = 'Thank you for contacting us'
        body = 'We will contact you for more clarification regarding your query.'
        from_email = settings.DEFAULT_FROM_EMAIL
        recipient_list = [message_obj.email]

        send_mail(subject, body, from_email, recipient_list)

        # Update the replied status
        message_obj.replied = True
        message_obj.save()

        messages.success(request, 'Reply sent successfully!')
        return redirect('view_messages')
    return redirect('view_messages')



# views.py
from django.shortcuts import render
from .models import Course  # Make sure the Course model is imported
from django.utils import timezone

def courses_view(request):
    # Fetch courses that haven't reached the start date
    available_courses = Course.objects.filter(starting_date__gt=timezone.now())
    return render(request, 'courses.html', {'courses': available_courses})

def course_detail_10(request):
    return render(request, 'course_detail_10.html')

def course_detail_higher_secondary(request):
    return render(request, 'course_detail_higher_secondary.html')

def teachers_view(request):
    return render(request, 'teachers.html')

def recover_view(request):
    return render(request, 'recover.html')

def features(request):
    return render(request, 'features.html')

def parent(request):
    return render(request, 'parent.html')

from django.shortcuts import render, redirect
from django.contrib import messages

def admin_login(request):
    if request.method == 'POST':
        username = request.POST['username']
        password = request.POST['password']
        
        # Check against predefined credentials
        if username == 'admin' and password == 'admin123':
            # Create a custom session for the admin user
            request.session['is_admin'] = True
            return redirect('admin_dashboard')  # Redirect to admin dashboard
        else:
            messages.error(request, "Invalid username or password. Please try again.")
    
    return render(request, 'admin_login.html')
# views.py
from django.shortcuts import render
from .models import CustomUser, Teacher, Course, Enrollment  # Import your models

def admin_dashboard(request):
    # Count the number of users, teachers, and courses
    total_users = CustomUser.objects.count()
    total_teachers = Teacher.objects.count()
    total_courses = Course.objects.count()

    # Prepare data for course enrollment visualization
    courses = Course.objects.all()
    course_enrollment_data = []

    for course in courses:
        enrollment_count = Enrollment.objects.filter(course=course).count()
        course_enrollment_data.append({
            'course_name': course.course_name,
            'enrollment_count': enrollment_count
        })

    context = {
        'total_users': total_users,
        'total_teachers': total_teachers,
        'total_courses': total_courses,
        'course_enrollment_data': course_enrollment_data,  # Pass enrollment data to template
    }

    return render(request, 'admin_dashboard.html', context)

# views.py
from django.shortcuts import render
from .models import Course, Enrollment  # Import your models

def course_enrollment_view(request):
    # Get enrollment counts for each course
    course_data = []
    courses = Course.objects.all()
    for course in courses:
        registration_count = Enrollment.objects.filter(course=course).count()  # Count enrollments for each course
        course_data.append({
            'course_name': course.course_name,
            'registration_count': registration_count
        })

    context = {
        'course_data': course_data,  # Pass course data to the template
    }

    return render(request, 'course_enrollment.html', context)

from django.shortcuts import render
from .models import CustomUser

def manage_students(request):
    students = CustomUser.objects.select_related('course').filter(is_active=1)

    past_students = CustomUser.objects.filter(is_active=0)
    return render(request, 'manage_students.html', {'students': students, 'past_students':past_students})


from django.shortcuts import get_object_or_404, redirect
from django.contrib import messages
from .models import CustomUser

def delete_student(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if request.method == 'POST':
        student.is_active = 0  
        student.save()  
        messages.success(request, 'Student marked as active successfully.')
        return redirect('manage_students')


# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import CustomUser

def toggle_student_status(request, student_id):
    student = get_object_or_404(CustomUser, id=student_id)
    if student.is_active:
        student.is_active = 0
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        student.is_active = True
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    student.save()

    # Send email to the student
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',
        [student.email],
        fail_silently=False,
    )

    messages.success(request, f"Student '{student.username}' status updated successfully.")
    return redirect('manage_students')


from django.contrib.auth import logout as auth_logout
from django.shortcuts import redirect

def logout(request):
    auth_logout(request)
    request.session.flush()  
    return redirect('login')


def recover(request):
    return render(request, 'recover.html')
from django.contrib.auth.views import PasswordResetView
from django.urls import reverse_lazy

class CustomPasswordResetView(PasswordResetView):
    template_name = 'password_reset.html'
    email_template_name = 'password_reset_email.html'
    subject_template_name = 'password_reset_subject.txt'
    success_url = reverse_lazy('password_reset_done')
from django.contrib.auth.views import PasswordResetDoneView

class CustomPasswordResetDoneView(PasswordResetDoneView):
    template_name = 'password_reset_done.html'
from django.contrib.auth.views import PasswordResetConfirmView
from django.urls import reverse_lazy

class CustomPasswordResetConfirmView(PasswordResetConfirmView):
    template_name = 'password_reset_confirm.html'
    success_url = reverse_lazy('password_reset_complete')
from django.contrib.auth.views import PasswordResetCompleteView

class CustomPasswordResetCompleteView(PasswordResetCompleteView):
    template_name = 'password_reset_complete.html'


from django.shortcuts import render, redirect
from django.contrib import messages
from .models import Teacher

def register_teacher(request):
    if request.method == 'POST':
        # Extract data from the POST request
        first_name = request.POST.get('first_name')
        last_name = request.POST.get('last_name')
        gender = request.POST.get('gender')
        dob = request.POST.get('dob')
        email = request.POST.get('email')
        contact = request.POST.get('contact')
        address_line1 = request.POST.get('address_line1')
        address_line2 = request.POST.get('address_line2')
        city = request.POST.get('city')
        state = request.POST.get('state')
        zip_code = request.POST.get('zip_code')
        qualification = request.POST.get('qualification')
        experience = request.POST.get('experience')

        # Handle file uploads
        qualification_certificate = request.FILES.get('qualification_certificate')
        experience_certificate = request.FILES.get('experience_certificate', None)  # This file is optional

        # Simple validation check for required fields
        if not first_name or not last_name or not email:
            messages.error(request, "First name, last name, and email are required.")
        else:
            # Create a new Teacher object and set the fields
            teacher = Teacher(
                first_name=first_name,
                last_name=last_name,
                gender=gender,
                date_of_birth=dob,
                email=email,
                contact=contact,
                address_line1=address_line1,
                address_line2=address_line2,
                city=city,
                state=state,
                zip_code=zip_code,
                qualification=qualification,
                experience=experience,
                qualification_certificate=qualification_certificate,
                experience_certificate=experience_certificate,  # This file is optional
                status='pending'  # Set the initial status of the teacher to 'pending'
            )
            teacher.save()

            # Add a success message
            messages.success(request, "You have successfully registered! Please wait for your approval.")

            # Redirect to the login page or another relevant page after successful registration
            return redirect('login')

    return render(request, 'register_teacher.html')



from django.shortcuts import render, redirect
from .models import Teacher  # Adjust the import based on your app structure

def view_profile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        print("No teacher ID found in session.")  # Debugging line
        return render(request, 'view_profile.html', {'error': 'Profile not found.'})
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        print(f"Logged in teacher's username: {teacher.auto_generated_username}")  # Debugging line
        context = {'teacher': teacher}
    except Teacher.DoesNotExist:
        print("No matching Teacher found.")  # Debugging line
        context = {'error': 'Profile not found.'}
    
    return render(request, 'view_profile.html', context)

from django.core.mail import send_mail
from django.contrib.auth.models import User
from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher
import random
import string

def manage_teachers(request):
    pending_teachers = Teacher.objects.filter(status='pending')
    approved_teachers = Teacher.objects.filter(status='approved')
    
    context = {
        'pending_teachers': pending_teachers,
        'approved_teachers': approved_teachers,
    }
    return render(request, 'manage_teachers.html', context)

# views.py
from django.core.mail import send_mail
from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher  # Import your Teacher model

def toggle_teacher_status(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    
    if teacher.is_active:
        teacher.is_active = False  # Set to False to deactivate
        email_subject = "Account Deactivated"
        email_message = "Your account has been deactivated. Please contact support if you have any questions."
    else:
        teacher.is_active = True  # Set to True to activate
        email_subject = "Account Activated"
        email_message = "Your account has been activated. You can now log in to the platform."

    teacher.save()

    # Send email to the teacher
    send_mail(
        email_subject,
        email_message,
        'divyaantony2025@mca.ajce.in',  # Replace with your email
        [teacher.email],  # Send email to the teacher
        fail_silently=False,
    )

    messages.success(request, f"Teacher '{teacher.first_name} {teacher.last_name}' status updated successfully.")
    return redirect('manage_teachers')


from django.shortcuts import render, get_object_or_404, redirect
from django.core.mail import send_mail
from django.utils.crypto import get_random_string
from .models import Teacher, Course, TeacherCourse  # Assuming TeacherCourse is the model for the relationship
from .forms import ApproveTeacherForm

def approve_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    courses = Course.objects.all()
    
    if request.method == 'POST':
        form = ApproveTeacherForm(request.POST, instance=teacher)
        
        if form.is_valid():
            if teacher.status != 'approved':
                teacher.status = 'approved'
                
                # Process selected courses and teaching areas
                selected_courses = request.POST.getlist('courses')  # This retrieves a list of selected course IDs
                for course_id in selected_courses:
                    course = get_object_or_404(Course, id=course_id)
                    teaching_area = request.POST.get(f'teaching_area_{course_id}', '')  # Get the teaching area for each course
                    
                    # Create or update the TeacherCourse relationship
                    TeacherCourse.objects.create(teacher=teacher, course=course, teaching_area=teaching_area)

                teacher.save()

                # Debugging: Check if teacher is approved
                print(f"Teacher {teacher.first_name} {teacher.last_name} is now approved")
             
                # Generate random username and password
                random_username = f"teacher_{get_random_string(8)}"  
                random_password = get_random_string(8)  

                # Debugging: Check if credentials are generated
                print(f"Generated credentials: {random_username}, {random_password}")
                
                teacher.auto_generated_username = random_username
                teacher.auto_generated_password = random_password  
                teacher.save()

                # Send email to the teacher
                subject = 'Your Teacher Account has been Approved'
                message = (
                    f"Dear {teacher.first_name} {teacher.last_name},\n\n"
                    f"Your account has been approved! Here are your login details:\n\n"
                    f"Username: {random_username}\n"
                    f"Password: {random_password}\n\n"
                    f"You have been assigned to teach the course(s).\n"
                    f"Teaching Area: **{teaching_area}**.\n\n"
                    "Please log in and change your password upon your first login.\n\n"
                    "Best Regards,\n"
                    "The Administration Team"
                )

                # Send the email
                send_mail(
                    subject,
                    message,
                    'divyaantony2025@mca.ajce.in',  
                    [teacher.email],
                    fail_silently=False,
                )

                return redirect('admin_dashboard')  
            else:
                print("Teacher is already approved")
        else:
            print(form.errors)

    else:
        form = ApproveTeacherForm(instance=teacher)

    context = {
        'form': form,
        'teacher': teacher,
        'courses': courses,
    }

    return render(request, 'approve_teacher.html', context)


from django.shortcuts import render, redirect, get_object_or_404
from .models import Teacher

def reject_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        teacher.delete()  # Or you can set the status to 'pending'
        return redirect('manage_teachers')  # Redirect after rejection

    return render(request, 'reject_teacher.html', {'teacher': teacher})

def approving_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Update the teacher's status to 'approved'
        teacher.status = 'approved'
        teacher.save()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been approved.')

        return redirect('manage_teachers')  # Redirect to the management page after approval

    # If GET request, render the approve form
    form = None  # You can add additional form processing logic if required
    return render(request, 'approve_teacher.html', {'teacher': teacher, 'form': form})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import Teacher

# View to handle deleting a teacher
def delete_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    if request.method == 'GET':
        teacher.delete()
        messages.success(request, f'Teacher {teacher.first_name} {teacher.last_name} has been deleted.')
        return redirect('manage_teachers')


from django.shortcuts import render
from .models import Teacher

def teacher_list(request):
    teachers = Teacher.objects.all()  # Fetch all teachers from the database
    return render(request, 'teacher_list.html', {'teachers': teachers})


from django.shortcuts import render, redirect
from django.core.mail import send_mail
from django.contrib import messages
from .models import Teacher, TeacherInterview  # Import the Interview model

def interview_teacher(request):
    # Fetch only teachers with status "pending"
    teachers = Teacher.objects.filter(status='pending')

    if request.method == 'POST':
        # Get form data
        teacher_id = request.POST.get('teacher_id')
        interview_date = request.POST.get('interview_date')
        starting_time = request.POST.get('starting_time')
        ending_time = request.POST.get('ending_time')
        meeting_link = request.POST.get('meeting_link')
        interviewer_name = request.POST.get('interviewer_name')
        notes = request.POST.get('notes')

        # Fetch the teacher based on teacher_id
        try:
            teacher = Teacher.objects.get(id=teacher_id)
            first_name = teacher.first_name
            last_name = teacher.last_name
            teacher_email = teacher.email  # Retrieve teacher's email
        except Teacher.DoesNotExist:
            messages.error(request, 'Teacher not found!')
            return redirect('interview_teacher')

        # Store interview details in the database
        interview = TeacherInterview.objects.create(
            teacher=teacher,
            interview_date=interview_date,
            starting_time=starting_time,
            ending_time=ending_time,
            meeting_link=meeting_link,
            interviewer_name=interviewer_name,
            notes=notes
        )

        # Email content
        subject = "Interview Scheduled for Teacher"
        message = f"""
        Dear {first_name} {last_name},

        You are scheduled for an interview.

        Interview Details:
        Date: {interview_date}
        Starting Time: {starting_time}
        Ending Time: {ending_time}
        Meeting Link: {meeting_link}

        Interviewer: {interviewer_name}
        Notes: {notes}

        Please make sure to attend the meeting on time.

        Best regards,
        {interviewer_name}
        """

        # Sending the email
        send_mail(
            subject,
            message,
            'divyaantony2025@mca.ajce.in',  # From email (replace with your configuration)
            [teacher_email],  # Recipient email
            fail_silently=False,
        )

        messages.success(request, f'Interview scheduled successfully, and email sent to {teacher_email}!')
        return redirect('admin_dashboard')

    return render(request, 'interview_teacher.html', {'teachers': teachers})



from datetime import datetime, timedelta

def add_course(request):
    if request.method == 'POST':
        course_name = request.POST.get('course_name')
        description = request.POST.get('description')
        duration = request.POST.get('duration')  # Duration is 1 hour
        price = request.POST.get('price')
        image = request.FILES.get('image')
        start_date = request.POST.get('start_date')  # getting the start date
        
        # Convert start_date string to a datetime object
        start_date = datetime.strptime(start_date, "%Y-%m-%d")
        
        # Calculate the end_date by adding one year to the start_date
        end_date = start_date + timedelta(days=365)

        # Create and save the new course object
        new_course = Course(
            course_name=course_name, 
            description=description, 
            duration=int(duration), 
            price=price, 
            image=image,
            starting_date=start_date,
            ending_date=end_date
        )
        new_course.save()

        messages.success(request, 'Course added successfully!')
        return redirect('course_list')
    else:
        return render(request, 'add_courses.html')



def course_list(request):
    courses = Course.objects.all()  
    return render(request, 'course_list.html', {'courses': courses})

from django.shortcuts import render, redirect
from .models import Course, ClassSchedule, Teacher, TeacherCourse
from datetime import date, datetime

def schedule_class(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    if not teacher_id:
        return redirect('login')  # Redirect to login if teacher ID is missing
    
    try:
        teacher = Teacher.objects.get(id=teacher_id)  # Fetch the logged-in teacher
    except Teacher.DoesNotExist:
        return redirect('login')  # Redirect to login if teacher is not found

    # Fetch the assigned courses from TeacherCourse model
    assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
    
    if not assigned_courses.exists():
        error_message = "No course is assigned to you."
        return render(request, 'schedule_class.html', {
            'today': date.today(),
            'error_message': error_message
        })

    today = date.today()
    error_message = None

    if request.method == 'POST':
        class_name = request.POST.get('class_name')
        selected_date_str = request.POST.get('date')
        start_time_str = request.POST.get('start_time')
        end_time_str = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')
        selected_course_id = request.POST.get('assigned_course')  # Get selected course ID from dropdown

        # Convert strings to appropriate date and time objects
        try:
            selected_date = datetime.strptime(selected_date_str, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time_str, '%H:%M').time()
            end_time = datetime.strptime(end_time_str, '%H:%M').time()
        except ValueError:
            error_message = "Invalid date or time format. Please use the correct format."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Validate that the selected date is not in the past
        if selected_date < today:
            error_message = "The selected date cannot be in the past."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Fetch the selected course to validate dates
        course = Course.objects.get(id=selected_course_id)  # Fetch the course by ID
        
        # Validate that the scheduled date is within the course's start and end dates
        if selected_date < course.starting_date or selected_date > course.ending_date:
            error_message = f"The scheduled date must be between {course.starting_date} and {course.ending_date}."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Check that the end time is after the start time
        if end_time <= start_time:
            error_message = "End time must be after the start time."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })
        
        # Overlap validation - Check if the new class conflicts with any existing scheduled classes for any teacher
        overlapping_classes = ClassSchedule.objects.filter(
            date=selected_date,
            start_time__lt=end_time,  # An overlap if the class starts before the new one ends
            end_time__gt=start_time   # And ends after the new one starts
        )

        if overlapping_classes.exists():
            error_message = "The scheduled class overlaps with an existing class scheduled by another teacher."
            return render(request, 'schedule_class.html', {
                'assigned_courses': assigned_courses,
                'today': today,
                'error_message': error_message
            })

        # Create the class schedule record in the database
        schedule = ClassSchedule(
            class_name=class_name,
            course_name=course,  # Use the fetched course
            date=selected_date,
            start_time=start_time,
            end_time=end_time,
            meeting_link=meeting_link,
            teacher=teacher
        )
        schedule.save()

        if schedule:
            return redirect('view_teacher_schedule_class')  # Redirect to teacher's dashboard on success

    # Fetch scheduled classes for this teacher
    current_datetime = datetime.now()
    scheduled_classes = ClassSchedule.objects.filter(
        teacher=teacher, date__gte=today, end_time__gt=current_datetime.time()
    )

    return render(request, 'schedule_class.html', {
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
        'assigned_courses': assigned_courses,
        'today': today,
        'scheduled_classes': scheduled_classes,
        'error_message': error_message
    })


from django.shortcuts import render, redirect
from django.utils import timezone
from django.db.models import Q
from .models import ClassSchedule, Teacher, Course

def view_teacher_schedule_class(request):
    # Get the teacher_id from session
    teacher_id = request.session.get('teacher_id')

    # Redirect to login if teacher_id is not present
    if not teacher_id:
        return redirect('login')

    try:
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        return redirect('login')

    # Get the current local time and date, automatically converted to IST if TIME_ZONE is set to 'Asia/Kolkata'
    current_datetime = timezone.localtime()

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")
    print(f"Current date (IST): {current_datetime.date()}, Current time (IST): {current_datetime.time()}")

    # Fetch ongoing and future classes based on the current IST date and time
    future_classes = ClassSchedule.objects.filter(
        teacher=teacher  # Filter classes by teacher
    ).filter(
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (scheduled after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Debug statement to check the ongoing and future classes query results
    print(f"Ongoing and future classes: {future_classes}")

    # Prepare the context for rendering
    context = {
        'future_classes': future_classes,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,  # Ongoing and future classes
    }

    return render(request, 'view_teacher_schedule_class.html', context)



from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from .models import ClassSchedule

def edit_class(request):
    courses = Course.objects.all()
    if request.method == 'POST':
        class_id = request.POST.get('class_id')
        course_id = request.POST.get('course_name')
        class_name = request.POST.get('class_name')
        date = request.POST.get('date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        meeting_link = request.POST.get('meeting_link')

        try:
            scheduled_class = ClassSchedule.objects.get(id=class_id)
            scheduled_class.course_name = course_id
            scheduled_class.class_name = class_name
            scheduled_class.date = date
            scheduled_class.start_time = start_time
            scheduled_class.end_time = end_time
            scheduled_class.meeting_link = meeting_link
            scheduled_class.save()

            messages.success(request, "Class details updated successfully.")
        except ClassSchedule.DoesNotExist:
            messages.error(request, "Class not found.")

        return redirect('view_teacher_schedule_class', {'courses': courses})  # Redirect to the view scheduled classes page

from django.contrib import messages # type: ignore
from django.shortcuts import redirect, render # type: ignore
from .models import ClassSchedule, CustomUser,Enrollment
from django.utils import timezone # type: ignore
from django.db.models import Q # type: ignore

def view_scheduled_classes(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
# Fetch the enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student)

    if not enrolled_courses.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    # Get the current date and time in IST
    current_datetime = timezone.localtime()  # This will automatically convert to IST if TIME_ZONE is set to 'Asia/Kolkata'

    # Debug statement to check the current IST date and time
    print(f"Current date and time (IST): {current_datetime}")

    # Fetch ongoing and future classes for the registered course
    ongoing_future_classes = ClassSchedule.objects.filter(
       course_name__in=enrolled_courses.values_list('course', flat=True)  # Correct field name here
    ).filter( 
        Q(date=current_datetime.date(), end_time__gt=current_datetime.time()) |  # Ongoing classes today that haven't ended
        Q(date__gt=current_datetime.date())  # Future classes (after today)
    ).order_by('date', 'start_time')  # Order by date and start time

    # Check if there are any ongoing or future classes
    if not ongoing_future_classes.exists():
        messages.info(request, "No ongoing or upcoming classes found.")

    # Pass the classes to the template
    return render(request, 'view_scheduled_classes.html', {'scheduled_classes': ongoing_future_classes})


from django.shortcuts import render, redirect # type: ignore
from .models import CustomUser, ClassSchedule, Parent

def view_class_schedule(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        
        return redirect('login')
    
    try:
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
       
        return redirect('error_page')

    child_username = parent.student_username
    child = CustomUser.objects.get(username=child_username)
    child_schedule = ClassSchedule.objects.filter(course_name=child.course)

    context = {
        'child': child,
        'child_schedule': child_schedule,
    }
    
    return render(request, 'view_class_schedule.html', context)


from .models import CustomUser, Parent, Teacher
from .models import Parent


def change_password(request):
    parent_id = request.session.get('parent_id')
    
    # Check if parent is logged in
    if not parent_id:
        messages.error(request, 'You are not logged in as a parent.')
        return redirect('login')

    try:
        # Fetch the parent instance
        parent = Parent.objects.get(id=parent_id)
    except Parent.DoesNotExist:
        messages.error(request, 'Parent not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('change_password')

        # Check if the old password is correct
        if old_password != parent.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('change_password')

        # Update the password in the database
        parent.auto_generated_password = new_password1
        parent.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'change_password.html')


from django.shortcuts import render, redirect # type: ignore
from django.contrib import messages # type: ignore
from django.contrib.auth import update_session_auth_hash # type: ignore
from .models import CustomUser, Parent, Teacher
from .models import Teacher

def teacher_changepassword(request):
    teacher_id = request.session.get('teacher_id')
    
    # Check if teacher is logged in
    if not teacher_id:
        messages.error(request, 'You are not logged in as a teacher.')
        return redirect('login')

    try:
        # Fetch the teacher instance
        teacher = Teacher.objects.get(id=teacher_id)
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found.')
        return redirect('login')

    if request.method == 'POST':
        old_password = request.POST.get('old_password')
        new_password1 = request.POST.get('new_password')
        new_password2 = request.POST.get('confirm_password')

        # Check if new passwords match
        if new_password1 != new_password2:
            messages.error(request, 'New passwords do not match.')
            return redirect('teacher_changepassword')

        # Check if the old password is correct
        if old_password != teacher.auto_generated_password:
            messages.error(request, 'Old password is incorrect.')
            return redirect('teacher_changepassword')

        # Update the password in the database
        teacher.auto_generated_password = new_password1
        teacher.save()
        messages.success(request, 'Your password has been successfully updated.')
        return redirect('login')

    return render(request, 'teacher_changepassword.html')



from .models import Teacher

def teacher_updateprofile(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        messages.error(request, 'No teacher ID found in session.')
        return redirect('login')

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        # Fetch the form data
        teacher.first_name = request.POST.get('first_name')
        teacher.last_name = request.POST.get('last_name')
        teacher.gender = request.POST.get('gender')
        teacher.date_of_birth = datetime.strptime(request.POST.get('date_of_birth'), '%Y-%m-%d').date()
        teacher.email = request.POST.get('email')
        teacher.contact = request.POST.get('contact')
        teacher.address_line1 = request.POST.get('address_line1')
        teacher.address_line2 = request.POST.get('address_line2')
        teacher.city = request.POST.get('city')
        teacher.state = request.POST.get('state')
        teacher.zip_code = request.POST.get('zip_code')
        teacher.qualification = request.POST.get('qualification')
        teacher.experience = request.POST.get('experience')

        # Handle file uploads
        if 'qualification_certificate' in request.FILES:
            teacher.qualification_certificate = request.FILES['qualification_certificate']
        if 'experience_certificate' in request.FILES:
            teacher.experience_certificate = request.FILES['experience_certificate']
        # Save the updated teacher details
        teacher.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('view_profile')  

    context = {'teacher': teacher,
               'first_name': teacher.first_name,
               'last_name': teacher.last_name,
        }
    return render(request, 'teacher_updateprofile.html', context)


from .models import Parent

def parent_update_profile(request):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        messages.error(request, 'No parent ID found in session.')
        return redirect('login')

    parent = get_object_or_404(Parent, id=parent_id)

    if request.method == 'POST':
        # Fetch the form data
        parent.auto_generated_username = request.POST.get('auto_generated_username')
        

        # Save the updated teacher details
        parent.save()
        messages.success(request, 'Profile updated successfully.')
        return redirect('parent_dashboard')  

    context = {'parent': parent}
    return render(request, 'parent_update_profile.html', context)
from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.contrib import messages  # For displaying messages
from .models import Material, Course, Teacher, TeacherCourse  # Import the TeacherCourse model

def upload_material(request):
    if request.method == 'POST':
        teacher_id = request.session.get('teacher_id')
        
        if not teacher_id:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'success': False, 'error': 'Not logged in as teacher'})
            else:
                messages.error(request, 'You are not logged in as a teacher.')
                return redirect('login')

        try:
            # Check if this is a note upload (AJAX request)
            note_id = request.POST.get('note_id')
            if note_id:
                # Get the note
                note = TeacherNote.objects.get(id=note_id, teacher_id=teacher_id)
                course = note.course

                # Create PDF from note content
                buffer = BytesIO()
                p = canvas.Canvas(buffer)
                
                # Add title
                p.setFont("Helvetica-Bold", 16)
                p.drawString(100, 800, note.title)
                
                # Add content
                p.setFont("Helvetica", 12)
                y = 750
                for line in note.content.split('\n'):
                    if y > 50:  # Ensure we don't write below page margin
                        p.drawString(100, y, line.strip())
                        y -= 20
                
                p.showPage()
                p.save()
                
                # Create PDF file
                pdf_file = ContentFile(buffer.getvalue())
                
                # Create material
                material = Material.objects.create(
                    teacher_id=teacher_id,
                    course=course,
                    description=f"Note: {note.title}",
                )
                
                # Save the PDF file
                material.file.save(f'note_{note.id}.pdf', pdf_file)

                return JsonResponse({'success': True})
            
            # Handle regular file upload
            else:
                course_id = request.POST.get('course')
                description = request.POST.get('description')
                file = request.FILES.get('file')

                # Verify course assignment
                assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
                if int(course_id) not in assigned_courses:
                    messages.error(request, 'You do not have permission to upload material for this course.')
                    return redirect('upload_material')

                if file:
                    course = get_object_or_404(Course, id=course_id)
                    teacher = get_object_or_404(Teacher, id=teacher_id)

                    material = Material.objects.create(
                        teacher=teacher,
                        course=course,
                        file=file,
                        description=description
                    )
                    messages.success(request, 'Material uploaded successfully!')
                    return redirect('teacher_dashboard')
                else:
                    messages.error(request, 'Please upload a file.')
                    return redirect('upload_material')

        except Exception as e:
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return JsonResponse({'success': False, 'error': str(e)})
            else:
                messages.error(request, f'Error saving material: {str(e)}')
                return redirect('upload_material')

    # GET request - show the upload form
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        messages.error(request, 'You must be logged in as a teacher.')
        return redirect('login')

    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)
    return render(request, 'upload_material.html', {'courses': courses})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.contrib import messages  # type: ignore
from .models import Material, CustomUser, Enrollment
from datetime import datetime
from django.utils import timezone

def view_materials(request):
    # Check if a CustomUser (student) is logged in by checking the session
    custom_user_id = request.session.get('custom_user_id')
    
    if not custom_user_id:
        messages.error(request, 'You must be logged in as a student to view materials.')
        return redirect('login')  # Redirect to login page if no session

    # Fetch the CustomUser (student) object using the session ID
    custom_user = get_object_or_404(CustomUser, id=custom_user_id)

    # Ensure the student is registered for any courses
    enrollments = Enrollment.objects.filter(student=custom_user)
    
    if not enrollments:
        messages.error(request, 'You are not registered for any course.')
        return redirect('student_dashboard')  # Redirect if no course is associated

    # Prepare a list to hold materials that meet the enrollment criteria
    materials = []
    
    # Loop through each enrollment to filter materials
    for enrollment in enrollments:
        # Combine enrollment_date and enrollment_time into a single datetime object
        if enrollment.enrollment_time:
            enrollment_datetime = datetime.combine(enrollment.enrollment_date, enrollment.enrollment_time)
        else:
            # Fallback if time is not present (you may adjust this as needed)
            enrollment_datetime = timezone.make_aware(datetime.combine(enrollment.enrollment_date, datetime.min.time()))
        
        # Make the combined datetime timezone-aware
        enrollment_datetime = timezone.make_aware(enrollment_datetime)
        
        # Fetch materials related to the course and filter by the enrollment datetime
        course_materials = Material.objects.filter(course=enrollment.course, uploaded_at__gte=enrollment_datetime)
        materials.extend(course_materials)

    # Render the materials in the template
    return render(request, 'view_materials.html', {'materials': materials})

# views.py
from .models import Material, Parent, CustomUser

def view_study_materials(request):

    parent = get_object_or_404(Parent, auto_generated_username=request.user.username)
    
    student = get_object_or_404(CustomUser, username=parent.student_username)

    materials = Material.objects.filter(course=student.course)
    
    return render(request, 'view_study_materials.html', {'materials': materials, 'student': student})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.utils import timezone  # type: ignore
from .models import Quizs, Course, TeacherCourse  # Import TeacherCourse for filtering

def create_quiz(request):
    teacher_id = request.session.get('teacher_id')  # Retrieve the teacher's ID from the session

    if not teacher_id:
        # Handle the case where the teacher is not logged in
        return redirect('login')

    # Fetch courses taught by the teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)  # Filter to only assigned courses

    if request.method == 'POST':
        course_id = request.POST.get('course')
        title = request.POST.get('title')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        duration = request.POST.get('duration')

        course = get_object_or_404(Course, id=course_id)
        
        # Create and save the quiz
        quiz = Quizs(course=course, teacher_id=teacher_id, title=title, 
                     start_date=start_date, end_date=end_date, 
                     start_time=start_time, end_time=end_time,duration=duration)
        quiz.save()

        return redirect('add_question', quiz_id=quiz.id)
    
    return render(request, 'create_quiz.html', {'courses': courses})


def add_question(request, quiz_id):
    quiz = get_object_or_404(Quizs, id=quiz_id)

    if request.method == 'POST':
        questions_data = request.POST.getlist('question')
        options_a = request.POST.getlist('option_a')
        options_b = request.POST.getlist('option_b')
        options_c = request.POST.getlist('option_c')
        options_d = request.POST.getlist('option_d')
        correct_options = request.POST.getlist('correct_option')

        # Iterate through all submitted questions
        for i in range(len(questions_data)):
            question_text = questions_data[i]
            option_a = options_a[i]
            option_b = options_b[i]
            option_c = options_c[i]
            option_d = options_d[i]
            correct_option = correct_options[i]

            # Create a new Question object for each set of data
            question = Question(
                quiz=quiz, 
                text=question_text,
                option_a=option_a,
                option_b=option_b,
                option_c=option_c,
                option_d=option_d,
                correct_option=correct_option
            )
            question.save()

        return redirect('add_question', quiz_id=quiz.id)

    # Add this return statement for GET requests
    return render(request, 'add_question.html', {
        'quiz': quiz,
        'quiz_id': quiz_id,
        'first_name': request.session.get('first_name', ''),
        'last_name': request.session.get('last_name', '')
    })

from django.shortcuts import render, get_object_or_404
from .models import Quizs, Question

def take_quiz(request, quiz_id):
    quiz = get_object_or_404(Quizs, id=quiz_id)
    questions = Question.objects.filter(quiz=quiz)

    if request.method == 'POST':
        user_answers = {}  # Store user's answers
        correct_count = 0  # Track the number of correct answers

        # Loop through each question in the quiz
        for question in questions:
            user_answer = request.POST.get(f'question_{question.id}')  # Get user's answer for the question
            user_answers[question.id] = user_answer  # Save the user's answer

            # Check if the user's answer matches the correct answer
            if user_answer == question.correct_option:
                correct_count += 1  # Increment correct count if answer is correct

        # After submission, show the result
        return render(request, 'quiz_result.html', {
            'quiz': quiz,
            'questions': questions,
            'user_answers': user_answers,
            'correct_count': correct_count,
            'total_questions': questions.count()
        })

    # Render the quiz page with questions
    return render(request, 'take_quiz.html', {
        'quiz': quiz,
        'questions': questions
    })

from django.shortcuts import render, redirect
from django.contrib import messages
from django.utils import timezone
from django.db.models import Q
from .models import CustomUser, Quizs, UserAnswers, Enrollment  # Ensure Enrollment is imported

def available_quizzes(request):
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    # Fetch enrollments for the student
    enrollments = Enrollment.objects.filter(student=student)

    if not enrollments.exists():
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    current_datetime = timezone.localtime()

    # Create a list of enrolled courses with their enrollment dates and times
    enrolled_courses = [
        (enrollment.course, enrollment.enrollment_date, enrollment.enrollment_time)
        for enrollment in enrollments
    ]

    available_quizzes = []

    # Fetch quizzes that are ongoing or upcoming and have not ended
    for course, enrollment_date, enrollment_time in enrolled_courses:
        # Combine enrollment date and time to create a full datetime object
        enrollment_datetime = timezone.make_aware(datetime.combine(enrollment_date, enrollment_time or datetime.min.time()))

        # Fetch quizzes based on the course and check the date and time constraints
        quizzes_for_course = Quizs.objects.filter(
            course=course,
            start_date__lte=current_datetime.date(),
            end_date__gte=current_datetime.date(),
        ).filter(
            Q(start_time__lte=current_datetime.time()) | Q(start_date__gt=current_datetime.date())
        ).exclude(
            Q(end_date__lt=current_datetime.date()) | 
            (Q(end_date=current_datetime.date()) & Q(end_time__lt=current_datetime.time()))
        ).order_by('start_date', 'start_time')

        # Filter out quizzes that the student has already attempted
        quizzes_for_course = [quiz for quiz in quizzes_for_course if not UserAnswers.objects.filter(user=student, question__quiz=quiz).exists()]

        # Add the filtered quizzes to the available_quizzes list
        available_quizzes.extend(quizzes_for_course)

    return render(request, 'available_quizzes.html', {
        'quizzes': available_quizzes,
        'today': current_datetime.date(),
        'current_time': current_datetime.time()
    })



def submit_quiz(request, quiz_id):
    # Get the quiz and its associated questions
    quiz = get_object_or_404(Quizs, id=quiz_id)
    questions = quiz.questions.all()

    # Retrieve the `custom_user_id` from the session
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You must be logged in to submit the quiz.")
        return redirect('login')

    # Fetch the CustomUser based on the custom user ID from the session
    try:
        custom_user = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "User not found.")
        return redirect('login')

    # Initialize counters
    correct_count = 0
    total_questions = questions.count()

    # Prepare to store user answers with questions
    question_data = []

    # Loop through the quiz questions
    for question in questions:
        # Get the user's selected answer from POST data
        user_answer = request.POST.get(f'question_{question.id}')

        # Only process if the user provided an answer
        if user_answer in ['A', 'B', 'C', 'D']:
            # Store the user's answer in the UserAnswers table
            UserAnswers.objects.update_or_create(
                user=custom_user,
                question=question,
                defaults={'selected_option': user_answer}
            )

            # Increment correct count if the answer is correct
            if user_answer == question.correct_option:
                correct_count += 1
        else:
            user_answer = None

        # Retrieve the stored answer from the database
        try:
            stored_answer = UserAnswers.objects.get(user=custom_user, question=question).selected_option
        except UserAnswers.DoesNotExist:
            stored_answer = "Not answered"

        # Append the question and user answer to the list
        question_data.append({
            'question': question,
            'user_answer': stored_answer,
            'correct_answer': question.correct_option,
        })

    # Render the result page, passing the structured data to the template
    context = {
        'quiz': quiz,
        'questions_data': question_data,
        'correct_count': correct_count,
        'total_questions': total_questions,
    }
    return render(request, 'quiz_result.html', context)


def quiz_result(request, quiz_id):
    # Ensure the user is logged in
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You must be logged in to view quiz results.")
        return redirect('login')

    # Get the quiz
    quiz = get_object_or_404(Quizs, id=quiz_id)

    # Retrieve user's answers for the specific quiz
    user_answers = UserAnswers.objects.filter(user_id=custom_user_id, question__quiz=quiz)

    # Prepare data for display
    results = []
    correct_count = 0
    total_questions = quiz.questions.count()

    # Create a dictionary to map question ids to user answers
    user_answers_dict = {user_answer.question.id: user_answer.selected_option for user_answer in user_answers}

    for question in quiz.questions.all():  # Iterate through all questions in the quiz
        selected_option = user_answers_dict.get(question.id, "Not answered")  # Get the selected option or default to "Not answered"
        correct_answer = question.correct_option
        is_correct = (selected_option == correct_answer)

        results.append({
            'question': question,
            'selected_option': selected_option,
            'correct_option': correct_answer,
            'is_correct': is_correct,  # Keep track of whether the answer is correct
        })

        if is_correct:
            correct_count += 1

    context = {
        'quiz': quiz,
        'results': results,
        'correct_count': correct_count,
        'total_questions': total_questions,
    }
    return render(request, 'quiz_result.html', context)


def quiz_questions(request, quiz_id):
    quiz = Quizs.objects.get(id=quiz_id)
    questions = Question.objects.filter(quizs=quiz)  # Adjust based on your models
    return render(request, 'quiz_questions.html', {'quiz': quiz, 'questions': questions})

from django.shortcuts import render, redirect, get_object_or_404  # type: ignore
from django.http import HttpResponse  # type: ignore
from .models import Assignment, Course, Teacher, TeacherCourse  # Import the TeacherCourse model
from django.core.exceptions import ValidationError  # type: ignore
from datetime import datetime

def create_assignment(request):
    if request.method == 'POST':
        # Fetch form data
        title = request.POST.get('title')
        description = request.POST.get('description')
        start_date = request.POST.get('start_date')
        end_date = request.POST.get('end_date')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        course_name_id = request.POST.get('course')
        file = request.FILES.get('file')

        # Validate and convert date and time
        try:
            start_date = datetime.strptime(start_date, '%Y-%m-%d').date()
            end_date = datetime.strptime(end_date, '%Y-%m-%d').date()
            start_time = datetime.strptime(start_time, '%H:%M').time()
            end_time = datetime.strptime(end_time, '%H:%M').time()

            teacher_id = request.session.get('teacher_id')
            teacher = get_object_or_404(Teacher, id=teacher_id)

            assignment = Assignment(
                title=title,
                description=description,
                start_date=start_date,
                end_date=end_date,
                start_time=start_time,
                end_time=end_time,
                file=file,
                course_name_id=course_name_id,
                teacher=teacher
            )

            if assignment.start_date > assignment.end_date:
                raise ValidationError("Start date cannot be after end date.")

            assignment.save()

            return redirect('teacher_dashboard')

        except ValidationError as e:
            return render(request, 'create_assignment.html', {
                'error': str(e),
                'courses': courses,  # Pass the courses to the template
                'assignment': request.POST
            })
        except Exception as e:
            return render(request, 'create_assignment.html', {
                'error': str(e),
                'courses': courses,
                'assignment': request.POST
            })

    # If not a POST request, get the assigned courses for the teacher
    teacher_id = request.session.get('teacher_id')
    teacher = get_object_or_404(Teacher, id=teacher_id)

    assigned_courses = TeacherCourse.objects.filter(teacher=teacher).select_related('course')
    courses = [tc.course for tc in assigned_courses]

    return render(request, 'create_assignment.html', {
        'courses': courses,
        'first_name': teacher.first_name,
        'last_name': teacher.last_name,
    })

from django.contrib import messages
from django.shortcuts import redirect, render
from django.utils import timezone
from datetime import datetime
from .models import Assignment, AssignmentSubmission, CustomUser, Enrollment  # Ensure to import Enrollment

def assignment_submission_view(request):
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    enrollments = Enrollment.objects.filter(student=student)

    if not enrollments:
        messages.error(request, "You are not registered for any course.")
        return redirect('student_dashboard')

    enrolled_courses = [enrollment.course for enrollment in enrollments]
    assignment_details = []

    # Handle file upload
    if request.method == 'POST':
        assignment_id = request.POST.get('assignment_id')
        file = request.FILES.get('file')
        
        if not file:
            messages.error(request, 'Please select a file to upload.')
            return redirect('assignment_detail')
            
        if not assignment_id:
            messages.error(request, 'Invalid assignment selection.')
            return redirect('assignment_detail')

        try:
            assignment = Assignment.objects.get(
                id=assignment_id, 
                course_name__in=enrolled_courses
            )
            
            # Check file type
            file_name = file.name.lower()
            if not file_name.endswith(('.pdf', '.doc', '.docx')):
                messages.error(request, 'Please upload only PDF or Word documents.')
                return redirect('assignment_detail')

            # Check file size (10MB limit)
            if file.size > 10 * 1024 * 1024:
                messages.error(request, 'File size should be less than 10MB.')
                return redirect('assignment_detail')

            # Check deadline
            current_datetime = timezone.now()
            assignment_end_datetime = timezone.make_aware(
                datetime.combine(assignment.end_date, assignment.end_time)
            )
            
            if assignment_end_datetime <= current_datetime:
                messages.error(request, "Assignment deadline has passed.")
                return redirect('assignment_detail')

            # Update or create submission
            try:
                submission = AssignmentSubmission.objects.get(
                    assignment=assignment,
                    student=student
                )
                submission.file = file
                submission.save()
                messages.success(request, 'Your submission has been updated successfully!')
            except AssignmentSubmission.DoesNotExist:
                AssignmentSubmission.objects.create(
                    assignment=assignment,
                    student=student,
                    file=file
                )
                messages.success(request, 'Your submission has been uploaded successfully!')
            
            return redirect('assignment_detail')

        except Assignment.DoesNotExist:
            messages.error(request, 'Invalid assignment.')
            return redirect('assignment_detail')
        except Exception as e:
            messages.error(request, f'Error submitting assignment: {str(e)}')
            return redirect('assignment_detail')

    # Prepare assignment details for display
    for enrollment in enrollments:
        enrollment_datetime = timezone.make_aware(
            datetime.combine(
                enrollment.enrollment_date,
                enrollment.enrollment_time or datetime.min.time()
            )
        )

        assignments = Assignment.objects.filter(
            course_name=enrollment.course,
            start_date__gte=enrollment.enrollment_date
        )

        for assignment in assignments:
            current_datetime = timezone.now()
            assignment_end_datetime = timezone.make_aware(
                datetime.combine(assignment.end_date, assignment.end_time)
            )

            submission = AssignmentSubmission.objects.filter(
                assignment=assignment,
                student=student
            ).first()

            assignment_details.append({
                'assignment': assignment,
                'submission': submission,
                'has_reached_deadline': assignment_end_datetime <= current_datetime,
                'can_submit': assignment_end_datetime > current_datetime and not submission,
                'submission_allowed': assignment_end_datetime > current_datetime
            })

    # Sort assignments by start date (newest first)
    assignment_details.sort(key=lambda x: x['assignment'].start_date, reverse=True)

    return render(request, 'assignment_detail.html', {
        'assignment_details': assignment_details
    })


from django.shortcuts import render, redirect  # type: ignore
from django.utils import timezone
from .models import Assignment, Course

def view_assignment(request):
    if 'teacher_id' in request.session:
        teacher_id = request.session['teacher_id']

        assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
        courses = Course.objects.filter(id__in=assigned_courses)

        # Get the selected course ID from the GET request
        selected_course_id = request.GET.get('course_id')

        # Fetch assignments for the selected course or all assignments if no course is selected
        if selected_course_id:
            assignments = Assignment.objects.filter(course_name__id=selected_course_id, teacher_id=teacher_id)
        else:
            assignments = Assignment.objects.filter(teacher_id=teacher_id)

        # Get the current date and time
        current_datetime = timezone.localtime()

        # Add a 'status' attribute to each assignment based on end date and time
        for assignment in assignments:
            if (assignment.end_date < current_datetime.date() or
                (assignment.end_date == current_datetime.date() and assignment.end_time <= current_datetime.time())):
                assignment.status = 'Completed'
            else:
                assignment.status = 'Ongoing'

        return render(request, 'view_assignment.html', {
            'assignments': assignments,
            'courses': courses,
            
            'selected_course_id': selected_course_id,
            
        })
    else:
        return redirect('login')



from django.db.models import Subquery, OuterRef  # type: ignore
from django.shortcuts import render  # type: ignore
from .models import AssignmentSubmission, Course, TeacherCourse  # Import the TeacherCourse model

def evaluate_assignments(request):
    # Check if the teacher is logged in
    if 'teacher_id' in request.session:
        teacher_id = request.session['teacher_id']
        selected_course_id = request.GET.get('course_id')

        # Fetch all courses assigned to the teacher from the TeacherCourse table
        assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
        courses = Course.objects.filter(id__in=assigned_courses)

        # Get the latest submission for each student
        latest_submissions = AssignmentSubmission.objects.filter(
            student=OuterRef('student'),
            assignment=OuterRef('assignment')
        ).order_by('-submitted_at')

        if selected_course_id:
            # Filter by course and ensure we only get the latest submission for each student
            submissions = AssignmentSubmission.objects.filter(
                assignment__course_name__id=selected_course_id,
                assignment__course_name__id__in=assigned_courses,  # Ensure the course is assigned to the teacher
                submitted_at=Subquery(latest_submissions.values('submitted_at')[:1])
            )
        else:
            # Get the latest submission for each student across all assigned courses
            submissions = AssignmentSubmission.objects.filter(
                assignment__course_name__id__in=assigned_courses,  # Ensure the course is assigned to the teacher
                submitted_at=Subquery(latest_submissions.values('submitted_at')[:1])
            )

        context = {
            'submissions': submissions,
            'courses': courses,
            'selected_course_id': selected_course_id,
            
            
        }
        return render(request, 'evaluate_assignment.html', context)
    else:
        return redirect('login')  # Redirect to login if the teacher is not authenticated

from django.shortcuts import get_object_or_404, redirect # type: ignore
from .models import AssignmentSubmission
from django.views.decorators.http import require_POST # type: ignore


@require_POST # type: ignore
def submit_grade(request, submission_id):
    submission = get_object_or_404(AssignmentSubmission, id=submission_id)
    grade = request.POST.get('grade')

    if grade is not None:
        submission.grade = grade  # Store the grade in the grade field
        submission.save()  # Save the submission to update the database

    return redirect('evaluate_assignment')

from django.shortcuts import render, redirect
from .models import FeedbackQuestion, Course

def add_feedback_question(request):
    courses = Course.objects.all()  # Fetch all courses from the Course table

    if request.method == 'POST':
        questions = []
        course_id = request.POST.get('course')  # Retrieve selected course ID
        release_date = request.POST.get('release_date')
        end_date = request.POST.get('end_date')

        course = Course.objects.get(id=course_id) if course_id else None

        for key, value in request.POST.items():
            if key.startswith('question_text_'):
                questions.append(value)

        if questions and release_date and end_date:
            for question_text in questions:
                # Save each question with an optional associated course, release date, and end date
                FeedbackQuestion.objects.create(
                    question_text=question_text,
                    course=course,
                    release_date=release_date,
                    end_date=end_date
                )
            return redirect('admin_dashboard')
        else:
            return render(request, 'add_feedback_question.html', {
                'courses': courses,
                'error': 'Please enter at least one question, and specify release and end dates.'
            })

    return render(request, 'add_feedback_question.html', {'courses': courses})




from django.shortcuts import render, redirect
from datetime import date
from .models import Feedback, FeedbackQuestion, Course, Enrollment  # Ensure your enrollment model is imported

def feedback_view(request):
    user_id = request.session.get('custom_user_id', 'Anonymous')  # Get the user ID from the session

    # Get the current date
    today = date.today()

    # Get the list of courses the student is enrolled in by querying the StudentCourseEnrollment model
    enrolled_courses = Enrollment.objects.filter(student_id=user_id).values_list('course_id', flat=True)

    # Get a list of questions that the user has already answered
    answered_questions = Feedback.objects.filter(user=user_id).values_list('question_id', flat=True)

    # Filter the feedback questions based on the enrolled courses, release date, and end date
    questions = FeedbackQuestion.objects.exclude(id__in=answered_questions).filter(
        release_date=today,
        end_date__gte=today,
        course_id__in=enrolled_courses  # Filter questions for the courses the student is enrolled in
    )

    if request.method == 'POST':
        for question in questions:
            response = request.POST.get(f'question_{question.id}')
            if response:  # Ensure there's a response before saving
                Feedback.objects.create(
                    user=user_id,
                    question=question,
                    response=response
                )
        return redirect('student_dashboard')  # Redirect to a thank you page or another view

    return render(request, 'feedback_form.html', {'questions': questions})


def thank_you_view(request):
    return render(request, 'thank_you.html')

from django.shortcuts import render, redirect
from .models import Feedback

def view_feedback_responses(request):
    course_id = request.GET.get('course')  # Get the selected course from the request
    feedback_responses = Feedback.objects.all()

    if course_id:
        feedback_responses = feedback_responses.filter(question__course_id=course_id)  # Filter by course

    # Prepare data for chart
    response_counts = {
        'strongly_agree': 0,
        'agree': 0,
        'neutral': 0,
        'disagree': 0,
        'strongly_disagree': 0,
    }

    for feedback in feedback_responses:
        response_counts[feedback.response] += 1

    return render(request, 'view_feedback_responses.html', {
        'feedback_responses': feedback_responses,
        'response_counts': response_counts,
        'courses': Course.objects.all(),  # Pass all courses for filtering
    })

from django.shortcuts import render, redirect
from django.utils import timezone  # Ensure you have this imported
from .models import CalendarEvent, Course, TeacherCourse  # Import TeacherCourse for filtering

def add_event(request):
    if request.method == 'POST':
        title = request.POST.get('title')
        start_time = request.POST.get('start_time')
        end_time = request.POST.get('end_time')
        description = request.POST.get('description', '')
        event_type = request.POST.get('event_type')
        course_id = request.POST.get('course_id')

        # Debugging: print the received values
        print(f'Title: {title}, Start Time: {start_time}, End Time: {end_time}, Course ID: {course_id}')
        
        # Check if course_id is provided
        if not course_id:
            courses = Course.objects.filter(id__in=TeacherCourse.objects.filter(teacher_id=request.session.get('teacher_id')).values_list('course_id', flat=True))
            return render(request, 'add_event.html', {'error': 'Course must be selected.', 'courses': courses})

        # Convert start_time and end_time
        start_time = timezone.make_aware(timezone.datetime.fromisoformat(start_time))
        end_time = timezone.make_aware(timezone.datetime.fromisoformat(end_time))

        # Create the CalendarEvent object
        event = CalendarEvent(
            title=title,
            start_time=start_time,
            end_time=end_time,
            description=description,
            event_type=event_type,
            course_id=course_id,
            created_by_id=request.session.get('teacher_id')  # Set created_by to the teacher
        )
        
        # Attempt to save the event
        try:
            event.save()
            print('Event saved successfully!')
        except Exception as e:
            print(f'Error saving event: {e}')
        
        return redirect('view_events')  # Redirect to the event list after saving

    # Fetch courses assigned to the teacher for the dropdown
    courses = Course.objects.filter(id__in=TeacherCourse.objects.filter(teacher_id=request.session.get('teacher_id')).values_list('course_id', flat=True))
    
    return render(request, 'add_event.html', {'courses': courses})  # Render the template with courses


from django.shortcuts import render, redirect
from .models import Course, CalendarEvent, TeacherCourse  # Import your models

def view_events(request):
    teacher_id = request.session.get('teacher_id')
    
    if not teacher_id:
        # Handle case where teacher_id is not found (e.g., redirect to login)
        return redirect('login')

    # Fetch courses assigned to the teacher
    assigned_courses = Course.objects.filter(id__in=TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True))
    
    # Optionally, get the event_type from the request (for example, from a GET parameter)
    event_type = request.GET.get('event_type', None)

    # Fetch events created by the teacher for the assigned courses, with optional event_type filter
    events = CalendarEvent.objects.filter(
        created_by_id=teacher_id, 
        course__in=assigned_courses
    )
    
    if event_type:
        events = events.filter(event_type=event_type)

    return render(request, 'view_events.html', {
        'events': events,
        'courses': assigned_courses,  # Pass only the courses assigned to the teacher to the template
        'event_type': event_type  # Optionally pass the selected event_type to the template for UI purposes
    })


from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from .models import CalendarEvent, Enrollment
from django.contrib.auth import get_user_model
from django.utils import timezone
from datetime import datetime

def student_events(request):
    CustomUser = get_user_model()
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, 'You must be logged in as a student to view events.')
        return redirect('login')

    # Fetch the CustomUser (student) object using the session ID
    student = get_object_or_404(CustomUser, id=custom_user_id)

    # Fetch the enrollments for the student
    enrollments = Enrollment.objects.filter(student=student)

    if not enrollments.exists():
        messages.error(request, 'You are not registered for any courses.')
        return redirect('student_dashboard')

    # Create a list of enrolled courses with their enrollment dates and times
    enrolled_courses = [(enrollment.course, enrollment.enrollment_date, enrollment.enrollment_time) 
                       for enrollment in enrollments]

    # Get current date and time
    current_datetime = timezone.localtime()

    # Fetch events for the enrolled courses
    events = CalendarEvent.objects.filter(
        course__in=[course for course, _, _ in enrolled_courses]
    ).order_by('-start_time')  # Sort by start time, newest first

    # Further filter events based on the enrollment date and time
    filtered_events = []
    for course, enrollment_date, enrollment_time in enrolled_courses:
        enrollment_datetime = timezone.make_aware(
            datetime.combine(enrollment_date, enrollment_time or datetime.min.time())
        )
        # Fetch events that start after the enrollment date and time
        course_events = events.filter(course=course, start_time__gte=enrollment_datetime)
        
        # Add event status based on start and end times
        for event in course_events:
            event.status = get_event_status(event, current_datetime)
            filtered_events.append(event)

    # Filter by event type if provided in the GET request
    event_type = request.GET.get('event_type', '')
    if event_type:
        filtered_events = [event for event in filtered_events if event.event_type == event_type]

    # Sort filtered events by start time (newest first)
    filtered_events.sort(key=lambda x: x.start_time, reverse=True)

    context = {
        'events': filtered_events,
        'courses': [course for course, _, _ in enrolled_courses],
        'now': current_datetime,  # Add current datetime to context
    }

    return render(request, 'student_event.html', context)

def get_event_status(event, current_datetime):
    """
    Determine the status of an event based on its start and end times
    """
    current_date = current_datetime.date()
    start_date = event.start_time.date()
    end_date = event.end_time.date()
    
    if current_date < start_date:
        return 'upcoming'
    elif current_date > end_date:
        return 'ended'
    else:
        # Check if event is currently ongoing
        if (current_datetime >= event.start_time and 
            current_datetime <= event.end_time):
            return 'ongoing'
        elif current_datetime > event.end_time:
            return 'ended'
        else:
            return 'upcoming'


from django.http import JsonResponse
from .models import CalendarEvent

def filtered_events(request):
    event_type = request.GET.get('event_type', '')
    course_id = request.GET.get('course_id', '')
    events = CalendarEvent.objects.all()

    if event_type:
        events = events.filter(event_type=event_type)
    if course_id:
        events = events.filter(course_id=course_id)    

    event_list = [{
        'title': event.title,
        'start_time': event.start_time.isoformat(),
        'end_time': event.end_time.isoformat(),
        'description': event.description,
        'event_type': event.event_type,
        'color': get_event_color(event.event_type),
    } for event in events]

    return JsonResponse(event_list, safe=False)

def get_event_color(event_type):
    color_map = {
        'class': '#ffcc00',  # Yellow for classes
        'assignment': '#00ccff',      # Blue for assignments
        'exam': '#ff6666',            # Red for exams
        'holiday': '#28a745',         # Green for holidays
        'personal_event': '#6c757d'   # Gray for personal events
    }
    return color_map.get(event_type, '#007bff')  # Default to Bootstrap primary color


from django.shortcuts import render, redirect, get_object_or_404
from .models import Quizs, Course, Question, TeacherCourse  # Import TeacherCourse for filtering

def view_quiz_questions(request):
    teacher_id = request.session.get('teacher_id')  # Retrieve the teacher_id from the session

    if not teacher_id:
        # If the teacher is not logged in or the session has expired, redirect to the login page
        return redirect('login')

    # Fetch the courses taught by the teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)  # Only get the courses assigned to the teacher
    
    course_id = request.GET.get('course_id')  # Get the course ID from the query parameters

    # Fetch quizzes created by this teacher for their assigned courses
    quizzes = Quizs.objects.filter(teacher_id=teacher_id, course_id__in=assigned_courses)

    # Fetch questions for quizzes of the selected course
    questions = Question.objects.none()  # Start with an empty QuerySet

    if course_id:  # If a course is selected, filter questions by that course
        # Filter quizzes based on the course and fetch related questions
        quizzes = quizzes.filter(course_id=course_id)
        questions = Question.objects.filter(quiz__in=quizzes) .order_by('-quiz__start_date') # Fetch questions for the selected course's quizzes

    context = {
        'courses': courses,
        'questions': questions,
        'selected_course': course_id,
    }
    return render(request, 'view_quiz_questions.html', context)

from django.shortcuts import render, redirect
from .models import Question, UserAnswers, Quizs

def view_student_answers(request):
    teacher_id = request.session.get('teacher_id')  # Retrieve the teacher_id from the session

    if not teacher_id:
        return redirect('login')

    # Fetch all quizzes related to the teacher
    quizzes = Quizs.objects.filter(teacher_id=teacher_id)  # Adjust this if needed

    # Fetch all questions related to those quizzes
    questions = Question.objects.filter(quiz__in=quizzes)

    # Fetch all student answers for those questions
    student_answers = UserAnswers.objects.filter(question__in=questions).select_related('user')

    context = {
        'quizzes': quizzes,
        'student_answers': student_answers,
    }

    return render(request, 'view_student_answers.html', context)
from django.shortcuts import render, get_object_or_404  # type: ignore
from .models import Material, Course, TeacherCourse  # Import the TeacherCourse model

def view_uploaded_materials(request):
    # Assuming you have a session variable for the logged-in teacher
    teacher_id = request.session.get('teacher_id')

    # Fetch courses taught by the teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)  # Get assigned course IDs
    courses = Course.objects.filter(id__in=assigned_courses)  # Get only assigned courses

    course_id = request.GET.get('course', None)  # Get the selected course from the dropdown filter
    materials = Material.objects.filter(teacher_id=teacher_id)  # Get materials uploaded by the teacher

    # Filter materials by the selected course if course_id is provided
    if course_id:
        materials = materials.filter(course_id=course_id)

    context = {
        'materials': materials,
        'courses': courses,
        'selected_course': course_id,
    }
    return render(request, 'view_uploaded_materials.html', context)

from django.shortcuts import redirect
from django.utils import timezone
from .models import Attendance, ClassSchedule, CustomUser

def mark_attendance(request, schedule_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You need to log in to join the class.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
        class_schedule = ClassSchedule.objects.get(id=schedule_id)
        current_time = timezone.localtime(timezone.now())

        # Check if student has registered face data
        try:
            face_data = StudentFaceData.objects.get(user=student)
            if not face_data.is_face_captured:
                messages.error(request, "Please register your face first.")
                return redirect('face_capture')
        except StudentFaceData.DoesNotExist:
            messages.error(request, "Please register your face first.")
            return redirect('face_capture')

        # Check if within class time
        if class_schedule.start_time <= current_time.time() <= class_schedule.end_time and current_time.date() == class_schedule.date:
            # Show face verification page
            return render(request, 'face_verification.html', {
                'schedule_id': schedule_id,
                'meeting_link': class_schedule.meeting_link
            })
        else:
            # If outside class time, mark as absent
            Attendance.objects.create(
                student=student,
                class_schedule=class_schedule,
                check_in_time=current_time,
                status='absent'
            )
            messages.error(request, "Class is not active at this time.")
            return redirect('view_scheduled_classes')

    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
    except ClassSchedule.DoesNotExist:
        messages.error(request, "Class schedule not found.")
        return redirect('student_dashboard')
    except Exception as e:
        messages.error(request, str(e))
        return redirect('view_scheduled_classes')


from myApp.models import TeacherCourse, Enrollment, Course

def student_list(request):
    teacher_id = request.session.get('teacher_id')  # Assuming you're using session to manage teacher login

    # Get the courses assigned to this teacher
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)

    # Get students enrolled in those courses
    enrolled_students = Enrollment.objects.filter(course_id__in=teacher_courses).select_related('student')

    context = {
        'enrolled_students': enrolled_students,
    }
    return render(request, 'student_list.html', context)

from django.shortcuts import render, get_object_or_404  # type: ignore
from myApp.models import TeacherCourse, Enrollment, Attendance, ClassSchedule

def view_attendance(request):
    teacher_id = request.session.get('teacher_id')
    # Get the courses assigned to this teacher
    teacher_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)

    # Get the class schedules for the courses
    class_schedules = ClassSchedule.objects.filter(course_name_id__in=teacher_courses)

    # Get the selected class schedule ID from the form
    selected_class_schedule_id = request.GET.get('class_schedule_id')

    # Initialize variables for attendance records and course name
    attendance_records = None
    course_name = None

    if selected_class_schedule_id:
        # Fetch attendance records based on the selected class schedule
        attendance_records = Attendance.objects.filter(class_schedule_id=selected_class_schedule_id)

        # Fetch the course name associated with the selected class schedule
        selected_class_schedule = ClassSchedule.objects.get(id=selected_class_schedule_id)
        course_name = selected_class_schedule.course_name  # Assuming 'course_name' is the field that holds the course

    context = {
        'class_schedules': class_schedules,
        'attendance_records': attendance_records,
        'selected_class_schedule': selected_class_schedule_id,
        'course_name': course_name,
    }

    return render(request, 'view_attendance.html', context)

import requests
import json
from datetime import datetime
from django.shortcuts import render

def create_zoom_meeting(request):
    # Replace this line with the function that retrieves your access token dynamically
    access_token = get_zoom_access_token()

    url = "https://api.zoom.us/v2/users/me/meetings"
    headers = {
        "Authorization": f"Bearer {access_token}",
        "Content-Type": "application/json"
    }

    payload = {
        "topic": "class_name",
        "type": 2,
        "start_time": datetime.utcnow().strftime('%Y-%m-%dT%H:%M:%SZ'),
        "duration": 60,
        "timezone": "UTC",
        "settings": {
            "host_video": True,
            "participant_video": True,
            "join_before_host": True,
            "mute_upon_entry": False
        }
    }

    response = requests.post(url, headers=headers, data=json.dumps(payload))

    if response.status_code == 201:
        meeting_info = response.json()
        return render(request, 'zoom_meeting_created.html', {'meeting_info': meeting_info})
    else:
        return render(request, 'error.html', {'error': response.text})

from requests.auth import HTTPBasicAuth

def get_zoom_access_token():
    client_id = '95w1T2nlSpOWyuxu5Gjh4w'
    client_secret = 'IXjEp1mgI5w5HEb3VHR3NS8taAmyCX5q'
    account_id = 'bs4hsQ6GRte8O-EHmG8uDQ'
    
    url = f'https://zoom.us/oauth/token?grant_type=account_credentials&account_id={account_id}'
    response = requests.post(url, auth=HTTPBasicAuth(client_id, client_secret))
    
    if response.status_code == 200:
        access_token = response.json().get('access_token')
        return access_token
    else:
        raise Exception(f"Failed to get Access Token. Status code: {response.status_code}, Response: {response.text}")

from django.shortcuts import render, redirect
from django.contrib import messages
from .models import LeaveRequest, CustomUser
from datetime import datetime

def apply_leave(request):
    custom_user_id = request.session.get('custom_user_id')
    
    if not custom_user_id:
        messages.error(request, 'You need to log in to apply for leave.')
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)

        if request.method == 'POST':
            leave_type = request.POST.get('leave_type')
            reason = request.POST.get('reason')
            start_date = request.POST.get('start_date')
            end_date = request.POST.get('end_date')

            # Validate dates
            start_date_obj = datetime.strptime(start_date, '%Y-%m-%d').date()
            end_date_obj = datetime.strptime(end_date, '%Y-%m-%d').date()
            today = datetime.now().date()

            if start_date_obj < today:
                messages.error(request, 'Start date cannot be in the past.')
                return render(request, 'apply_leave.html')

            if end_date_obj < start_date_obj:
                messages.error(request, 'End date must be after start date.')
                return render(request, 'apply_leave.html')
            
            # Create leave request
            leave_request = LeaveRequest.objects.create(
                student=student,
                leave_type=leave_type,
                reason=reason,
                start_date=start_date,
                end_date=end_date,
                status='pending'  # Default status
            )
            
            messages.success(request, 'Leave request submitted successfully!')
            return redirect('student_leave_requests')  # Changed redirect to student_leave_requests
        
        return render(request, 'apply_leave.html')

    except CustomUser.DoesNotExist:
        messages.error(request, 'Student profile not found.')
        return redirect('login')
    except ValueError as e:
        messages.error(request, 'Invalid date format.')
        return render(request, 'apply_leave.html')
    except Exception as e:
        messages.error(request, f'An error occurred: {str(e)}')
        return render(request, 'apply_leave.html')


from django.shortcuts import render, get_object_or_404
from .models import LeaveRequest

def manage_leave_requests(request):
    leave_requests = LeaveRequest.objects.filter(status='pending')
    return render(request, 'manage_leave.html', {'leave_requests': leave_requests})

def update_leave_status(request, leave_id, status):
    leave_request = get_object_or_404(LeaveRequest, id=leave_id)
    leave_request.status = status
    leave_request.save()
    return redirect('manage_leave')

def student_leave_requests(request):
    custom_user_id = request.session.get('custom_user_id')
    student = CustomUser.objects.get(id=custom_user_id)
    leave_requests = LeaveRequest.objects.filter(student=student)

    return render(request, 'student_leave_requests.html', {
        'leave_requests': leave_requests
    })

from .models import Group, Message, TeacherMessage, Course, CustomUser, Teacher
from django.shortcuts import render, get_object_or_404, redirect

def group_chat_view(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    
    # Fetch the user based on the session variable custom_user_id
    custom_user_id = request.session.get('custom_user_id') 
    if not custom_user_id:
        return redirect('login')  # Redirect to login if not authenticated via session

    user = get_object_or_404(CustomUser, id=custom_user_id)

    # Get or create the group for the course
    group, created = Group.objects.get_or_create(course=course)

    # Check if the user is a student or a teacher, and add them to the respective group if not already added
    if user:  
        if not group.students.filter(id=user.id).exists():
            group.students.add(user)

    # Handle new message submission
    if request.method == 'POST':
        content = request.POST.get('message')
        if content:
            Message.objects.create(group=group, sender=user, content=content)
            return redirect('group_chat', course_id=course_id)

    # Get all messages related to the group
    student_messages = group.messages.all().order_by('timestamp')
    teacher_messages = group.teacher_messages.all().order_by('timestamp')

    # Combine messages from both models
    all_messages = sorted(list(student_messages) + list(teacher_messages), key=lambda m: m.timestamp)

    return render(request, 'group_chat.html', {
        'group': group,
        'messages': all_messages,
        'user': user,  # Pass the user to the template if needed
    })



def discussion_forum(request):
    custom_user_id = request.session.get('custom_user_id')
    enrolled_courses = Course.objects.filter(enrollments__student_id=custom_user_id)
    

    return render(request,'discussion_forum.html',{'enrolled_courses':enrolled_courses})


from .models import Group, TeacherMessage, Course
from django.shortcuts import render, get_object_or_404, redirect

def teacher_group_chat_view(request, course_id):
    teacher_id = request.session.get('teacher_id')

    if not teacher_id:
        return redirect('login')

    course = get_object_or_404(Course, id=course_id)
    group = get_object_or_404(Group, course=course)

    # Get all messages and combine them
    student_messages = group.messages.all()
    teacher_messages = group.teacher_messages.all()

    # Combine and sort all messages
    all_messages = list(student_messages) + list(teacher_messages)
    all_messages.sort(key=lambda x: x.timestamp)  # Sort by timestamp

    teacher = get_object_or_404(Teacher, id=teacher_id)

    if request.method == 'POST':
        content = request.POST.get('message')
        if content:
            # Create a message for teachers
            TeacherMessage.objects.create(group=group, teacher=teacher, content=content)
            return redirect('teacher_group_chat', course_id=course_id)

    return render(request, 'teacher_group_chat.html', {
        'group': group,
        'messages': all_messages,  # Pass the combined messages to the template
        'teacher': teacher,
    })


from django.shortcuts import render, get_object_or_404
from .models import TeacherCourse

def teacher_discussion_forum(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')  # Redirect to login if not authenticated

    # Get all courses assigned to the teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id)

    return render(request, 'teacher_discussion_forum.html', {'assigned_courses': assigned_courses})


# myproject/myApp/views.py
from django.shortcuts import render, get_object_or_404, redirect
from .models import Course

def edit_course(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    
    if request.method == 'POST':
        # Manually update the course fields
        course.course_name = request.POST.get('course_name')
        course.description = request.POST.get('description')
        course.duration = request.POST.get('duration')
        course.price = request.POST.get('price')
        course.starting_date = request.POST.get('starting_date')
        course.ending_date = request.POST.get('ending_date')
        
        # Handle file upload for the image
        if request.FILES.get('image'):
            course.image = request.FILES['image']
        
        # Save the updated course
        course.save()
        return redirect('course_list')  # Redirect to the course list or another page
    
    return render(request, 'edit_course.html', {'course': course})


from django.shortcuts import render, get_object_or_404, redirect
from .models import Course, CustomUser, Enrollment
from django.contrib import messages
import razorpay
from django.conf import settings

def enrollment_details(request, course_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, 'You need to log in to view enrollment details.')
        return redirect('login')  # Redirect to login if session doesn't have a custom_user_id

    user = get_object_or_404(CustomUser, id=custom_user_id)
    course = get_object_or_404(Course, id=course_id)
    
    client = razorpay.Client(auth=(settings.RAZORPAY_API_KEY, settings.RAZORPAY_API_SECRET))

    # Convert course price to paise (multiply by 100) as Razorpay expects the amount in paise
    amount = int(course.price * 100)

    # Create an order with Razorpay for tracking (optional)
    razorpay_order = client.order.create({
        "amount": amount,
        "currency": "INR",
        "payment_capture": "1"  # Auto-capture
    })

    # Pass the necessary data to the template
    context = {
        'user': user,
        'course': course,
        'enrolled': False,  # Default: user not enrolled
        'razorpay_key_id': settings.RAZORPAY_API_KEY,
        'razorpay_order_id': razorpay_order['id'],  # Order ID for tracking
        'amount': amount
    }

    # After the user confirms the payment, Razorpay will send a response with payment details
    if request.method == "POST":
        payment_id = request.POST.get('razorpay_payment_id')
        order_id = request.POST.get('razorpay_order_id')
        signature = request.POST.get('razorpay_signature')

        # Verify the payment signature
        params_dict = {
            'razorpay_order_id': order_id,
            'razorpay_payment_id': payment_id,
            'razorpay_signature': signature
        }
        
        try:
            # Verify the signature
            client.utility.verify_payment_signature(params_dict)

            # Create enrollment record and save payment details
            enrollment = Enrollment.objects.create(
                student=user,
                course=course,
                payment_amount=amount / 100,  # Converting paise back to INR
                payment_id=payment_id,
                payment_status="Completed",  # Or any status based on your business logic
            )

            # Update context and mark the user as enrolled
            context.update({'enrolled': True},{enrollment,})
            messages.success(request, f"Enrollment confirmed! Thank you, {user.first_name}. You have been successfully enrolled in {course.course_name}.")

        except Exception as e:
            messages.error(request, "Payment verification failed. Please try again.")
            return redirect('enrollment_details', course_id=course.id)

    return render(request, 'enrollment_details.html', context)


def confirm_enrollment(request, course_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, 'You need to log in to confirm enrollment.')
        return redirect('login')  # Redirect to login if session doesn't have a custom_user_id

    user = get_object_or_404(CustomUser, id=custom_user_id)
    course = get_object_or_404(Course, id=course_id)

    # Create an enrollment record
    enrollment = Enrollment.objects.create(
        student=user,
        course=course,
        payment_amount=course.price,  # Assuming price is already in INR
        
    )

    messages.success(request, f"You have successfully confirmed your enrollment in {course.course_name}.")
    
    # Redirect to enrollment details to display Razorpay button
    return redirect('available_courses')


from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from .models import CalendarEvent, EventRegistration
from django.conf import settings

def register_event(request, event_id):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')  # Redirect to login page if user ID is not in session

    event = get_object_or_404(CalendarEvent, id=event_id)

    if request.method == 'POST':
        contact_number = request.POST.get('contact_number')

        # Check if the user is already registered for the event
        if EventRegistration.objects.filter(event=event, user_id=custom_user_id).exists():
            messages.warning(request, "You are already registered for this event.")
        else:
            # Save the registration
            EventRegistration.objects.create(
                event=event,
                user_id=custom_user_id,
                contact_number=contact_number,
                status='registered'
            )
            messages.success(request, "You have successfully registered for the event.")
        return redirect('student_event')  # Replace 'events' with your event listing URL name

    return render(request, 'event_registration.html', {'event': event})

# views.py
from django.shortcuts import render, get_object_or_404
from .models import TeacherInterview

def view_interview_details(request, teacher_id):
    interview = get_object_or_404(TeacherInterview, teacher_id=teacher_id)
    return render(request, 'interview_details.html', {'interview': interview})

from django.http import JsonResponse
from .models import Course

def check_course_name(request):
    course_name = request.GET.get('course_name', None)
    if course_name:
        course_exists = Course.objects.filter(course_name=course_name).exists()
        return JsonResponse({'exists': course_exists})
    return JsonResponse({'exists': False})


def assign_students_to_teacher(request):
    teachers = Teacher.objects.all()  # Fetch all teachers
    return render(request, 'assign_students_to_teacher.html', {'teachers': teachers})
def select_course_for_teacher(request, teacher_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    teacher_courses = TeacherCourse.objects.filter(teacher=teacher)
    return render(request, 'select_course_for_teacher.html', {'teacher': teacher, 'teacher_courses': teacher_courses})
def assign_students(request, teacher_id, course_id):
    teacher = get_object_or_404(Teacher, id=teacher_id)
    course = get_object_or_404(Course, id=course_id)
    
    # Use the Enrollment model or appropriate model linking students and courses
    students = CustomUser.objects.filter(enrollment__course=course).distinct()

    if request.method == 'POST':
        selected_student_ids = request.POST.getlist('students')
        for student_id in selected_student_ids:
            student = CustomUser.objects.get(id=student_id)
            TeacherStudent.objects.create(teacher=teacher, student=student)
        return redirect('assign_students_to_teacher')

    return render(request, 'assign_students.html', {'teacher': teacher, 'course': course, 'students': students})




from django.http import JsonResponse
from .models import Course  # Import the Course model

def check_course_name(request):
    if request.method == 'GET':
        course_name = request.GET.get('course_name', '').strip()  # Get the course name from the request
        if course_name:  # Ensure course_name is not empty
            if Course.objects.filter(course_name__iexact=course_name).exists():  # Check for case-insensitive match
                return JsonResponse({'status': 'error', 'message': 'Course name already exists'})
            else:
                return JsonResponse({'status': 'success', 'message': 'Course name is valid'})
        else:
            return JsonResponse({'status': 'error', 'message': 'Course name is required'})
    return JsonResponse({'status': 'error', 'message': 'Invalid request method'})


from django.shortcuts import render, redirect
from .models import UserAnswers, TeacherCourse, Course, Quizs, Question

def evaluate_answers(request):
    if 'teacher_id' in request.session:
        teacher_id = request.session['teacher_id']
        selected_course_id = request.GET.get('course_id')
        selected_quiz_id = request.GET.get('quiz_id')

        # Fetch all courses assigned to the teacher
        assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
        courses = Course.objects.filter(id__in=assigned_courses)

        # Get quizzes for the selected course
        quizzes = Quizs.objects.filter(course_id__in=assigned_courses)

        if selected_course_id:
            quizzes = quizzes.filter(course_id=selected_course_id)  # Filter quizzes by selected course

        if selected_quiz_id:
            quizzes = quizzes.filter(id=selected_quiz_id)  # Filter quizzes by selected quiz

        # Get the UserAnswers for the selected quiz and course
        answers = UserAnswers.objects.filter(
            question__quiz__course_id__in=assigned_courses
        ).select_related('question')
        # Get the UserAnswers for the selected quiz and course that haven't been graded yet
        answers = UserAnswers.objects.filter(
            question__quiz__course_id__in=assigned_courses,
            marks_obtained__isnull=True  # Only include answers without assigned marks
        ).select_related('question', 'user')

        if selected_course_id:
            answers = answers.filter(question__quiz__course_id=selected_course_id)

        if selected_quiz_id:
            answers = answers.filter(question__quiz_id=selected_quiz_id)

        # Prepare a list to store the evaluation results
        evaluation_results = []
        correct_count = 0  # Counter to track correct answers

        for answer in answers:
            # Get the correct option from the Question model
            correct_option = answer.question.correct_option
            
            # Check if the selected option is correct
            is_correct = answer.selected_option == correct_option
            
            # Assign grade based on correctness
            if is_correct:
                grade = 'A+'
                correct_count += 1
            else:
                grade = 'F'

            # Calculate percentage for the specific answer (1 if correct, 0 if incorrect)
            percentage = 100 if is_correct else 0

            # Store the percentage in the UserAnswer model
            answer.percentage = percentage
            answer.marks_obtained = 1 if is_correct else 0  # Assign marks based on correctness
            answer.grade = grade
            answer.save()

            # Add the result to the list for display
            evaluation_results.append({
                'answer_id': answer.id,
                'student': answer.user.username,
                'question': answer.question.text,
                'selected_option': answer.selected_option,
                'correct_option': correct_option,
                'is_correct': is_correct,
                'marks': answer.marks_obtained,
                'grade': grade,
                'attempt_date': answer.attempt_date,
                'percentage': answer.percentage,  # Include the percentage in the results
            })

        # Calculate the total number of questions
        total_questions = answers.count()

        # Calculate the overall percentage
        percentage = (correct_count / total_questions) * 100 if total_questions else 0

        # Determine final grade based on the percentage (optional)
        if percentage >= 90:
            final_grade = 'A+'
        elif percentage >= 80:
            final_grade = 'A'
        elif percentage >= 70:
            final_grade = 'B'
        elif percentage >= 60:
            final_grade = 'C'
        else:
            final_grade = 'F'

        # Process marks update for each student answer (if POST request)
        if request.method == 'POST':
            # Iterate through the answers and save the marks
            for answer_id, marks in request.POST.items():
                if answer_id.startswith('marks_'):  # Marks input fields have names like "marks_<answer_id>"
                    user_answer_id = answer_id.split('_')[1]
                    try:
                        user_answer = UserAnswers.objects.get(id=user_answer_id)
                        user_answer.marks_obtained = marks
                        user_answer.grade = 'A+' if int(marks) > 0 else 'F'  # Set grade based on marks
                        user_answer.save()
                    except UserAnswers.DoesNotExist:
                        pass  # Handle invalid IDs gracefully

            # Redirect to the same page to reflect updated marks
            return redirect(request.path + f"?course_id={selected_course_id}&quiz_id={selected_quiz_id}" if selected_course_id and selected_quiz_id else request.path)

        context = {
            'evaluation_results': evaluation_results,
            'courses': courses,
            'selected_course_id': selected_course_id,
            'selected_quiz_id': selected_quiz_id,
            'quizzes': quizzes,
            'total_questions': total_questions,
            'correct_count': correct_count,
            'percentage': percentage,
            'final_grade': final_grade,
        }

        # Pass the evaluation results to the template
        return render(request, 'evaluate_answers.html', context)
    else:
        return redirect('login')



from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import logging
import google.generativeai as genai

# Configure logging
logger = logging.getLogger(__name__)

# Set your Gemini API Key
API_KEY = os.getenv('API_KEY')

# Configure the Generative AI client
genai.configure(api_key=API_KEY)

@csrf_exempt
def chatbot_response(request):
    if request.method == 'POST':
        user_message = request.POST.get('message')

        if not user_message:
            return JsonResponse({'error': 'No message provided'}, status=400)

        try:
            # Create the model
            model = genai.GenerativeModel("gemini-1.5-flash")

            # Generate content using the user message
            response = model.generate_content(user_message)

            # Extract the AI response
            ai_message = response.text

            # Log the AI response for debugging
            logger.info(f"Gemini API response: {ai_message}")

            return JsonResponse({'response': ai_message})

        except Exception as e:
            logger.error(f"Error fetching response from Gemini API: {str(e)}")
            return JsonResponse({'error': 'Failed to get response from AI'}, status=500)

    return JsonResponse({'error': 'Invalid request method'}, status=400)


import os  # Ensure this import is present
import time  # Ensure this import is present
import base64  # Ensure this import is present
import requests  # Ensure this import is present
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import logging

# Replace with your actual Hugging Face API key
HUGGING_FACE_API_KEY = os.getenv('HUGGING_FACE_API_KEY')
# Configure logging
logger = logging.getLogger(__name__)

@csrf_exempt
def generate_image(request):
    if request.method == 'POST':
        user_description = request.POST.get('description')

        if not user_description:
            return JsonResponse({'error': 'No description provided'}, status=400)

        try:
            logger.info(f"Generating image with description: {user_description}")

            # API request headers
            headers = {
                'Authorization': f'Bearer {HUGGING_FACE_API_KEY}',
                'Content-Type': 'application/json'
            }

            # Payload with user input and optional parameters
            payload = {
                'text_prompts': [{'text': user_description}],
                'cfg_scale': 7.5,  # Prompt adherence (lower = more creative)
                'width': 1024,  # Allowed width
                'height': 1024,  # Allowed height
                'samples': 1  # Number of images to generate
            }

            # Make API request
            response = requests.post(
                'https://api.stability.ai/v1/generation/stable-diffusion-xl-1024-v1-0/text-to-image',
                headers=headers,
                json=payload
            )

            logger.info(f"Response status code: {response.status_code}")

            # Handle API response
            if response.status_code == 200:
                response_data = response.json()

                if 'artifacts' in response_data and response_data['artifacts']:
                    # Extract base64 image data
                    image_data = response_data['artifacts'][0].get('base64')
                    if image_data:
                        # Ensure the directory exists
                        os.makedirs('media/generated_images', exist_ok=True)  # Create directory if it doesn't exist
                        
                        # Save the image to a file
                        image_file_path = f'media/generated_images/image_{int(time.time())}.png'
                        with open(image_file_path, 'wb') as image_file:
                            image_file.write(base64.b64decode(image_data))
                        
                        # Return the URL of the saved image
                        return JsonResponse({'image_url': f'/{image_file_path}'})
                else:
                    logger.error("No artifacts found in response.")
                    return JsonResponse({'error': 'No images generated'}, status=500)
            else:
                logger.error(f"API returned an error: {response.text}")
                return JsonResponse({'error': 'Failed to generate image'}, status=response.status_code)

        except Exception as e:
            logger.error(f"Error generating image: {e}")
            return JsonResponse({'error': 'An error occurred'}, status=500)

    return JsonResponse({'error': 'Invalid request method'}, status=400)
from django.core.files.storage import FileSystemStorage

from datetime import datetime
from django.shortcuts import render, get_object_or_404, redirect
from django.http import HttpResponse
from xhtml2pdf import pisa
from .models import Course, CustomUser, Enrollment
from django.template.loader import render_to_string

def generate_certificate(request, course_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    course = get_object_or_404(Course, id=course_id)
    student = get_object_or_404(CustomUser, id=custom_user_id)

    # Convert course ending date and today's date to date objects to avoid comparison error
    today = datetime.today().date()  # Today's date (no time part)
    end_date = course.ending_date  # No need to call .date() on already a date object

    # Check if the course is completed (end date passed)
    if today < end_date:
        return HttpResponse("Course is not yet completed.", status=403)

    enrollment = Enrollment.objects.filter(student=student, course=course).first()
    if not enrollment:
        return HttpResponse("You are not enrolled in this course.", status=403)

    template_path = 'certificate_template.html'
    context = {'student': student, 'course': course}

    response = HttpResponse(content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{student.username}_certificate.pdf"'

    html = render_to_string(template_path, context)
    pisa_status = pisa.CreatePDF(html, dest=response)

    if pisa_status.err:
        return HttpResponse("Error generating PDF", status=500)
    return response



def student_performance_view(request):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    # Fetch quiz performance
    quiz_data = UserAnswers.objects.filter(user_id=student.id).order_by('attempt_date')
    quiz_performance = []
    for quiz_answer in quiz_data:
        course_name = (
            quiz_answer.quiz.course.course_name
            if hasattr(quiz_answer, 'quiz') and hasattr(quiz_answer.quiz, 'course')
            else 'Unknown'
        )
        quiz_performance.append({
            'course_name': course_name,
            'date': quiz_answer.attempt_date,
            'marks': quiz_answer.marks_obtained,
        })

    # Fetch assignment performance
    assignment_data = AssignmentSubmission.objects.filter(student_id=student.id).order_by('submitted_at')
    assignment_performance = []
    for assignment_submission in assignment_data:
        course_name = (
            assignment_submission.assignment.course.course_name
            if hasattr(assignment_submission, 'assignment') and hasattr(assignment_submission.assignment, 'course')
            else 'Unknown'
        )
        assignment_performance.append({
            'course_name': course_name,
            'date': assignment_submission.submitted_at,
            'grade': assignment_submission.grade,
        })

    # Calculate averages
    average_quiz_marks = quiz_data.aggregate(avg_marks=Avg('marks_obtained'))['avg_marks'] or 0
    average_assignment_grade = assignment_data.aggregate(avg_grade=Avg('grade'))['avg_grade'] or 0

    # Enrolled courses
    enrollments = Enrollment.objects.filter(student=student)
    enrolled_courses_with_dates = [
        (enrollment.course, enrollment.enrollment_date, enrollment.enrollment_time)
        for enrollment in enrollments
    ]

    context = {
        'student': student,
        'quiz_performance': quiz_performance,
        'assignment_performance': assignment_performance,
        'average_quiz_marks': average_quiz_marks,
        'average_assignment_grade': average_assignment_grade,
        'enrolled_courses_with_dates': enrolled_courses_with_dates,
    }

    return render(request, 'student_performance.html', context)



from django.utils.timezone import now
from django.shortcuts import render, redirect
from django.contrib import messages
from .models import UserAnswers, CustomUser, Enrollment, Quizs, Material, Course

def quiz_marks_view(request):
    # Retrieve the custom user ID from the session
    custom_user_id = request.session.get('custom_user_id')

    if not custom_user_id:
        messages.error(request, "You are not logged in.")
        return redirect('login')

    try:
        # Fetch the student (CustomUser) based on the custom user ID
        student = CustomUser.objects.get(id=custom_user_id)
    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')

    # Fetch enrolled courses for the student
    enrolled_courses = Enrollment.objects.filter(student=student).select_related('course')

    if not enrolled_courses.exists():
        messages.error(request, "You are not enrolled in any courses.")
        return redirect('student_dashboard')

    # Get course IDs for enrolled courses
    enrolled_course_ids = enrolled_courses.values_list('course_id', flat=True)

    # Fetch quizzes associated with enrolled courses
    quizzes = Quizs.objects.filter(course_id__in=enrolled_course_ids).select_related('course', 'teacher')

    # Fetch UserAnswers for the student related to these quizzes
    user_answers = UserAnswers.objects.filter(user=student, question__quiz__in=quizzes).select_related('question__quiz')

    # Organize data by quiz
    quiz_marks = {}
    passed_quizzes = []  # To track passed quizzes for recommendations

    for quiz in quizzes:
        quiz_questions = user_answers.filter(question__quiz=quiz)
        total_marks = sum(answer.marks_obtained or 0 for answer in quiz_questions)
        total_questions = quiz_questions.count()
        percentage = (total_marks / (total_questions * 1)) * 100 if total_questions > 0 else 0
        grade = calculate_grade(percentage)  # A helper function for grade calculation

        # Track quizzes the student has passed
        if grade != "F":
            passed_quizzes.append(quiz)

        # Fetch study materials if the student failed the quiz
        study_materials = []
        if grade == "F":
            study_materials = Material.objects.filter(course=quiz.course)

        quiz_marks[quiz] = {
            'questions': quiz_questions,
            'total_marks': total_marks,
            'total_questions': total_questions,
            'percentage': percentage,
            'grade': grade,
            'study_materials': study_materials,  # Add study materials to the context
        }

    # Recommend courses for passed quizzes
    recommended_courses = Course.objects.exclude(id__in=enrolled_course_ids).filter(
        starting_date__gte=now().date()
    )

    context = {
        'quiz_marks': quiz_marks,
        'recommended_courses': recommended_courses,  # Add recommended courses to the context
    }
    return render(request, 'quiz_marks.html', context)


def calculate_grade(percentage):
    """
    Helper function to calculate grade based on percentage.
    """
    if percentage >= 90:
        return "A+"
    elif percentage >= 80:
        return "A"
    elif percentage >= 70:
        return "B+"
    elif percentage >= 60:
        return "B"
    elif percentage >= 50:
        return "C"
    else:
        return "F"

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.http import JsonResponse, HttpResponse
from django.views.decorators.csrf import csrf_exempt
import json
import PyPDF2
from docx import Document
from pptx import Presentation
from django.conf import settings
import os
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from .models import GeneratedQuestionPaper
from django.core.files.base import ContentFile
import io
from django.utils import timezone

try:
    import google.generativeai as genai
    # Configure Gemini AI
    GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
    genai.configure(api_key=GEMINI_API_KEY)
    
    # Test available models
    for m in genai.list_models():
        if 'generateContent' in m.supported_generation_methods:
            print(m.name)
            
except ImportError:
    print("Please install google-generativeai package using: pip install google-generativeai")
    raise


def question_generator(request):
    teacher_id = request.session.get('teacher_id')  # Get teacher ID from session
    
    if not teacher_id:
        return redirect('login')
    
    # Get courses assigned to this teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    courses = Course.objects.filter(id__in=assigned_courses)

    context = {
        'courses': courses,
        'teacher_id': teacher_id,
        'first_name': request.session.get('first_name', ''),
        'last_name': request.session.get('last_name', '')
    }
    return render(request, 'question_generator.html', context)

def generate_pdf(questions_data, filename):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )

    # Container for the 'Flowable' objects
    elements = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    styles.add(ParagraphStyle(
        name='QuestionStyle',
        fontSize=12,
        spaceAfter=20,
        spaceBefore=20
    ))
    
    # Add title
    title = Paragraph(filename.split('_')[0], styles['Heading1'])
    elements.append(title)
    elements.append(Spacer(1, 12))

    # Process each section
    for section_index, section in enumerate(questions_data):
        # Add section header
        section_title = Paragraph(
            f"Section {section_index + 1} ({section['marks']} marks per question)", 
            styles['Heading2']
        )
        elements.append(section_title)
        elements.append(Spacer(1, 12))

        # Process each question in the section
        for q_index, question in enumerate(section['questions']):
            # Question text
            q_text = f"Q{q_index + 1}. {question['question']}"
            elements.append(Paragraph(q_text, styles['QuestionStyle']))

            # Mark distribution table
            mark_data = [[Paragraph("Component", styles['Heading4']), 
                         Paragraph("Marks", styles['Heading4'])]]
            
            for dist in question['marks_distribution']:
                component = dist.split('(')[0].strip()
                marks = dist.split('(')[1].replace(')', '').strip()
                mark_data.append([Paragraph(component, styles['Normal']), 
                                Paragraph(marks, styles['Normal'])])

            table = Table(mark_data, colWidths=[4*inch, 1*inch])
            table.setStyle(TableStyle([
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('TOPPADDING', (0, 0), (-1, 0), 12),
            ]))
            elements.append(table)
            elements.append(Spacer(1, 20))

    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

@csrf_exempt
def generate_question(request):
    if request.method == 'POST':
        try:
            # Get teacher ID and course ID
            teacher_id = request.session.get('teacher_id')
            course_id = request.POST.get('course_id')
            document = request.FILES.get('document')
            paper_title = request.POST.get('title', '')
            configs = json.loads(request.POST.get('configs', '[]'))
            
            if not teacher_id or not course_id:
                return JsonResponse({'error': 'Teacher ID or Course ID missing'})

            # Get the course object
            course = get_object_or_404(Course, id=course_id)
            
            # Extract text from document
            text_content = extract_text_from_document(document)
            all_questions = []
            total_marks = 0
            difficulty = 'medium'  # default difficulty
            
            for config in configs:
                question_count = int(config['questionCount'])
                marks = int(config['marksPerQuestion'])
                difficulty = config['difficulty']
                total_marks += marks * question_count
                
                if marks >= 10:
                    difficulty = "hard"
                elif marks >= 5:
                    difficulty = "medium" if difficulty == "easy" else difficulty
                
                questions = generate_questions_from_text(
                    text_content,
                    question_count,
                    marks,
                    difficulty
                )
                
                all_questions.append({
                    'marks': marks,
                    'questions': questions
                })

            # Generate PDFs
            timestamp = timezone.now().strftime('%Y%m%d_%H%M%S')
            filename = f'{paper_title}_{timestamp}'
            
            question_pdf = generate_pdf(all_questions, filename)
            answer_key_pdf = generate_answer_key_pdf(all_questions, filename)

            # Save to database
            question_paper = GeneratedQuestionPaper(
                title=paper_title,
                total_marks=total_marks,
                difficulty=difficulty,
                teacher_id=teacher_id,
                course=course
            )
            
            # Save question paper PDF
            question_paper.pdf_file.save(f'{filename}_questions.pdf', 
                                       ContentFile(question_pdf.getvalue()))
            
            # Save answer key PDF
            question_paper.answer_key_file.save(f'{filename}_answers.pdf', 
                                              ContentFile(answer_key_pdf.getvalue()))
            
            question_paper.save()

            return JsonResponse({
                'success': True,
                'questions': all_questions,
                'pdf_url': question_paper.pdf_file.url,
                'answer_key_url': question_paper.answer_key_file.url
            })
            
        except Exception as e:
            return JsonResponse({'error': str(e)}, status=500)
    
    return JsonResponse({'error': 'Invalid request'}, status=400)

def extract_text_from_document(document):
    """Extract text from various document formats"""
    file_extension = document.name.split('.')[-1].lower()
    text_content = ""
    
    try:
        if file_extension == 'pdf':
            pdf_reader = PyPDF2.PdfReader(document)
            for page in pdf_reader.pages:
                text_content += page.extract_text()
                
        elif file_extension in ['doc', 'docx']:
            doc = Document(document)
            for para in doc.paragraphs:
                text_content += para.text + "\n"
                
        elif file_extension in ['ppt', 'pptx']:
            prs = Presentation(document)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
    except Exception as e:
        raise Exception(f"Error processing document: {str(e)}")
    
    return text_content

def generate_questions_from_text(text, count, marks, difficulty):
    try:
        model = genai.GenerativeModel('gemini-1.5-pro')
        
        # Calculate mark distribution
        concept_marks = marks // 3
        explanation_marks = marks // 3
        example_marks = marks - (2 * (marks // 3))
        
        prompt = f"""
        Generate {count} questions based on the following text. Format as a JSON array.

        Each question object must have exactly this structure, no variations:
        {{
            "question": "Question text",
            "marks_distribution": [
                "concept ({concept_marks} marks)",
                "explanation ({explanation_marks} marks)",
                "example ({example_marks} marks)"
            ],
            "answer": {{
                "main_points": ["point 1", "point 2"],
                "keywords": ["keyword1", "keyword2"],
                "explanation": "Brief explanation",
                "examples": ["example1"]
            }}
        }}

        Text: {text[:3000]}  # Limiting text length to avoid token limits

        Rules:
        - Difficulty level: {difficulty}
        - Generate exactly {count} questions
        - Use only double quotes, not single quotes
        - Keep answers concise
        - Return only the JSON array
        """
        
        response = model.generate_content(prompt)
        response_text = response.text.strip()
        
        # Clean up the response
        if '```' in response_text:
            response_text = response_text.split('```')[1].split('```')[0]
        
        # Extract just the JSON array
        start_idx = response_text.find('[')
        end_idx = response_text.rfind(']') + 1
        if start_idx != -1 and end_idx != 0:
            response_text = response_text[start_idx:end_idx]
        
        # Clean up common JSON issues
        response_text = (response_text
            .replace('\n', ' ')
            .replace('  ', ' ')
            .replace('} {', '},{')
            .replace('}}]"', '}}]')
            .replace('"}"}', '}}')
            .replace('\\"', '"')
            .replace('""', '"')
            .strip())
        
        try:
            questions = json.loads(response_text)
            if not isinstance(questions, list):
                questions = [questions]
            
            # Basic validation
            for q in questions:
                if not all(key in q for key in ['question', 'marks_distribution', 'answer']):
                    raise Exception("Question missing required fields")
                if not all(key in q['answer'] for key in ['main_points', 'keywords', 'explanation', 'examples']):
                    raise Exception("Answer missing required fields")
            
            return questions
            
        except json.JSONDecodeError as e:
            print(f"JSON Error: {str(e)}")
            print(f"Response text: {response_text[:200]}")
            raise Exception("Failed to parse AI response as JSON")
            
    except Exception as e:
        raise Exception(f"Error generating questions: {str(e)}")

def view_question_papers(request):
    """View to list all generated question papers"""
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')

    # Get papers for courses assigned to this teacher
    assigned_courses = TeacherCourse.objects.filter(teacher_id=teacher_id).values_list('course_id', flat=True)
    papers = GeneratedQuestionPaper.objects.filter(
        course_id__in=assigned_courses
    ).order_by('-created_at')
    
    return render(request, 'view_question_papers.html', {'papers': papers})

def delete_question_paper(request, paper_id):
    """View to delete a question paper"""
    if request.method == 'POST':
        paper = get_object_or_404(GeneratedQuestionPaper, id=paper_id)
        # Delete the file from storage
        if paper.pdf_file:
            if os.path.isfile(paper.pdf_file.path):
                os.remove(paper.pdf_file.path)
        # Delete the database record
        paper.delete()
        return JsonResponse({'status': 'success'})
    return JsonResponse({'status': 'error'}, status=400)

def download_question_paper(request, paper_id):
    """View to download a question paper"""
    paper = get_object_or_404(GeneratedQuestionPaper, id=paper_id)
    if paper.pdf_file:
        response = HttpResponse(paper.pdf_file, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{paper.pdf_file.name}"'
        return response
    return HttpResponse('PDF not found', status=404)

#pip install google-generativeai
#pip install PyPDF2 python-docx python-pptx
#pip install reportlab

from django.shortcuts import get_object_or_404, redirect
from django.http import HttpResponse
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.colors import tan, black, white
from .models import Course, CustomUser, Enrollment
import os
import io
from django.conf import settings

def generate_certificate(request, course_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    course = get_object_or_404(Course, id=course_id)
    student = get_object_or_404(CustomUser, id=custom_user_id)

    # Check if the course is completed
    today = datetime.today().date()
    if today < course.ending_date:
        return HttpResponse("Course is not yet completed.", status=403)

    enrollment = Enrollment.objects.filter(student=student, course=course).first()
    if not enrollment:
        return HttpResponse("You are not enrolled in this course.", status=403)

    # Create a buffer to hold the PDF data
    buffer = io.BytesIO()

    # Create a canvas
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    # Draw a border
    c.setStrokeColor(tan)
    c.setLineWidth(10)
    c.rect(30, 30, width - 60, height - 60)

    # Draw a background rectangle
    c.setFillColor(white)
    c.rect(40, 40, width - 80, height - 80, fill=True)

    # Set up the certificate layout
    c.setFont("Helvetica-Bold", 30)
    c.setFillColor(tan)
    c.drawCentredString(width / 2, height - 100, "Certificate of Completion")

    c.setFont("Helvetica", 18)
    c.setFillColor(black)
    c.drawCentredString(width / 2, height - 150, "This certificate is presented to")
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width / 2, height - 200, f"{student.first_name} {student.last_name}")

    c.setFont("Helvetica", 18)
    c.drawCentredString(width / 2, height - 250, "For successfully completing the course")
    c.setFont("Helvetica-Bold", 22)
    c.drawCentredString(width / 2, height - 300, f"{course.course_name}")

    c.setFont("Helvetica", 16)
    c.drawCentredString(width / 2, height - 350, f"On: {course.ending_date}")

    # Include the image
    image_path = os.path.join(settings.BASE_DIR, 'myApp/static/images/academics.jpg')
    c.drawImage(image_path, 50, height - 150, width=80, height=80, preserveAspectRatio=True, mask='auto')

    # Finalize the PDF
    c.showPage()
    c.save()

    # Get the PDF data from the buffer
    buffer.seek(0)

    # Create a response
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{student.username}_certificate.pdf"'

    return response

from django.shortcuts import get_object_or_404, redirect
from django.http import HttpResponse
from reportlab.lib.pagesizes import landscape, letter
from reportlab.pdfgen import canvas
from reportlab.lib.colors import blue, black, white
from .models import Course, CustomUser, Enrollment
import os
import io
from django.conf import settings
from datetime import datetime

def generate_certificate(request, course_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return redirect('login')

    course = get_object_or_404(Course, id=course_id)
    student = get_object_or_404(CustomUser, id=custom_user_id)

    # Check if the course is completed
    today = datetime.today().date()
    if today < course.ending_date:
        return HttpResponse("Course is not yet completed.", status=403)

    enrollment = Enrollment.objects.filter(student=student, course=course).first()
    if not enrollment:
        return HttpResponse("You are not enrolled in this course.", status=403)

    # Create a buffer to hold the PDF data
    buffer = io.BytesIO()

    # Create a canvas in landscape mode
    c = canvas.Canvas(buffer, pagesize=landscape(letter))
    width, height = landscape(letter)

    # Draw a border
    c.setStrokeColor(blue)
    c.setLineWidth(10)
    c.rect(30, 30, width - 60, height - 60)

    # Draw a background rectangle
    c.setFillColor(white)
    c.rect(40, 40, width - 80, height - 80, fill=True)

    # Include the image at the top center
    image_path = os.path.join(settings.BASE_DIR, 'myApp/static/images/academics.jpg')
    c.drawImage(image_path, width / 2 - 40, height - 120, width=80, height=80, preserveAspectRatio=True, mask='auto')

    # Set up the certificate layout
    c.setFont("Helvetica-Bold", 30)
    c.setFillColor(blue)
    c.drawCentredString(width / 2, height - 150, "Certificate of Completion")

    c.setFont("Helvetica", 18)
    c.setFillColor(black)
    c.drawCentredString(width / 2, height - 200, "This certificate is presented to")
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(width / 2, height - 250, f"{student.first_name} {student.last_name}")

    c.setFont("Helvetica", 18)
    c.drawCentredString(width / 2, height - 300, "For successfully completing the course")
    c.setFont("Helvetica-Bold", 22)
    c.drawCentredString(width / 2, height - 350, f"{course.course_name}")

    c.setFont("Helvetica", 16)
    c.drawCentredString(width / 2, height - 400, f"On: {course.ending_date}")

    # Finalize the PDF
    c.showPage()
    c.save()

    # Get the PDF data from the buffer
    buffer.seek(0)

    # Create a response
    response = HttpResponse(buffer, content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{student.username}_certificate.pdf"'

    return response

import requests
from django.conf import settings
from django.shortcuts import render

def search_books(request):
    query = request.GET.get('q', '')
    books = []
    error_message = None

    if query:
        try:
            # Add fields parameter to get more information including viewability
            url = f"https://www.googleapis.com/books/v1/volumes?q={query}&maxResults=20&fields=items(volumeInfo/title,volumeInfo/authors,volumeInfo/publishedDate,volumeInfo/imageLinks,volumeInfo/previewLink,volumeInfo/infoLink,volumeInfo/canonicalVolumeLink,accessInfo/pdf,accessInfo/viewability,accessInfo/webReaderLink)"
            
            response = requests.get(url)
            response.raise_for_status()
            
            data = response.json()
            books = data.get("items", [])
            
            # Process image URLs and add reading options
            for book in books:
                if 'imageLinks' in book.get('volumeInfo', {}):
                    for key in ['thumbnail', 'smallThumbnail']:
                        if key in book['volumeInfo']['imageLinks']:
                            image_url = book['volumeInfo']['imageLinks'][key]
                            # Convert HTTP to HTTPS and remove edge=curl
                            if image_url.startswith('http:'):
                                image_url = image_url.replace('http:', 'https:')
                            image_url = image_url.replace('&edge=curl', '')
                            book['volumeInfo']['imageLinks'][key] = image_url + "&zoom=1"
                
                # Add reading options based on availability
                book['reading_options'] = {
                    'preview_available': bool(book.get('volumeInfo', {}).get('previewLink')),
                    'web_reader_available': bool(book.get('accessInfo', {}).get('webReaderLink')),
                    'pdf_available': book.get('accessInfo', {}).get('pdf', {}).get('isAvailable', False),
                    'viewability': book.get('accessInfo', {}).get('viewability', 'NO_PAGES')
                }
            
        except requests.RequestException as e:
            error_message = f"An error occurred while searching: {str(e)}"
            print(f"API Error: {str(e)}")
        except Exception as e:
            error_message = "An unexpected error occurred"
            print(f"Unexpected Error: {str(e)}")

    return render(request, "search_books.html", {
        "books": books,
        "query": query,
        "error_message": error_message
    })

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.conf import settings
import os
from googletrans import Translator
import PyPDF2
import io
from pathlib import Path

@csrf_exempt
def translate_material(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            file_url = data.get('file_url')
            target_language = data.get('target_language', 'es')  # Default to Spanish
            
            # Map language codes to model names
            model_map = {
                'es': 'Helsinki-NLP/opus-mt-en-es',
                'fr': 'Helsinki-NLP/opus-mt-en-fr',
                'de': 'Helsinki-NLP/opus-mt-en-de',
                'it': 'Helsinki-NLP/opus-mt-en-it',
                'pt': 'Helsinki-NLP/opus-mt-en-pt',
            }
            
            model_name = model_map.get(target_language)
            if not model_name:
                return JsonResponse({
                    'success': False,
                    'error': 'Unsupported target language'
                })
            
            # Get the file path from the URL
            file_path = os.path.join(settings.MEDIA_ROOT, file_url.split('/media/')[-1])

            # Read PDF content
            pdf_text = ""
            try:
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        pdf_text += page.extract_text()
            except Exception as e:
                return JsonResponse({
                    'success': False,
                    'error': f'Error reading PDF: {str(e)}'
                })

            try:
                # Load the translation model and tokenizer
                tokenizer = MarianTokenizer.from_pretrained(model_name)
                model = MarianMTModel.from_pretrained(model_name)
                
                # Set device (GPU if available, else CPU)
                device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
                model = model.to(device)
                
                # Split text into smaller chunks (MarianMT has a max length limit)
                max_chunk_size = 128  # Tokens, not characters
                chunks = []
                words = pdf_text.split()
                current_chunk = []
                
                for word in words:
                    current_chunk.append(word)
                    if len(current_chunk) >= max_chunk_size:
                        chunks.append(' '.join(current_chunk))
                        current_chunk = []
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                
                # Translate chunks
                translated_chunks = []
                for chunk in chunks:
                    # Tokenize and translate
                    inputs = tokenizer(chunk, return_tensors="pt", padding=True, truncation=True, max_length=512)
                    inputs = inputs.to(device)
                    
                    translated = model.generate(**inputs)
                    translated_text = tokenizer.batch_decode(translated, skip_special_tokens=True)[0]
                    translated_chunks.append(translated_text)

                translated_text = ' '.join(translated_chunks)

                # Create the translated PDF
                original_filename = Path(file_path).stem
                translated_filename = f"{original_filename}_{target_language}_translated.pdf"
                translations_dir = os.path.join(settings.MEDIA_ROOT, 'translations')
                os.makedirs(translations_dir, exist_ok=True)
                translated_file_path = os.path.join(translations_dir, translated_filename)

                # Create PDF with formatting
                doc = SimpleDocTemplate(
                    translated_file_path,
                    pagesize=letter,
                    rightMargin=72,
                    leftMargin=72,
                    topMargin=72,
                    bottomMargin=72
                )

                styles = getSampleStyleSheet()
                normal_style = ParagraphStyle(
                    'CustomNormal',
                    parent=styles['Normal'],
                    fontSize=11,
                    leading=14,
                    spaceBefore=6,
                    spaceAfter=6
                )

                story = []
                for para in translated_text.split('\n'):
                    if para.strip():
                        cleaned_text = para.strip().replace('&', '&amp;')
                        p = Paragraph(cleaned_text, normal_style)
                        story.append(p)

                # Build the PDF
                doc.build(story)
                
                return JsonResponse({
                    'success': True,
                    'translated_file_url': f'/media/translations/{translated_filename}'
                })

            except Exception as e:
                print(f"Translation error: {str(e)}")
                return JsonResponse({
                    'success': False,
                    'error': f'Translation error: {str(e)}'
                })

        except Exception as e:
            print(f"General error: {str(e)}")
            return JsonResponse({
                'success': False,
                'error': f'General error: {str(e)}'
            })

    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

from django.http import JsonResponse
import google.generativeai as genai
from django.conf import settings
from PyPDF2 import PdfReader
import io

def generate_summary(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            material_id = data.get('material_id')
            regenerate = data.get('regenerate', False)  # New parameter
            material = Material.objects.get(id=material_id)

            # Check if summary exists and regenerate is not requested
            existing_summary = MaterialSummary.objects.filter(material=material).first()
            if existing_summary and not regenerate:
                return JsonResponse({
                    'success': True,
                    'summary': existing_summary.summary_text,
                    'key_points': existing_summary.key_points.split('\n')
                })

            # Read PDF content
            pdf_content = ""
            pdf_file = material.file
            pdf_reader = PdfReader(io.BytesIO(pdf_file.read()))
            for page in pdf_reader.pages:
                pdf_content += page.extract_text()

            # Configure Google's Generative AI
            GOOGLE_API_KEY = os.getenv('GOOGLE_API_KEY')
            genai.configure(api_key=GOOGLE_API_KEY)
            model = genai.GenerativeModel('gemini-1.5-pro')

            # ... rest of the code remains the same until saving ...

            summary_prompt = f"""
            Please analyze the following educational material and provide a comprehensive summary. 
            Focus on:
            1. Main concepts and theories
            2. Key arguments and their supporting evidence
            3. Important definitions and terminology
            4. Practical applications and examples
            5. Relationships between different concepts

            Text to summarize:
            {pdf_content}

            Please structure the summary in a clear, educational format suitable for students.
            """

            key_points_prompt = f"""
            From the following educational material, extract the most important learning points.
            For each point:
            1. Focus on core concepts that students must understand
            2. Include any critical formulas, definitions, or principles
            3. Highlight practical applications
            4. Note any common misconceptions or important clarifications
            5. Emphasize connections between different topics

            Text to analyze:
            {pdf_content}

            Please provide the key points in a bullet-point format, with each point being clear and actionable for learning.
            """

            try:
                # Generate summary with error handling and retries
                max_retries = 3
                for attempt in range(max_retries):
                    try:
                        summary_response = model.generate_content(summary_prompt)
                        summary_text = summary_response.text
                        break
                    except Exception as e:
                        if attempt == max_retries - 1:
                            raise Exception(f"Failed to generate summary after {max_retries} attempts")
                        continue

                # Generate key points with error handling and retries
                for attempt in range(max_retries):
                    try:
                        key_points_response = model.generate_content(key_points_prompt)
                        key_points = key_points_response.text.split('\n')
                        # Filter out empty lines and clean up bullet points
                        key_points = [point.strip().lstrip('•-*').strip() for point in key_points if point.strip()]
                        break
                    except Exception as e:
                        if attempt == max_retries - 1:
                            raise Exception(f"Failed to generate key points after {max_retries} attempts")
                        continue

                # Post-process the summary
                summary_text = summary_text.replace('\n', '<br>')
                
                # Update or create summary
                if existing_summary:
                    existing_summary.summary_text = summary_text
                    existing_summary.key_points = '\n'.join(key_points)
                    existing_summary.save()
                else:
                    MaterialSummary.objects.create(
                        material=material,
                        summary_text=summary_text,
                        key_points='\n'.join(key_points)
                    )

                return JsonResponse({
                    'success': True,
                    'summary': summary_text,
                    'key_points': key_points
                })

            except Exception as e:
                return JsonResponse({
                    'success': False,
                    'error': f"Processing Error: {str(e)}"
                })

        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': f"System Error: {str(e)}"
            })

    return JsonResponse({'success': False, 'error': 'Invalid request method'})

from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json

@csrf_exempt
def log_violation(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        quiz_id = data.get('quiz_id')
        violation_type = data.get('violation_type')
        
        # Log the violation (you can store this in your database)
        # Example: Violation.objects.create(
        #     quiz_id=quiz_id,
        #     student=request.user,
        #     violation_type=violation_type
        # )
        
        return JsonResponse({'status': 'success'})
    return JsonResponse({'status': 'error'}, status=400)

def question_bank(request, course_id):
    # Check if user is logged in
    if 'custom_user_id' not in request.session:
        return redirect('login')
    
    custom_user_id = request.session['custom_user_id']
    
    # Check if student is enrolled in this course
    if not Enrollment.objects.filter(student_id=custom_user_id, course_id=course_id).exists():
        messages.error(request, "You are not enrolled in this course.")
        return redirect('student_dashboard')
    
    # Get the course and its generated questions
    course = get_object_or_404(Course, id=course_id)
    question_papers = GeneratedQuestionPaper.objects.filter(course_id=course_id)
    
    context = {
        'course': course,
        'question_papers': question_papers,
        'custom_user': CustomUser.objects.get(id=custom_user_id)
    }
    
    return render(request, 'question_bank.html', context)

def generate_answer_key_pdf(questions_data, filename):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=72,
        leftMargin=72,
        topMargin=72,
        bottomMargin=72
    )

    elements = []
    styles = getSampleStyleSheet()
    
    # Add title
    title = Paragraph(f"Answer Key: {filename.split('_')[0]}", styles['Heading1'])
    elements.append(title)
    elements.append(Spacer(1, 12))

    # Process each section
    for section_index, section in enumerate(questions_data):
        section_title = Paragraph(
            f"Section {section_index + 1} ({section['marks']} marks per question)", 
            styles['Heading2']
        )
        elements.append(section_title)
        elements.append(Spacer(1, 12))

        # Process each question
        for q_index, question in enumerate(section['questions']):
            # Question text
            q_text = f"Q{q_index + 1}. {question['question']}"
            elements.append(Paragraph(q_text, styles['Heading3']))
            
            # Answer components
            answer = question['answer']
            
            # Keywords
            elements.append(Paragraph("Keywords:", styles['Heading4']))
            keywords_text = ", ".join(answer['keywords'])
            elements.append(Paragraph(keywords_text, styles['Normal']))
            elements.append(Spacer(1, 6))
            
            # Main points
            elements.append(Paragraph("Main Points:", styles['Heading4']))
            for point in answer['main_points']:
                elements.append(Paragraph(f"• {point}", styles['Normal']))
            elements.append(Spacer(1, 6))
            
            # Detailed explanation
            elements.append(Paragraph("Detailed Explanation:", styles['Heading4']))
            elements.append(Paragraph(answer['explanation'], styles['Normal']))
            elements.append(Spacer(1, 6))
            
            # Examples (if any)
            if 'examples' in answer and answer['examples']:
                elements.append(Paragraph("Examples:", styles['Heading4']))
                for example in answer['examples']:
                    elements.append(Paragraph(f"• {example}", styles['Normal']))
            
            elements.append(Spacer(1, 12))

    # Build PDF
    doc.build(elements)
    buffer.seek(0)
    return buffer

from django.utils import timezone
from datetime import timedelta
import random
from transformers import pipeline
from django.db.models import Q

def check_final_exam_eligibility(request):
    """Check if student is eligible for final exam and return exam details"""
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        return JsonResponse({'eligible': False, 'message': 'Please log in'})

    try:
        student = CustomUser.objects.get(id=custom_user_id)
        enrollments = Enrollment.objects.filter(student=student)
        
        eligible_exams = []
        for enrollment in enrollments:
            course = enrollment.course
            course_end_date = course.ending_date
            today = timezone.now().date()
            
            # Check if course is completed and 7 days have passed
            if course_end_date and today >= (course_end_date + timedelta(days=7)):
                # Check if student hasn't taken final exam yet
                if not FinalExam.objects.filter(student=student, course=course).exists():
                    eligible_exams.append({
                        'course_id': course.id,
                        'course_name': course.course_name,
                        'exam_date': (course_end_date + timedelta(days=7)).strftime('%Y-%m-%d')
                    })

        return JsonResponse({
            'eligible': len(eligible_exams) > 0,
            'exams': eligible_exams
        })

    except CustomUser.DoesNotExist:
        return JsonResponse({'eligible': False, 'message': 'Student not found'})

def final_exam_setup(request, course_id):
    course = get_object_or_404(Course, id=course_id)
    custom_user_id = request.session.get('custom_user_id')
    
    if not custom_user_id:
        return redirect('login')
    
    student = get_object_or_404(CustomUser, id=custom_user_id)
    
    # Check if student has already taken the exam
    existing_exam = FinalExam.objects.filter(course=course, student=student).first()
    if existing_exam and existing_exam.completed:
        messages.warning(request, 'You have already completed this exam.')
        return redirect('student_dashboard')
    
    context = {
        'course': course,
        'student': student
    }
    
    return render(request, 'final_exam_setup.html', context)

def start_final_exam(request, course_id):
    try:
        course = get_object_or_404(Course, id=course_id)
        custom_user_id = request.session.get('custom_user_id')
        
        if not custom_user_id:
            messages.error(request, 'Please log in to take the exam.')
            return redirect('login')
            
        student = get_object_or_404(CustomUser, id=custom_user_id)
        
        # Check if student has already taken the exam
        existing_exam = FinalExam.objects.filter(
            course=course, 
            student=student, 
            completed=True
        ).first()
        
        if existing_exam:
            messages.warning(request, 'You have already completed this exam.')
            return redirect('student_dashboard')
        
        # Get all question papers for this course
        question_papers = GeneratedQuestionPaper.objects.filter(course=course)
        
        if not question_papers.exists():
            messages.error(request, 'No questions available for this exam.')
            return redirect('student_dashboard')
        
        # Create new exam instance
        exam = FinalExam.objects.create(
            course=course,
            student=student,
            start_time=timezone.now()
        )
        
        # Collect all questions from all papers
        all_questions = []
        for paper in question_papers:
            try:
                # Extract text from PDF
                pdf_reader = PyPDF2.PdfReader(paper.pdf_file.path)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                
                # Split text into individual questions (assuming questions are numbered)
                question_pattern = r'(?:\d+\.|Q\d+\.)(.*?)(?=(?:\d+\.|Q\d+\.|$))'
                questions = re.findall(question_pattern, text, re.DOTALL)
                
                # Add each question with its source paper
                for q in questions:
                    q = q.strip()
                    if q:  # Only add non-empty questions
                        all_questions.append({
                            'text': q,
                            'paper': paper,
                            'answer_key': extract_answer_key(paper.answer_key_file.path) if paper.answer_key_file else ""
                        })
                
            except Exception as e:
                print(f"Error processing paper {paper.id}: {str(e)}")
                continue
        
        if len(all_questions) < 10:
            exam.delete()
            messages.error(request, 'Not enough questions available for the exam.')
            return redirect('student_dashboard')
        
        # Randomly select exactly 10 questions
        selected_questions = random.sample(all_questions, 10)
        
        # Create question instances
        exam_questions = []
        for question in selected_questions:
            exam_question = FinalExamQuestion.objects.create(
                exam=exam,
                question_paper=question['paper'],
                question_text=question['text'],
                answer_key=question['answer_key']
            )
            exam_questions.append(exam_question)
        
        return render(request, 'take_final_exam.html', {
            'exam': exam,
            'exam_questions': exam_questions,
            'time_remaining': 3600,  # 1 hour in seconds
            'course': course
        })
        
    except Exception as e:
        print(f"Error starting exam: {str(e)}")
        messages.error(request, 'An error occurred while starting the exam.')
        return redirect('student_dashboard')

def extract_answer_key(answer_key_path):
    try:
        pdf_reader = PyPDF2.PdfReader(answer_key_path)
        answer_text = ""
        for page in pdf_reader.pages:
            answer_text += page.extract_text()
        return answer_text.strip()
    except Exception as e:
        print(f"Error extracting answer key: {str(e)}")
        return ""

def take_final_exam(request, exam_id):
    exam = get_object_or_404(FinalExam, id=exam_id)
    
    if request.method == 'POST':
        # Handle exam submission
        answers = request.POST.getlist('answers[]')
        questions = FinalExamQuestion.objects.filter(exam=exam)
        
        # Initialize AI model for evaluation
        qa_evaluator = pipeline("question-answering")
        
        total_score = 0
        for question, answer in zip(questions, answers):
            # Evaluate answer using AI
            result = qa_evaluator(
                question=question.question_text,
                context=answer,
                answer_key=question.answer_key
            )
            
            # Calculate score based on AI evaluation
            score = result['score'] * 10  # Convert to 10-point scale
            question.marks_obtained = score
            question.student_answer = answer
            question.is_correct = score >= 4  # 40% threshold
            question.save()
            
            total_score += score

        # Calculate final score
        exam.score = total_score / len(questions)
        exam.completed = True
        exam.end_time = timezone.now()
        exam.save()

        # Generate certificate if score >= 40%
        if exam.score >= 40:
            generate_certificate(request, exam.course.id)
            messages.success(request, 'Congratulations! You have passed the exam.')
        else:
            messages.warning(request, 'Sorry, you did not achieve the passing score of 40%.')

        return redirect('exam_results', exam_id=exam.id)

    # GET request - show exam
    if exam.completed:
        return redirect('exam_results', exam_id=exam.id)

    questions = FinalExamQuestion.objects.filter(exam=exam)
    time_remaining = 3600 - (timezone.now() - exam.start_time).seconds

    return render(request, 'take_final_exam.html', {
        'exam': exam,
        'questions': questions,
        'time_remaining': time_remaining
    })

def exam_results(request, exam_id):
    exam = get_object_or_404(FinalExam, id=exam_id)
    questions = FinalExamQuestion.objects.filter(exam=exam)
    
    return render(request, 'exam_results.html', {
        'exam': exam,
        'questions': questions
    })

def extract_questions_from_pdf(pdf_file):
    questions = []
    try:
        # Open PDF file
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        # Extract text from all pages
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()

        # Split text into questions (assuming questions start with numbers)
        # Modify this pattern based on your PDF structure
        question_pattern = r'\d+\.(.*?)(?=\d+\.|$)'
        extracted_questions = re.findall(question_pattern, text, re.DOTALL)
        
        # Clean up questions
        questions = [q.strip() for q in extracted_questions if q.strip()]
        
    except Exception as e:
        print(f"Error extracting questions: {str(e)}")
    
    return questions

from django.http import HttpResponse

def ping(request):
    """Simple endpoint for connection testing"""
    return HttpResponse("pong")

def submit_exam(request, exam_id):
    exam = get_object_or_404(FinalExam, id=exam_id)
    
    # Get all questions for this exam
    questions = FinalExamQuestion.objects.filter(exam=exam)
    
    total_marks = 0
    for question in questions:  # Changed from 'for question in question:'
        answer = request.POST.get(f'answer_{question.id}')
        
        # Check if answer exists
        if not answer or answer.strip() == "":
            question.marks_obtained = 0
        else:
            question.marks_obtained = grade_exam_question(question, answer)
            
        question.student_answer = answer
        question.save()
        total_marks += question.marks_obtained
    
    exam.score = total_marks
    exam.completed = True
    exam.end_time = timezone.now()
    exam.save()
    
    return redirect('exam_results', exam_id=exam.id)

def grade_exam_question(question, answer):
    if not answer or answer.strip() == "":
        return 0  # Return 0 if no answer provided
        
    # ... rest of grading logic ...

def check_existing_exam(request, course_id):
    try:
        custom_user_id = request.session.get('custom_user_id')
        if not custom_user_id:
            return JsonResponse({'exists': False, 'message': 'Please log in'})
            
        student = get_object_or_404(CustomUser, id=custom_user_id)
        course = get_object_or_404(Course, id=course_id)
        
        # Check if student has already started or completed an exam
        existing_exam = FinalExam.objects.filter(
            course=course,
            student=student
        ).exists()
        
        return JsonResponse({
            'exists': existing_exam,
            'message': 'Exam already exists' if existing_exam else 'No existing exam'
        })
        
    except Exception as e:
        return JsonResponse({
            'exists': False,
            'message': str(e)
        })


import numpy as np
from django.http import JsonResponse
from .models import StudentFaceData

def face_capture_view(request):
    # Check if user is logged in and hasn't registered face yet
    if not request.session.get('custom_user_id'):
        return redirect('login')
    
    user = CustomUser.objects.get(id=request.session['custom_user_id'])
    face_data = StudentFaceData.objects.filter(user=user).first()
    
    if face_data and face_data.is_face_captured:
        return redirect('available_courses')
        
    return render(request, 'face_capture.html')

@csrf_exempt
def save_face_data(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'})
    
    try:
        # Get the uploaded image
        face_image = request.FILES.get('face_image')
        if not face_image:
            return JsonResponse({'success': False, 'error': 'No image provided'})

        # Convert image to numpy array
        image_array = np.frombuffer(face_image.read(), np.uint8)
        image = cv2.imdecode(image_array, cv2.IMREAD_COLOR)

        # Initialize MediaPipe Face Detection
        mp_face_detection = mp.solutions.face_detection
        mp_drawing = mp.solutions.drawing_utils

        with mp_face_detection.FaceDetection(min_detection_confidence=0.5) as face_detection:
            # Convert the BGR image to RGB
            image_rgb = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
            results = face_detection.process(image_rgb)

            if not results.detections:
                return JsonResponse({'success': False, 'error': 'No face detected'})
            
            if len(results.detections) > 1:
                return JsonResponse({'success': False, 'error': 'Multiple faces detected'})

            # Get the face bounding box
            detection = results.detections[0]
            bboxC = detection.location_data.relative_bounding_box
            ih, iw, _ = image.shape
            x, y, w, h = int(bboxC.xmin * iw), int(bboxC.ymin * ih), \
                        int(bboxC.width * iw), int(bboxC.height * ih)

            # Extract and save face region
            face_roi = image[y:y+h, x:x+w]
            
            # Save the face image
            _, img_encoded = cv2.imencode('.jpg', face_roi)
            face_image_content = ContentFile(img_encoded.tobytes())

            # Save to database
            user = CustomUser.objects.get(id=request.session.get('custom_user_id'))
            face_data, created = StudentFaceData.objects.get_or_create(user=user)
            face_data.face_encoding = img_encoded.tobytes()
            face_data.is_face_captured = True
            face_data.save()

            return JsonResponse({'success': True})

    except Exception as e:
        print(f"Error saving face data: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

def get_face_encoding(image):
    """Extract face encoding using OpenCV and face-recognition"""
    # Convert to grayscale
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    
    # Load face detector
    face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
    
    # Detect faces
    faces = face_cascade.detectMultiScale(gray, 1.3, 5)
    
    if len(faces) == 0:
        raise ValueError("No face detected")
    if len(faces) > 1:
        raise ValueError("Multiple faces detected")
        
    # Get the face ROI
    x, y, w, h = faces[0]
    face_roi = gray[y:y+h, x:x+w]
    
    # Resize to standard size and flatten
    face_encoding = cv2.resize(face_roi, (128, 128)).flatten()
    
    # Normalize the encoding
    return face_encoding / np.linalg.norm(face_encoding)

def check_face_similarity(new_encoding, stored_encoding, threshold=0.7):
    """Compare face encodings using cosine similarity"""
    similarity = 1 - cosine(new_encoding, stored_encoding)
    return similarity > threshold

def is_face_already_registered(new_encoding):
    """Check if face is already registered by any user"""
    all_face_data = StudentFaceData.objects.filter(is_face_captured=True)
    
    for face_data in all_face_data:
        stored_image = cv2.imdecode(
            np.frombuffer(face_data.face_encoding, np.uint8),
            cv2.IMREAD_COLOR
        )
        stored_encoding = get_face_encoding(stored_image)
        
        if check_face_similarity(new_encoding, stored_encoding):
            return True
    return False

def compare_faces(face1, face2, threshold=0.5):
    """
    Compare two face images using OpenCV's built-in methods
    Returns True if faces match, False otherwise
    """
    # Initialize SIFT detector
    sift = cv2.SIFT_create()
    
    # Detect keypoints and descriptors
    kp1, des1 = sift.detectAndCompute(face1, None)
    kp2, des2 = sift.detectAndCompute(face2, None)
    
    if des1 is None or des2 is None or len(des1) == 0 or len(des2) == 0:
        return False
    
    # Create BFMatcher object
    bf = cv2.BFMatcher()
    matches = bf.knnMatch(des1, des2, k=2)
    
    # Apply ratio test
    good_matches = []
    for m, n in matches:
        if m.distance < threshold * n.distance:
            good_matches.append(m)
    
    # If we have enough good matches, consider it the same face
    return len(good_matches) > 10

@csrf_exempt
def save_face_data(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'})
    
    try:
        # Get the uploaded image
        face_image = request.FILES.get('face_image')
        if not face_image:
            return JsonResponse({'success': False, 'error': 'No image provided'})

        # Convert uploaded image to numpy array
        image_array = np.frombuffer(face_image.read(), np.uint8)
        new_face = cv2.imdecode(image_array, cv2.IMREAD_GRAYSCALE)  # Convert to grayscale
        
        if new_face is None:
            return JsonResponse({'success': False, 'error': 'Invalid image data'})

        # Detect face in the new image
        face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        faces = face_cascade.detectMultiScale(new_face, 1.3, 5)
        
        if len(faces) == 0:
            return JsonResponse({'success': False, 'error': 'No face detected in image'})
        if len(faces) > 1:
            return JsonResponse({'success': False, 'error': 'Multiple faces detected'})

        # Get the face region
        x, y, w, h = faces[0]
        new_face_roi = new_face[y:y+h, x:x+w]
        
        # Check against all existing faces
        existing_faces = StudentFaceData.objects.filter(is_face_captured=True)
        for face_data in existing_faces:
            # Convert stored face data back to image
            stored_face_bytes = np.frombuffer(face_data.face_encoding, np.uint8)
            stored_face = cv2.imdecode(stored_face_bytes, cv2.IMREAD_GRAYSCALE)
            
            if compare_faces(new_face_roi, stored_face):
                return JsonResponse({
                    'success': False,
                    'error': 'This face is already registered with another account'
                })

        # If we get here, the face is unique - save it
        user = CustomUser.objects.get(id=request.session.get('custom_user_id'))
        face_data, created = StudentFaceData.objects.get_or_create(user=user)
        
        # Store the face ROI
        _, img_encoded = cv2.imencode('.jpg', new_face_roi)
        face_data.face_encoding = img_encoded.tobytes()
        face_data.is_face_captured = True
        face_data.save()

        return JsonResponse({'success': True})

    except Exception as e:
        print(f"Error saving face data: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

def face_capture_view(request):
    if not request.session.get('custom_user_id'):
        return redirect('login')
    
    user = CustomUser.objects.get(id=request.session['custom_user_id'])
    face_data = StudentFaceData.objects.filter(user=user).first()
    
    if face_data and face_data.is_face_captured:
        return redirect('available_courses')
        
    return render(request, 'face_capture.html')

def check_face_similarity(new_face, stored_face, threshold=0.7):
    """
    Compare two face images and return similarity score
    """
    try:
        # Resize images to same size
        new_face = cv2.resize(new_face, (128, 128))
        stored_face = cv2.resize(stored_face, (128, 128))
        
        # Calculate similarity using normalized correlation
        similarity = cv2.matchTemplate(new_face, stored_face, cv2.TM_CCOEFF_NORMED)[0][0]
        
        return similarity > threshold
    except Exception as e:
        print(f"Error in face comparison: {str(e)}")
        return False

@csrf_exempt
def save_face_data(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'})
    
    try:
        # Get the uploaded image
        face_image = request.FILES.get('face_image')
        if not face_image:
            return JsonResponse({'success': False, 'error': 'No image provided'})

        # Convert image to numpy array
        image_array = np.frombuffer(face_image.read(), np.uint8)
        new_face = cv2.imdecode(image_array, cv2.IMREAD_GRAYSCALE)
        
        if new_face is None:
            return JsonResponse({'success': False, 'error': 'Invalid image data'})

        # Detect face in the new image
        face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
        faces = face_cascade.detectMultiScale(new_face, 1.1, 4)
        
        if len(faces) == 0:
            return JsonResponse({'success': False, 'error': 'No face detected in image'})
        if len(faces) > 1:
            return JsonResponse({'success': False, 'error': 'Multiple faces detected'})

        # Extract the face region
        x, y, w, h = faces[0]
        new_face_roi = new_face[y:y+h, x:x+w]

        # Check against all existing faces in database
        existing_faces = StudentFaceData.objects.filter(is_face_captured=True).exclude(
            user_id=request.session.get('custom_user_id')
        )

        for existing_face in existing_faces:
            try:
                # Convert stored face data back to image
                stored_face_bytes = np.frombuffer(existing_face.face_encoding, np.uint8)
                stored_face = cv2.imdecode(stored_face_bytes, cv2.IMREAD_GRAYSCALE)
                
                if stored_face is not None and check_face_similarity(new_face_roi, stored_face):
                    return JsonResponse({
                        'success': False,
                        'error': 'This face is already registered with another account. Each student must have a unique face.'
                    })
            except Exception as e:
                print(f"Error processing stored face: {str(e)}")
                continue

        # If we get here, the face is unique - save it
        try:
            user = CustomUser.objects.get(id=request.session.get('custom_user_id'))
            face_data, created = StudentFaceData.objects.get_or_create(user=user)
            
            # Store the face ROI
            _, img_encoded = cv2.imencode('.jpg', new_face_roi)
            face_data.face_encoding = img_encoded.tobytes()
            face_data.is_face_captured = True
            face_data.save()

            return JsonResponse({
                'success': True,
                'message': 'Face registered successfully!'
            })

        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': f'Error saving face data: {str(e)}'
            })

    except Exception as e:
        print(f"Error in face capture: {str(e)}")
        return JsonResponse({
            'success': False,
            'error': f'Face capture failed: {str(e)}'
        })

def face_capture_view(request):
    """View for the face capture page"""
    if not request.session.get('custom_user_id'):
        return redirect('login')
    
    try:
        user = CustomUser.objects.get(id=request.session['custom_user_id'])
        face_data = StudentFaceData.objects.filter(user=user).first()
        
        if face_data and face_data.is_face_captured:
            return redirect('available_courses')
            
        return render(request, 'face_capture.html')
        
    except CustomUser.DoesNotExist:
        return redirect('login')
    except Exception as e:
        print(f"Error in face_capture_view: {str(e)}")
        return redirect('login')
from django.shortcuts import render

def whiteboard(request):
    session_id = request.GET.get('session') or str(uuid.uuid4())
    teacher_id = request.session.get('teacher_id')
    
    if teacher_id:
        teacher = get_object_or_404(Teacher, id=teacher_id)
        context = {
            'is_teacher': True,
            'session_id': session_id,
            'teacher_courses': TeacherCourse.objects.filter(teacher=teacher)
        }
    else:
        context = {
            'is_teacher': False,
            'session_id': session_id
        }
    
    return render(request, 'whiteboard.html', context)


from django.http import JsonResponse
import base64
import pytesseract
from PIL import Image
import io

def process_ocr(request):
    if request.method == 'POST':
        image_data = request.POST.get('image_data')
        if image_data:
            # Decode the image data
            image_data = image_data.split(',')[1]  # Remove the data URL part
            image = Image.open(io.BytesIO(base64.b64decode(image_data)))

            # Perform OCR
            text = pytesseract.image_to_string(image)

            return JsonResponse({'success': True, 'text': text})
    
    return JsonResponse({'success': False, 'error': 'Invalid request'})

import cv2
import numpy as np

from scipy.spatial.distance import cosine
from django.shortcuts import render, redirect
from .models import Attendance, StudentFaceData, ClassSchedule
from django.utils import timezone

def mark_attendance(request, schedule_id):
    custom_user_id = request.session.get('custom_user_id')
    if not custom_user_id:
        messages.error(request, "You need to log in to join the class.")
        return redirect('login')

    try:
        student = CustomUser.objects.get(id=custom_user_id)
        class_schedule = ClassSchedule.objects.get(id=schedule_id)
        current_time = timezone.localtime(timezone.now())

        # Check if student has registered face data
        try:
            face_data = StudentFaceData.objects.get(user=student)
            if not face_data.is_face_captured:
                messages.error(request, "Please register your face first.")
                return redirect('face_capture')
        except StudentFaceData.DoesNotExist:
            messages.error(request, "Please register your face first.")
            return redirect('face_capture')

        # Check if within class time
        if class_schedule.start_time <= current_time.time() <= class_schedule.end_time and current_time.date() == class_schedule.date:
            # Show face verification page
            return render(request, 'face_verification.html', {
                'schedule_id': schedule_id,
                'meeting_link': class_schedule.meeting_link
            })
        else:
            # If outside class time, mark as absent
            Attendance.objects.create(
                student=student,
                class_schedule=class_schedule,
                check_in_time=current_time,
                status='absent'
            )
            messages.error(request, "Class is not active at this time.")
            return redirect('view_scheduled_classes')

    except CustomUser.DoesNotExist:
        messages.error(request, "Student not found.")
        return redirect('login')
    except ClassSchedule.DoesNotExist:
        messages.error(request, "Class schedule not found.")
        return redirect('student_dashboard')
    except Exception as e:
        messages.error(request, str(e))
        return redirect('view_scheduled_classes')

@csrf_exempt
def verify_face_and_mark_attendance(request, schedule_id):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'})

    try:
        # Get current user and schedule
        student = CustomUser.objects.get(id=request.session.get('custom_user_id'))
        class_schedule = ClassSchedule.objects.get(id=schedule_id)
        stored_face_data = StudentFaceData.objects.get(user=student)

        # Get captured face image
        face_image = request.FILES.get('face_image')
        if not face_image:
            return JsonResponse({'success': False, 'error': 'No image provided'})

        # Convert to numpy array
        image_array = np.frombuffer(face_image.read(), np.uint8)
        current_image = cv2.imdecode(image_array, cv2.IMREAD_COLOR)

        # Initialize MediaPipe Face Detection
        mp_face_detection = mp.solutions.face_detection
        with mp_face_detection.FaceDetection(min_detection_confidence=0.5) as face_detection:
            # Process current image
            current_rgb = cv2.cvtColor(current_image, cv2.COLOR_BGR2RGB)
            current_results = face_detection.process(current_rgb)

            if not current_results.detections:
                return JsonResponse({'success': False, 'error': 'No face detected in current image'})

            # Get stored face image
            stored_face_bytes = np.frombuffer(stored_face_data.face_encoding, np.uint8)
            stored_image = cv2.imdecode(stored_face_bytes, cv2.IMREAD_COLOR)

            # Process stored image
            stored_rgb = cv2.cvtColor(stored_image, cv2.COLOR_BGR2RGB)
            stored_results = face_detection.process(stored_rgb)

            if not stored_results.detections:
                return JsonResponse({'success': False, 'error': 'No face detected in stored image'})

            # Extract and compare face landmarks
            current_landmarks = current_results.detections[0].location_data.relative_keypoints
            stored_landmarks = stored_results.detections[0].location_data.relative_keypoints

            # Calculate similarity based on landmark positions
            similarity_score = calculate_landmark_similarity(current_landmarks, stored_landmarks)
            
            threshold = 0.75
            if similarity_score >= threshold:
                # Mark attendance as present
                Attendance.objects.create(
                    student=student,
                    class_schedule=class_schedule,
                    check_in_time=timezone.now(),
                    status='present'
                )
                return JsonResponse({
                    'success': True,
                    'meeting_link': class_schedule.meeting_link
                })
            else:
                # Mark as absent
                Attendance.objects.create(
                    student=student,
                    class_schedule=class_schedule,
                    check_in_time=timezone.now(),
                    status='absent'
                )
                return JsonResponse({
                    'success': False,
                    'error': f'Face verification failed. Score: {similarity_score:.2f}'
                })

    except Exception as e:
        print(f"Error in face verification: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

def calculate_landmark_similarity(landmarks1, landmarks2):
    """Calculate similarity between two sets of facial landmarks"""
    try:
        # Convert landmarks to numpy arrays
        points1 = np.array([[p.x, p.y] for p in landmarks1])
        points2 = np.array([[p.x, p.y] for p in landmarks2])
        
        # Calculate Euclidean distances between corresponding points
        distances = np.linalg.norm(points1 - points2, axis=1)
        
        # Convert distances to similarity score (inverse relationship)
        similarity = 1 / (1 + np.mean(distances))
        
        return similarity
    except Exception as e:
        print(f"Error calculating landmark similarity: {str(e)}")
        return 0.0



import os
from django.conf import settings
from ml_code.create_db import create_data
from ml_code.face_recognition import face_recognize
from .models import StudentFaceData, CustomUser

@csrf_exempt
def save_face_data(request):
    if request.method != 'POST':
        return JsonResponse({'success': False, 'error': 'Invalid request method'})
    
    try:
        user = CustomUser.objects.get(id=request.session.get('custom_user_id'))
        
        # Create directory for user if it doesn't exist
        user_dir = f"ml_code/database/{user.username}"
        if not os.path.exists(user_dir):
            os.makedirs(user_dir)

        # Call create_data function to capture and save face images
        face_image_path = create_data(str(user.id))

        if face_image_path and os.path.exists(face_image_path):
            # Save to StudentFaceData model
            with open(face_image_path, 'rb') as f:
                face_data, created = StudentFaceData.objects.get_or_create(user=user)
                face_data.face_image.save(
                    f'{user.username}_face.jpg',
                    ContentFile(f.read()),
                    save=True
                )
                face_data.is_face_captured = True
                face_data.save()

            return JsonResponse({
                'success': True, 
                'message': 'Face registered successfully'
            })
        else:
            return JsonResponse({
                'success': False,
                'error': 'Face image capture failed'
            })

    except Exception as e:
        print(f"Error saving face data: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

@csrf_exempt
def verify_face_and_mark_attendance(request, schedule_id):
    if request.method != 'POST':
    
        return JsonResponse({'success': False, 'error': 'Invalid request method'})

    try:
        print("verify_face_and_mark_attendance")
        student = CustomUser.objects.get(id=request.session.get('custom_user_id'))
        class_schedule = ClassSchedule.objects.get(id=schedule_id)

        # Get the student's face data
        # try:
        #     face_data = StudentFaceData.objects.get(user=student)
        #     if not face_data.is_face_captured:
        #         return JsonResponse({
        #             'success': False,
        #             'error': 'Face data not registered. Please register your face first.'
        #         })
        # except StudentFaceData.DoesNotExist:
        #     return JsonResponse({
        #         'success': False,
        #         'error': 'Face data not found. Please register your face first.'
        #     })

        # Verify face using face_recognition.py
        verification_result = face_recognize(str(student.id))
        
        if verification_result:
            # Mark attendance as present
            # attendance = Attendance.objects.create(
            #     student=student,
            #     class_schedule=class_schedule,
            #     check_in_time=timezone.now(),
            #     status='present'
            # )
            return JsonResponse({
                'success': True,
                'meeting_link': class_schedule.meeting_link,
                'message': 'Attendance marked successfully'
            })
        else:
            # Mark as absent
            # attendance = Attendance.objects.create(
            #     student=student,
            #     class_schedule=class_schedule,
            #     check_in_time=timezone.now(),
            #     status='absent'
            # )
            return JsonResponse({
                'success': False,
                'error': 'Face verification failed. Attendance marked as absent.'
            })

    except Exception as e:
        print(f"Error in face verification: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

from django.shortcuts import render, redirect, get_object_or_404
from django.contrib import messages
from django.http import JsonResponse
from django.core.paginator import Paginator
from .models import ParentTeacherMessage, Teacher, Parent
from django.db.models import Q
from django.utils import timezone

def parent_message_center(request):
    parent_id = request.session.get('parent_id')
    if not parent_id:
        return redirect('login')

    parent = get_object_or_404(Parent, id=parent_id)
    
    # Get all messages (both sent and received) for this parent
    messages_list = ParentTeacherMessage.objects.filter(
        parent=parent
    ).select_related('teacher').order_by('-date')
    
    # Get unique teachers this parent has communicated with
    teachers = Teacher.objects.filter(
        Q(parent_messages__parent=parent) | 
        Q(teachercourse__course__students__username=parent.student_username)
    ).distinct()

    # Add sender_type to each message for template rendering
    for message in messages_list:
        message.sender_type = 'teacher' if message.message_type == 'teacher_to_parent' else 'parent'

    # Pagination
    paginator = Paginator(messages_list, 10)
    page = request.GET.get('page')
    messages_page = paginator.get_page(page)

    context = {
        'messages': messages_page,  # Changed from messages_page to match template
        'teachers': teachers,
        'parent': parent,
        'courses': parent.student.enrollments.all(),  # Add courses for the message form
        'unread_count': messages_list.filter(is_read=False, message_type='teacher_to_parent').count()
    }
    return render(request, 'parent_message_center.html', context)

def send_parent_message(request):
    if request.method == 'POST':
        parent_id = request.session.get('parent_id')
        if not parent_id:
            return JsonResponse({'status': 'error', 'message': 'Not logged in'})

        parent = get_object_or_404(Parent, id=parent_id)
        teacher_id = request.POST.get('teacher_id')
        content = request.POST.get('content')
        subject = request.POST.get('subject', 'No Subject')

        if not all([teacher_id, content]):
            return JsonResponse({'status': 'error', 'message': 'Missing required fields'})

        try:
            teacher = Teacher.objects.get(id=teacher_id)
            message = ParentTeacherMessage.objects.create(
                teacher=teacher,
                parent=parent,
                content=content,
                subject=subject,
                message_type='parent_to_teacher'
            )

            return JsonResponse({
                'status': 'success',
                'message': 'Message sent successfully',
                'data': {
                    'id': message.id,
                    'date': message.date.strftime('%Y-%m-%d %H:%M'),
                    'content': message.content,
                    'subject': message.subject,
                    'teacher_name': f"{teacher.first_name} {teacher.last_name}"
                }
            })
        except Exception as e:
            return JsonResponse({'status': 'error', 'message': str(e)})

    return JsonResponse({'status': 'error', 'message': 'Invalid request method'})

def mark_message_read(request, message_id):
    if request.method == 'POST':
        message = get_object_or_404(ParentTeacherMessage, id=message_id)
        message.is_read = True
        message.save()
        return JsonResponse({'status': 'success'})
    return JsonResponse({'status': 'error'})

def delete_message(request, message_id):
    if request.method == 'POST':
        message = get_object_or_404(ParentTeacherMessage, id=message_id)
        message.delete()
        return JsonResponse({'status': 'success'})
    return JsonResponse({'status': 'error'})

from django.shortcuts import get_object_or_404
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from .models import Group, Message, CustomUser

@csrf_exempt
def upload_media_message(request, group_id):
    if request.method == 'POST' and request.FILES.get('media_file'):
        try:
            group = get_object_or_404(Group, id=group_id)
            custom_user_id = request.session.get('custom_user_id')
            user = get_object_or_404(CustomUser, id=custom_user_id)
            
            media_file = request.FILES['media_file']
            message_type = request.POST.get('message_type')
            
            message = Message.objects.create(
                group=group,
                sender=user,
                message_type=message_type,
                media_file=media_file,
                content=None  # No text content for media messages
            )
            
            return JsonResponse({
                'success': True,
                'message_id': message.id
            })
            
        except Exception as e:
            print(f"Error in upload_media_message: {str(e)}")
            return JsonResponse({
                'success': False,
                'error': str(e)
            }, status=400)
            
    return JsonResponse({
        'success': False,
        'error': 'Invalid request'
    }, status=400)


from django.http import JsonResponse
from django.views.decorators.http import require_POST

from .models import Course, CustomUser, WhiteboardShare, Teacher, Notification, TeacherCourse, Enrollment
import json

@require_POST
def share_whiteboard(request):
    try:
        # Get the logged-in teacher from session
        teacher_id = request.session.get('teacher_id')
        if not teacher_id:
            return JsonResponse({
                'success': False,
                'error': 'Teacher not logged in'
            }, status=401)
            
        teacher = get_object_or_404(Teacher, id=teacher_id)
        data = json.loads(request.body)
        course_id = data.get('course_id')
        session_id = data.get('session_id')
        whiteboard_url = data.get('whiteboard_url')

        # Get the course and verify teacher is assigned to it
        course = get_object_or_404(Course, id=course_id)
        teacher_course = TeacherCourse.objects.filter(teacher=teacher, course=course).first()
        if not teacher_course:
            return JsonResponse({
                'success': False,
                'error': 'You are not assigned to this course'
            }, status=403)

        # Get enrolled students - just get all enrollments for the course
        enrolled_students = Enrollment.objects.filter(
            course=course
        ).select_related('student')

        # Create whiteboard share record
        whiteboard_share = WhiteboardShare.objects.create(
            teacher=teacher,
            course=course,
            session_id=session_id,
            whiteboard_url=whiteboard_url
        )

        # Create notifications for enrolled students
        for enrollment in enrolled_students:
            Notification.objects.create(
                user=enrollment.student,
                title=f"New Whiteboard Session - {course.course_name}",
                message=f"Teacher {teacher.first_name} {teacher.last_name} has started a whiteboard session",
                link=whiteboard_url
            )

        return JsonResponse({
            'success': True,
            'course_name': course.course_name,
            'student_count': enrolled_students.count(),
            'teacher_name': f"{teacher.first_name} {teacher.last_name}"
        })

    except json.JSONDecodeError:
        return JsonResponse({
            'success': False,
            'error': 'Invalid JSON data'
        }, status=400)
    except Exception as e:
        print(f"Error in share_whiteboard: {str(e)}")  # Add debugging
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)

def get_notifications(request):
    if not request.session.get('custom_user_id'):
        return JsonResponse({'error': 'Not logged in'}, status=401)
        
    user_id = request.session.get('custom_user_id')
    notifications = Notification.objects.filter(
        user_id=user_id,
        title__contains='Whiteboard Session'
    ).order_by('-created_at')[:5]  # Get last 5 whiteboard notifications
    
    notifications_data = [{
        'id': notif.id,
        'title': notif.title,
        'message': notif.message,
        'link': notif.link,
        'created_at': notif.created_at.isoformat(),
        'is_read': notif.is_read
    } for notif in notifications]
    
    return JsonResponse(notifications_data, safe=False)

from django.http import JsonResponse
from .models import EventSuggestion, CalendarEvent

@require_POST
def send_event_suggestion(request):
    try:
        event_id = request.POST.get('event_id')
        suggestion_text = request.POST.get('suggestion_text')
        parent_id = request.session.get('parent_id')

        if not all([event_id, suggestion_text, parent_id]):
            return JsonResponse({
                'success': False,
                'message': 'Missing required information'
            })

        parent = Parent.objects.get(id=parent_id)
        event = CalendarEvent.objects.get(id=event_id)
        student = CustomUser.objects.get(username=parent.student_username)

        suggestion = EventSuggestion.objects.create(
            parent=parent,
            student=student,
            event=event,
            suggestion_text=suggestion_text
        )

        return JsonResponse({
            'success': True,
            'message': 'Suggestion sent successfully',
            'suggestion': {
                'id': suggestion.id,
                'created_at': suggestion.created_at.strftime('%Y-%m-%d %H:%M:%S')
            }
        })

    except Exception as e:
        return JsonResponse({
            'success': False,
            'message': str(e)
        })

from django.http import JsonResponse
from .models import EventSuggestion

def get_event_suggestions(request, event_id):
    parent_id = request.session.get('parent_id')
    
    if not parent_id:
        return JsonResponse({
            'success': False,
            'message': 'Not authenticated'
        })

    try:
        parent = Parent.objects.get(id=parent_id)
        suggestions = EventSuggestion.objects.filter(
            event_id=event_id,
            parent=parent
        ).order_by('-created_at')
        
        suggestions_data = [{
            'suggestion_text': suggestion.suggestion_text,
            'created_at': suggestion.created_at.strftime('%b %d, %Y %H:%M'),
            'is_read': suggestion.is_read
        } for suggestion in suggestions]
        
        return JsonResponse({
            'success': True,
            'suggestions': suggestions_data
        })
    except Exception as e:
        return JsonResponse({
            'success': False,
            'message': str(e)
        })
from django.http import JsonResponse
from .models import MindMap
import json

def save_mind_map(request):
    if not request.session.get('custom_user_id'):
        return JsonResponse({
            'success': False,
            'error': 'Not authenticated'
        })
        
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            student = CustomUser.objects.get(id=request.session.get('custom_user_id'))
            
            # Validate the data
            if not data.get('name') or not data.get('data'):
                return JsonResponse({
                    'success': False,
                    'error': 'Missing required fields'
                })
            
            # Create new mind map
            mind_map = MindMap.objects.create(
                student=student,
                name=data['name'],
                data=data['data']
            )
            
            return JsonResponse({
                'success': True,
                'map_id': mind_map.id
            })
            
        except CustomUser.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': 'User not found'
            })
        except json.JSONDecodeError:
            return JsonResponse({
                'success': False,
                'error': 'Invalid JSON data'
            })
        except Exception as e:
            print(f"Error saving mind map: {str(e)}")  # Add logging
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
            
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

def get_mind_maps(request):
    if not request.session.get('custom_user_id'):
        return JsonResponse({
            'success': False,
            'error': 'Not authenticated',
            'maps': []
        })
    
    try:
        student = CustomUser.objects.get(id=request.session.get('custom_user_id'))
        maps = MindMap.objects.filter(student=student)
        maps_data = [{
            'id': map.id,
            'name': map.name,
            'data': map.data,
            'created_at': map.created_at.isoformat()
        } for map in maps]
        return JsonResponse({
            'success': True,
            'maps': maps_data
        })
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e),
            'maps': []
        })

def delete_mind_map(request, map_id):
    if not request.session.get('custom_user_id'):
        return JsonResponse({
            'success': False,
            'error': 'Not authenticated'
        })
    
    if request.method == 'DELETE':
        try:
            student = CustomUser.objects.get(id=request.session.get('custom_user_id'))
            mind_map = MindMap.objects.get(id=map_id, student=student)
            mind_map.delete()
            return JsonResponse({'success': True})
        except MindMap.DoesNotExist:
            return JsonResponse({
                'success': False,
                'error': 'Mind map not found'
            })
        except Exception as e:
            return JsonResponse({
                'success': False,
                'error': str(e)
            })
    return JsonResponse({
        'success': False,
        'error': 'Invalid request method'
    })

from gtts import gTTS
import os
from django.http import FileResponse
from django.shortcuts import get_object_or_404
from .models import Material
import PyPDF2
import docx
import tempfile

def text_to_speech(request, material_id):
    try:
        material = get_object_or_404(Material, id=material_id)
        text = extract_text_from_file(material.file.path)
        
        if not text:
            return JsonResponse({'success': False, 'error': 'No text found'})

        # Initialize models and processor
        processor = SpeechT5Processor.from_pretrained("microsoft/speecht5_tts")
        model = SpeechT5ForTextToSpeech.from_pretrained("microsoft/speecht5_tts")
        vocoder = SpeechT5HifiGan.from_pretrained("microsoft/speecht5_hifigan")

        # Load speaker embeddings
        embeddings_dataset = load_dataset("Matthijs/cmu-arctic-xvectors", split="validation")
        speaker_embeddings = torch.tensor(embeddings_dataset[7306]["xvector"]).unsqueeze(0)

        # Process text
        inputs = processor(text=text, return_tensors="pt")

        # Generate speech
        speech = model.generate_speech(inputs["input_ids"], speaker_embeddings, vocoder=vocoder)

        # Create audio directory if it doesn't exist
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'audio')
        os.makedirs(audio_dir, exist_ok=True)

        # Save the audio file
        audio_filename = f'speech_{material_id}.wav'
        audio_path = os.path.join(audio_dir, audio_filename)
        sf.write(audio_path, speech.numpy(), samplerate=16000)

        # Return the URL to the audio file
        audio_url = f'{settings.MEDIA_URL}audio/{audio_filename}'
        return JsonResponse({
            'success': True,
            'audio_url': audio_url
        })

    except Exception as e:
        print(f"Error in text_to_speech: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

def extract_text_from_file(file_path):
    """Extract text from different file types"""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            return extract_from_pdf(file_path)
        elif file_extension in ['.doc', '.docx']:
            return extract_from_word(file_path)
        elif file_extension == '.txt':
            return extract_from_txt(file_path)
        else:
            raise ValueError('Unsupported file type')
    except Exception as e:
        print(f"Error extracting text: {str(e)}")
        return ""

def extract_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_from_word(file_path):
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return '\n'.join(text)

def extract_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

import pyttsx3
import threading
from django.conf import settings

def text_to_speech(request, material_id):
    try:
        material = get_object_or_404(Material, id=material_id)
        text = extract_text_from_file(material.file.path)
        
        if not text:
            return JsonResponse({'success': False, 'error': 'No text found'})

        # Create audio directory if it doesn't exist
        audio_dir = os.path.join(settings.MEDIA_ROOT, 'audio')
        os.makedirs(audio_dir, exist_ok=True)

        # Set audio file path
        audio_filename = f'speech_{material_id}.mp3'
        audio_path = os.path.join(audio_dir, audio_filename)

        def generate_speech():
            engine = pyttsx3.init()
            # Configure voice properties
            engine.setProperty('rate', 150)    # Speed of speech
            engine.setProperty('volume', 0.9)  # Volume (0-1)
            
            # Get available voices and set a female voice if available
            voices = engine.getProperty('voices')
            for voice in voices:
                if "female" in voice.name.lower():
                    engine.setProperty('voice', voice.id)
                    break

            # Generate and save audio
            engine.save_to_file(text, audio_path)
            engine.runAndWait()

        # Run speech generation in a separate thread
        thread = threading.Thread(target=generate_speech)
        thread.start()
        thread.join()  # Wait for completion

        # Return the URL to the audio file
        audio_url = f'{settings.MEDIA_URL}audio/{audio_filename}'
        return JsonResponse({
            'success': True,
            'audio_url': audio_url
        })

    except Exception as e:
        print(f"Error in text_to_speech: {str(e)}")
        return JsonResponse({'success': False, 'error': str(e)})

def create_note(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
        
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        
        if request.method == 'POST':
            title = request.POST.get('title')
            content = request.POST.get('content')
            course_id = request.POST.get('course')
            
            course = Course.objects.get(id=course_id)
            TeacherNote.objects.create(
                teacher=teacher,
                title=title,
                content=content,
                course=course
            )
            messages.success(request, 'Note created successfully!')
            return redirect('view_notes')
        
        # Get courses assigned to this teacher
        courses = Course.objects.filter(course_teachers__teacher=teacher)
        return render(request, 'create_note.html', {
            'form_title': 'Create New Note',
            'courses': courses,
            'first_name': teacher.first_name,
            'last_name': teacher.last_name
        })
        
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found')
        return redirect('login')

def view_notes(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
        
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        notes = TeacherNote.objects.filter(teacher=teacher).order_by('-created_at')
        return render(request, 'view_notes.html', {
            'notes': notes,
            'first_name': teacher.first_name,
            'last_name': teacher.last_name
        })
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found')
        return redirect('login')

def edit_note(request, note_id):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
        
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        note = get_object_or_404(TeacherNote, id=note_id, teacher=teacher)
        
        if request.method == 'POST':
            note.title = request.POST.get('title')
            note.content = request.POST.get('content')
            note.course_id = request.POST.get('course')
            note.save()
            messages.success(request, 'Note updated successfully!')
            return redirect('view_notes')
        
        # Fix the course filter here too
        courses = Course.objects.filter(course_teachers__teacher=teacher)
        return render(request, 'create_note.html', {
            'form_title': 'Edit Note',
            'note': note,
            'courses': courses,
            'first_name': teacher.first_name,
            'last_name': teacher.last_name
        })
        
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found')
        return redirect('login')

def delete_note(request, note_id):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
        
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        note = get_object_or_404(TeacherNote, id=note_id, teacher=teacher)
        note.delete()
        messages.success(request, 'Note deleted successfully!')
        return redirect('view_notes')
        
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found')
        return redirect('login')

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from io import BytesIO
from html import unescape
import re

def download_note_pdf(request, note_id):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
        
    try:
        teacher = Teacher.objects.get(id=teacher_id)
        note = get_object_or_404(TeacherNote, id=note_id, teacher=teacher)
        
        # Create a file-like buffer to receive PDF data
        buffer = BytesIO()
        
        # Create the PDF object using the buffer as its "file"
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        
        # Create the story containing all the elements
        story = []
        
        # Define styles
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            spaceAfter=30
        )
        
        info_style = ParagraphStyle(
            'CourseInfo',
            parent=styles['Normal'],
            fontSize=12,
            textColor=colors.gray,
            spaceAfter=20
        )
        
        content_style = ParagraphStyle(
            'CustomBody',
            parent=styles['Normal'],
            fontSize=12,
            leading=14,
            spaceAfter=12
        )
        
        # Add title
        story.append(Paragraph(note.title, title_style))
        
        # Add course info and dates
        course_info = f"""
        Course: {note.course.course_name}<br/>
        Created: {note.created_at.strftime('%B %d, %Y')}<br/>
        Last Updated: {note.updated_at.strftime('%B %d, %Y')}
        """
        story.append(Paragraph(course_info, info_style))
        story.append(Spacer(1, 20))
        
        # Process the HTML content
        content = note.content
        # Remove HTML tags but preserve line breaks
        content = re.sub(r'<br\s*/?>', '\n', content)
        content = re.sub(r'<[^>]+>', '', content)
        content = unescape(content)  # Convert HTML entities to characters
        
        # Split content into paragraphs and add them
        paragraphs = content.split('\n')
        for para in paragraphs:
            if para.strip():
                story.append(Paragraph(para.strip(), content_style))
        
        # Build the PDF
        doc.build(story)
        
        # Get the value of the BytesIO buffer and write it to the response
        pdf = buffer.getvalue()
        buffer.close()
        
        # Create the HTTP response with PDF content
        response = HttpResponse(content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{note.title.replace(" ", "_")}.pdf"'
        response.write(pdf)
        
        return response
        
    except Teacher.DoesNotExist:
        messages.error(request, 'Teacher not found')
        return redirect('login')
    except Exception as e:
        messages.error(request, f'Error generating PDF: {str(e)}')
        return redirect('view_notes')

from django.shortcuts import render
from .models import CalendarEvent, EventRegistration
import csv
from django.http import HttpResponse

def view_event_registrations(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')

    # Get events created by this teacher
    events = CalendarEvent.objects.filter(created_by_id=teacher_id)
    
    # Filter by selected event
    selected_event = request.GET.get('event')
    if selected_event:
        events = events.filter(id=selected_event)

    # Filter by registration status
    status = request.GET.get('status')
    if status:
        events = events.filter(registrations__status=status)

    context = {
        'events': events,
        'selected_event': selected_event,
        'status': status,
    }
    return render(request, 'view_event_registrations.html', context)

def export_registrations(request):
    event_id = request.GET.get('event')
    if not event_id:
        messages.error(request, 'Please select an event to export')
        return redirect('view_event_registrations')

    try:
        event = get_object_or_404(CalendarEvent, id=event_id)
        registrations = event.registrations.all().select_related('user')

        response = HttpResponse(content_type='text/csv')
        response['Content-Disposition'] = f'attachment; filename="{event.title}_registrations.csv"'

        writer = csv.writer(response)
        writer.writerow(['Student Name', 'Email', 'Contact', 'Registration Date', 'Status'])

        for registration in registrations:
            writer.writerow([
                f"{registration.user.first_name} {registration.user.last_name}",
                registration.user.email,
                registration.contact_number or 'Not provided',
                registration.registration_date.strftime('%Y-%m-%d %H:%M'),
                registration.status
            ])

        return response

    except Exception as e:
        messages.error(request, f'Error exporting registrations: {str(e)}')
        return redirect('view_event_registrations')
    

def teacher_messages(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return redirect('login')
    
    try:
        # Fetch the logged-in teacher
        teacher = Teacher.objects.get(id=teacher_id)
        
        # Fetch assigned courses from TeacherCourse model
        assigned_courses = TeacherCourse.objects.filter(teacher=teacher)
        
        # Get messages where the logged-in teacher is the recipient
        messages = ParentTeacherMessage.objects.filter(
            teacher_id=teacher_id,
            message_type='parent_to_teacher'
        ).select_related('parent')
        
        # Check for existing replies and mark messages as read
        for message in messages:
            message.has_reply = ParentTeacherMessage.objects.filter(
                teacher_id=teacher_id,
                parent=message.parent,
                message_type='teacher_to_parent',
                subject__startswith=f"Re: {message.subject}"
            ).exists()
            
            if not message.is_read:
                message.is_read = True
                message.save()
        
        context = {
            'messages': messages,
            'assigned_courses': assigned_courses,
            'first_name': request.session.get('first_name', ''),
            'last_name': request.session.get('last_name', ''),
        }
        
        return render(request, 'teacher_messages.html', context)
        
    except Exception as e:
        print(f"Error in teacher_messages view: {e}")
        return redirect('login')

# Add new view to get students for a course
def get_course_students(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return JsonResponse({'error': 'Not authenticated'})
    
    try:
        course_id = request.GET.get('course_id')
        if not course_id:
            return JsonResponse({'error': 'Course ID is required'})
        
        # Verify teacher has access to this course
        teacher_course = TeacherCourse.objects.filter(
            teacher_id=teacher_id,
            course_id=course_id
        ).exists()
        
        if not teacher_course:
            return JsonResponse({'error': 'Course not assigned to you'})
        
        # Get enrolled students
        enrolled_students = Enrollment.objects.filter(
            course_id=course_id
        ).select_related('student')
        
        students_data = [{
            'id': enrollment.student.id,
            'name': f"{enrollment.student.first_name} {enrollment.student.last_name}",
            'username': enrollment.student.username
        } for enrollment in enrolled_students]
        
        return JsonResponse({
            'success': True,
            'students': students_data
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)})

# Add new view to send message to parent
@require_POST
def send_message_to_parent(request):
    teacher_id = request.session.get('teacher_id')
    if not teacher_id:
        return JsonResponse({'error': 'Not authenticated'})
    
    try:
        data = json.loads(request.body)
        student_id = data.get('student_id')
        content = data.get('content')
        subject = data.get('subject', 'Message from Teacher')
        
        # Get student's parent
        student = CustomUser.objects.get(id=student_id)
        parent = Parent.objects.get(student_username=student.username)
        
        # Create new message
        message = ParentTeacherMessage.objects.create(
            teacher_id=teacher_id,
            parent=parent,
            subject=subject,
            content=content,
            message_type='teacher_to_parent'
        )
        
        return JsonResponse({
            'success': True,
            'message_id': message.id
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        })
    
@csrf_exempt
def send_teacher_reply(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message_id = data.get('message_id')
            content = data.get('content')
            teacher_id = request.session.get('teacher_id')

            if not all([message_id, content, teacher_id]):
                return JsonResponse({'success': False, 'error': 'Missing required data'})

            # Get the original message
            original_message = ParentTeacherMessage.objects.get(id=message_id)

            # Create reply message
            reply = ParentTeacherMessage.objects.create(
                teacher_id=teacher_id,
                parent=original_message.parent,
                content=content,
                subject=f"Re: {original_message.subject}",
                message_type='teacher_to_parent'
            )

            return JsonResponse({'success': True})

        except Exception as e:
            return JsonResponse({'success': False, 'error': str(e)})

    return JsonResponse({'success': False, 'error': 'Invalid request method'})

def get_message_history(request, message_id):
    if not request.session.get('teacher_id'):
        return JsonResponse({'error': 'Not authenticated'}, status=401)
        
    try:
        # Get the original message
        original_message = ParentTeacherMessage.objects.get(id=message_id)
        
        # Get all messages between this teacher and parent
        messages = ParentTeacherMessage.objects.filter(
            teacher_id=original_message.teacher_id,
            parent=original_message.parent
        ).order_by('date')
        
        messages_data = [{
            'content': msg.content,
            'date': msg.date.strftime('%b %d, %Y %H:%M'),
            'message_type': msg.message_type,
            'subject': msg.subject
        } for msg in messages]
        
        return JsonResponse({
            'success': True,
            'messages': messages_data
        })
        
    except Exception as e:
        return JsonResponse({
            'success': False,
            'error': str(e)
        }, status=500)




from django.http import JsonResponse

def check_material_exists(request):
    title = request.GET.get('title')
    course_id = request.GET.get('course')
    
    exists = Material.objects.filter(
        description__in=[title, f"Note: {title}"],
        course_id=course_id
    ).exists()
    
    return JsonResponse({'exists': exists})


from .plagiarism_checker import PlagiarismAIChecker
from django.core.cache import cache

# Initialize the checker
checker = PlagiarismAIChecker()

@require_POST
def check_plagiarism(request, submission_id):
    submission = get_object_or_404(AssignmentSubmission, id=submission_id)
    
    try:
        # Extract text from current submission
        current_text = extract_text_from_file(submission.file)
        if not current_text:
            return JsonResponse({'success': False, 'error': 'Could not extract text from file'})
        
        # Get all other submissions for the same assignment
        other_submissions = AssignmentSubmission.objects.filter(
            assignment=submission.assignment
        ).exclude(id=submission_id)
        
        highest_similarity = 0
        plagiarism_details = []
        
        # Compare with each other submission
        for other_submission in other_submissions:
            other_text = extract_text_from_file(other_submission.file)
            if other_text:
                result = checker.check_plagiarism(current_text, other_text)
                
                if result['plagiarism_score'] > highest_similarity:
                    highest_similarity = result['plagiarism_score']
                    
                if result['plagiarism_score'] > 0.3:  # Threshold for recording matches
                    plagiarism_details.append({
                        'compared_with': other_submission.student.username,
                        'similarity_score': result['plagiarism_score'],
                        'matching_passages': result['matching_passages']
                    })
        
        # Store results
        submission.plagiarism_percentage = highest_similarity * 100
        submission.plagiarism_details = plagiarism_details
        submission.save()
        
        return JsonResponse({
            'success': True,
            'percentage': submission.plagiarism_percentage,
            'details': plagiarism_details
        })
        
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)

@require_POST
def check_ai_content(request, submission_id):
    submission = get_object_or_404(AssignmentSubmission, id=submission_id)
    
    try:
        # Extract text
        text = extract_text_from_file(submission.file)
        if not text:
            return JsonResponse({'success': False, 'error': 'Could not extract text from file'})
        
        # Check for cached results
        cache_key = f'ai_detection_{submission_id}'
        result = cache.get(cache_key)
        
        if not result:
            # Perform AI content detection
            result = checker.detect_ai_content(text)
            # Cache the result for 24 hours
            cache.set(cache_key, result, 60 * 60 * 24)
        
        # Store results
        submission.ai_indicators = result
        submission.save()
        
        return JsonResponse({
            'success': True,
            'ai_probability': result['ai_probability'],
            'metrics': result['metrics'],
            'interpretation': result['interpretation']
        })
        
    except Exception as e:
        return JsonResponse({'success': False, 'error': str(e)}, status=500)